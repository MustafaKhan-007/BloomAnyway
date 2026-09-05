"""Slide decks, drawn into pages so a .pptx reads like any other PDF here.

PowerPoint's own renderer isn't on the host and never will be — there is no
LibreOffice to shell out to on a plain Python dyno — so a deck is read with
python-pptx and painted into a PDF with fpdf2: one page per slide, at the
deck's own page size, with every shape drawn where the deck puts it. Text,
pictures, tables and solid fills come across, which is what a teaching deck is
made of. Charts, SmartArt, gradients and animations do not, so a deck built
around those is better exported to PDF in PowerPoint and uploaded as that.

The conversion happens once, on upload, and the PDF is what gets stored: from
there on the reader, the page counter, the bookmarks and the download are the
ones every other document already has.
"""
from __future__ import annotations

import logging
import os
from io import BytesIO

from fpdf import FPDF

log = logging.getLogger(__name__)

#: A deck this big is nearly always a video someone dropped onto a slide, and
#: reading it whole into memory to draw it is not worth the risk.
MAX_BYTES = 40 * 1024 * 1024

#: PowerPoint measures in English Metric Units; PDFs here are drawn in points.
EMU_PER_PT = 12700

#: Fallbacks for text a deck leaves to its theme, in points.
TITLE_PT = 30.0
BODY_PT = 16.0

#: Core PDF fonts are Latin-1 only, and a deck is full of typographer's
#: punctuation, so the same gentle transliteration the keepsake PDF uses.
_MAP = {
    "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
    "\u2013": "-", "\u2014": "-", "\u2026": "...", "\u2022": "-",
    "\u00a0": " ", "\u2039": "<", "\u203a": ">", "\u00b7": "-",
    "\u2192": "->", "\u2190": "<-", "\u2713": "v", "\u00ab": "<<",
    "\u00bb": ">>",
}


class SlideError(ValueError):
    pass


#: Both office drawers answer to the same few names, so the upload path can
#: hand a file to whichever one knows how to read it.
Error = SlideError


def is_deck(filename: str) -> bool:
    return os.path.splitext(filename or "")[1].lower() == ".pptx"


def pdf_name(filename: str) -> str:
    """What the drawn deck is called once it is a PDF."""
    base = os.path.splitext(os.path.basename(filename or ""))[0] or "slides"
    return f"{base}.pdf"


def title_from(filename: str) -> str:
    """The deck's own name, for the piece list, without the extension on it."""
    return (os.path.splitext(os.path.basename(filename or ""))[0]
            .replace("_", " ").strip()[:160] or "Slides")


def _t(text) -> str:
    s = str(text or "")
    for bad, plain in _MAP.items():
        s = s.replace(bad, plain)
    return s.encode("latin-1", "replace").decode("latin-1")


def _pt(emu) -> float:
    try:
        return float(emu) / EMU_PER_PT
    except (TypeError, ValueError):
        return 0.0


class _Frame:
    """Maps a shape's own coordinates onto the page.

    At the top level that is only a unit change. Inside a group it also
    carries the group's offset and scale, because a grouped shape is placed in
    the group's own coordinate space rather than the slide's.
    """

    def __init__(self, ox=0.0, oy=0.0, sx=1.0 / EMU_PER_PT, sy=1.0 / EMU_PER_PT,
                 cx0=0, cy0=0):
        self.ox, self.oy, self.sx, self.sy = ox, oy, sx, sy
        self.cx0, self.cy0 = cx0, cy0

    def box(self, shape) -> tuple[float, float, float, float] | None:
        try:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
        except (AttributeError, ValueError):
            return None
        if None in (left, top, width, height):
            return None
        return (self.ox + (int(left) - self.cx0) * self.sx,
                self.oy + (int(top) - self.cy0) * self.sy,
                int(width) * self.sx,
                int(height) * self.sy)


def to_pdf(data: bytes) -> bytes:
    """Draw a .pptx into a PDF. Raises :class:`SlideError` if it can't be read."""
    if not data:
        raise SlideError("That file was empty.")
    if len(data) > MAX_BYTES:
        raise SlideError(
            f"That deck is over {MAX_BYTES // (1024 * 1024)} MB. Export it as "
            "a PDF from PowerPoint and upload that instead.")
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - dependency is in requirements
        raise SlideError("Slide decks can't be read on this server yet.") from exc
    try:
        deck = Presentation(BytesIO(data))
    except Exception as exc:
        raise SlideError(
            "That doesn't open as a PowerPoint file. If it came from Keynote "
            "or Google Slides, export it as a PDF and upload that.") from exc

    width = _pt(deck.slide_width) or 720.0
    height = _pt(deck.slide_height) or 540.0
    pdf = FPDF(unit="pt", format=(width, height))
    pdf.set_auto_page_break(False)
    pdf.set_margins(0, 0, 0)
    pdf.c_margin = 0

    slides = list(deck.slides)
    if not slides:
        raise SlideError("There are no slides in that deck.")
    for slide in slides:
        pdf.add_page()
        _paint_background(pdf, slide, width, height)
        frame, sizes = _Frame(), _theme_sizes(slide)
        # The layout's own furniture first — a banner or a logo sits behind
        # what the slide itself says, the same way PowerPoint stacks them.
        for shape in _layout_furniture(slide):
            _draw(pdf, shape, frame, sizes)
        for shape in slide.shapes:
            _draw(pdf, shape, frame, sizes)
    return bytes(pdf.output())


def _layout_furniture(slide) -> list:
    """Non-placeholder shapes from the slide's layout: banners, rules, logos."""
    try:
        shapes = list(slide.slide_layout.shapes)
    except Exception:
        return []
    return [s for s in shapes if not getattr(s, "is_placeholder", False)]


def _theme_sizes(slide) -> dict:
    """How big the deck's master says its text is, by outline level.

    Almost no slide states a font size: it is inherited from the master, and
    without reading it every title would come out the same middling size as
    the body under it.
    """
    out = {"title": TITLE_PT, "other": BODY_PT, "body": {},
           "title_align": "L", "body_align": "L"}
    try:
        from pptx.oxml.ns import qn
        styles = slide.slide_layout.slide_master.element.find(qn("p:txStyles"))
        if styles is None:
            return out
        title_style = styles.find(qn("p:titleStyle"))
        title = _style_size(title_style, 0)
        if title:
            out["title"] = title
        out["title_align"] = _style_align(title_style, 0) or "L"
        other = _style_size(styles.find(qn("p:otherStyle")), 0)
        if other:
            out["other"] = other
        body = styles.find(qn("p:bodyStyle"))
        out["body_align"] = _style_align(body, 0) or "L"
        for level in range(9):
            size = _style_size(body, level)
            if size:
                out["body"][level] = size
    except Exception:
        pass
    return out


def _level_rules(style, level: int):
    from pptx.oxml.ns import qn
    if style is None:
        return None
    return style.find(qn(f"a:lvl{level + 1}pPr"))


def _style_size(style, level: int) -> float | None:
    try:
        from pptx.oxml.ns import qn
        rules = _level_rules(style, level)
        props = rules.find(qn("a:defRPr")) if rules is not None else None
        raw = props.get("sz") if props is not None else None
        return float(raw) / 100 if raw else None
    except Exception:
        return None


def _style_align(style, level: int) -> str | None:
    try:
        rules = _level_rules(style, level)
        return {"ctr": "C", "r": "R", "just": "J"}.get(
            (rules.get("algn") or "") if rules is not None else "")
    except Exception:
        return None


# --- paint ------------------------------------------------------------------

def _rgb(color) -> tuple[int, int, int] | None:
    """A colour only when the deck states one outright.

    Theme colours would mean resolving the theme part and its overrides, which
    is a rabbit hole; anything left to the theme falls back to plain ink.
    """
    try:
        if color is None or color.type is None:
            return None
        rgb = color.rgb
    except Exception:
        return None
    try:
        return (rgb[0], rgb[1], rgb[2])
    except (TypeError, IndexError):
        try:
            value = int(str(rgb), 16)
        except ValueError:
            return None
        return ((value >> 16) & 255, (value >> 8) & 255, value & 255)


def _fill_rgb(holder) -> tuple[int, int, int] | None:
    """The solid colour something is filled with, if it is filled at all.

    Reading a fill is not a safe question in python-pptx: asking an unfilled
    shape for its colour raises rather than answering, so every step of it is
    guarded.
    """
    try:
        fill = holder.fill
        if fill.type != 1:  # MSO_FILL.SOLID
            return None
        return _rgb(fill.fore_color)
    except Exception:
        return None


def _paint_background(pdf, slide, width, height) -> None:
    """White unless the slide, its layout or the master says otherwise."""
    layout = getattr(slide, "slide_layout", None)
    colour = None
    for holder in (slide, layout, getattr(layout, "slide_master", None)):
        if holder is None:
            continue
        colour = _fill_rgb(getattr(holder, "background", None))
        if colour:
            break
    pdf.set_fill_color(*(colour or (255, 255, 255)))
    pdf.rect(0, 0, width, height, "F")


def _is_dark(colour) -> bool:
    r, g, b = colour
    return (0.299 * r + 0.587 * g + 0.114 * b) < 128


def _draw(pdf, shape, frame: _Frame, sizes: dict) -> None:
    """One shape, wherever the deck put it. Never raises: a slide with one odd
    shape on it should still come out as a page."""
    try:
        if _is_group(shape):
            _draw_group(pdf, shape, frame, sizes)
            return
        box = frame.box(shape)
        if box is None or box[2] <= 0 or box[3] <= 0:
            return
        image = _image_blob(shape)
        if image is not None:
            _draw_image(pdf, image, box)
            return
        if getattr(shape, "has_table", False):
            _draw_table(pdf, shape.table, box)
            return
        fill = _fill_rgb(shape)
        if fill:
            pdf.set_fill_color(*fill)
            pdf.rect(box[0], box[1], box[2], box[3], "F")
        if getattr(shape, "has_text_frame", False):
            _draw_text(pdf, shape, box, sizes,
                       on_dark=bool(fill and _is_dark(fill)))
    except Exception:
        log.debug("slides: skipped a shape that wouldn't draw", exc_info=True)


def _is_group(shape) -> bool:
    try:
        return shape.shape_type == 6  # MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _draw_group(pdf, group, frame: _Frame, sizes: dict) -> None:
    box = frame.box(group)
    if box is None:
        return
    x, y, w, h = box
    child = _Frame(x, y, frame.sx, frame.sy)
    off, ext = _child_space(group)
    if off and ext and ext[0] and ext[1]:
        child = _Frame(x, y, w / ext[0], h / ext[1], off[0], off[1])
    for inner in group.shapes:
        _draw(pdf, inner, child, sizes)


def _child_space(group):
    """A group's own coordinate origin and size, straight off its XML."""
    try:
        from pptx.oxml.ns import qn
        props = group._element.find(qn("p:grpSpPr"))
        xfrm = props.find(qn("a:xfrm")) if props is not None else None
        if xfrm is None:
            return None, None
        off = xfrm.find(qn("a:chOff"))
        ext = xfrm.find(qn("a:chExt"))
        if off is None or ext is None:
            return None, None
        return ((int(off.get("x")), int(off.get("y"))),
                (int(ext.get("cx")), int(ext.get("cy"))))
    except Exception:
        return None, None


def _image_blob(shape):
    try:
        return shape.image.blob
    except Exception:
        return None


def _draw_image(pdf, blob: bytes, box) -> None:
    """The picture, whole and undistorted, centred in the space it was given."""
    x, y, w, h = box
    try:
        from PIL import Image
        with Image.open(BytesIO(blob)) as img:
            iw, ih = img.size
            fmt = (img.format or "").upper()
            if fmt not in ("PNG", "JPEG", "GIF", "BMP"):
                buf = BytesIO()
                img.convert("RGBA" if "A" in img.getbands() else "RGB").save(
                    buf, format="PNG")
                blob = buf.getvalue()
    except Exception:
        iw = ih = 0
    if iw and ih:
        scale = min(w / iw, h / ih)
        draw_w, draw_h = iw * scale, ih * scale
        x += (w - draw_w) / 2
        y += (h - draw_h) / 2
        w, h = draw_w, draw_h
    pdf.image(BytesIO(blob), x=x, y=y, w=w, h=h)


# --- words ------------------------------------------------------------------

_ALIGN = {1: "C", 2: "R", 3: "J"}  # PP_ALIGN centre / right / justify


#: Placeholder kinds worth knowing apart (PP_PLACEHOLDER values).
_TITLE_HOLDERS = (1, 3, 5)  # TITLE, CENTER_TITLE, VERTICAL_TITLE
_BULLET_HOLDERS = (2, 6, 7)  # BODY, VERTICAL_BODY, OBJECT


def _holder_type(shape) -> int | None:
    try:
        if not shape.is_placeholder:
            return None
        return int(shape.placeholder_format.type)
    except Exception:
        return None


def _run_size(paragraph, fallback: float) -> float:
    for run in paragraph.runs:
        try:
            if run.font.size is not None:
                return run.font.size.pt
        except Exception:
            continue
    try:
        if paragraph.font.size is not None:
            return paragraph.font.size.pt
    except Exception:
        pass
    return fallback


def _run_style(paragraph) -> str:
    bold = italic = False
    for run in paragraph.runs:
        try:
            bold = bold or bool(run.font.bold)
            italic = italic or bool(run.font.italic)
        except Exception:
            continue
    return ("B" if bold else "") + ("I" if italic else "")


def _run_colour(paragraph):
    for run in paragraph.runs:
        found = _rgb(getattr(run.font, "color", None))
        if found:
            return found
    return None


def _is_title(shape) -> bool:
    return _holder_type(shape) in _TITLE_HOLDERS


def _bulleted(shape, paragraph) -> bool:
    """Whether a line carries a bullet, as far as the file admits to it."""
    try:
        from pptx.oxml.ns import qn
        props = paragraph._p.find(qn("a:pPr"))
        if props is not None:
            if props.find(qn("a:buNone")) is not None:
                return False
            if (props.find(qn("a:buChar")) is not None
                    or props.find(qn("a:buAutoNum")) is not None):
                return True
    except Exception:
        return False
    # Left to the layout: the body of a slide bullets its lines, while a
    # title, a subtitle or a caption is written straight out.
    return _holder_type(shape) in _BULLET_HOLDERS


def _lines(shape, sizes: dict) -> list[dict]:
    """The shape's paragraphs, as the pieces needed to draw each one."""
    title = _is_title(shape)
    holder = _holder_type(shape)
    out = []
    for paragraph in shape.text_frame.paragraphs:
        text = _t(paragraph.text).strip()
        try:
            level = int(paragraph.level or 0)
        except (TypeError, ValueError):
            level = 0
        if title:
            inherited = sizes["title"]
        elif holder in _BULLET_HOLDERS:
            inherited = sizes["body"].get(level, BODY_PT)
        else:
            inherited = sizes["other"]
        size = _run_size(paragraph, inherited)
        out.append({
            "text": text,
            "size": max(6.0, min(96.0, size)),
            "style": _run_style(paragraph),
            "colour": _run_colour(paragraph),
            "align": _ALIGN.get(
                getattr(paragraph.alignment, "value", None),
                sizes["title_align"] if title else sizes["body_align"]),
            "indent": level * 14.0,
            "bullet": bool(text) and not title and _bulleted(shape, paragraph),
        })
    while out and not out[-1]["text"]:
        out.pop()
    return out


def _block_height(lines) -> float:
    return sum(line["size"] * 1.32 for line in lines)


def _draw_text(pdf, shape, box, sizes: dict, *, on_dark: bool) -> None:
    x, y, w, h = box
    lines = _lines(shape, sizes)
    if not any(line["text"] for line in lines):
        return
    pad = min(6.0, w / 12)
    x, w = x + pad, max(12.0, w - pad * 2)

    # PowerPoint shrinks text to fit its box; without that a wordy slide would
    # run off the bottom of the page instead of over the edge of its shape.
    needed = _block_height(lines)
    if h > 0 and needed > h:
        squeeze = max(0.5, h / needed)
        for line in lines:
            line["size"] = max(6.0, line["size"] * squeeze)

    cursor = y + _anchor_offset(shape, h, _block_height(lines))
    default = (255, 255, 255) if on_dark else (26, 26, 26)
    for line in lines:
        step = line["size"] * 1.32
        if not line["text"]:
            cursor += step
            continue
        pdf.set_font("Helvetica", line["style"], line["size"])
        pdf.set_text_color(*(line["colour"] or default))
        left = x + line["indent"]
        width = max(12.0, w - line["indent"])
        if line["bullet"]:
            # Hung off to the side, so a line that wraps sits under its own
            # words rather than back under the dash.
            marker = line["size"] * 0.75
            pdf.set_xy(left, cursor)
            pdf.cell(marker, step, text="-")
            left += marker
            width = max(12.0, width - marker)
        pdf.set_xy(left, cursor)
        try:
            pdf.multi_cell(width, step, text=line["text"],
                           align=line["align"], new_x="LEFT", new_y="NEXT")
            cursor = pdf.get_y()
        except Exception:
            cursor += step
        if cursor > pdf.h:
            break


def _anchor_offset(shape, box_height: float, text_height: float) -> float:
    """Top, middle or bottom of the shape, as the deck asks."""
    if box_height <= 0 or text_height >= box_height:
        return 0.0
    try:
        anchor = getattr(shape.text_frame.vertical_anchor, "value", None)
    except Exception:
        anchor = None
    if anchor == 3:  # MSO_ANCHOR.MIDDLE
        return (box_height - text_height) / 2
    if anchor == 4:  # MSO_ANCHOR.BOTTOM
        return box_height - text_height
    return 0.0


def _draw_table(pdf, table, box) -> None:
    x, y, w, h = box
    columns = [max(1, int(col.width or 0)) for col in table.columns] or [1]
    total = sum(columns) or 1
    widths = [w * (c / total) for c in columns]
    rows = [max(1, int(row.height or 0)) for row in table.rows] or [1]
    tall = sum(rows) or 1
    heights = [h * (r / tall) for r in rows]

    pdf.set_draw_color(190, 190, 190)
    pdf.set_line_width(0.6)
    top = y
    for r, row in enumerate(table.rows):
        left = x
        for c, cell in enumerate(row.cells):
            cw = widths[c] if c < len(widths) else 0
            ch = heights[r] if r < len(heights) else 0
            if cw <= 0 or ch <= 0:
                continue
            fill = _fill_rgb(cell)
            if fill:
                pdf.set_fill_color(*fill)
                pdf.rect(left, top, cw, ch, "F")
            pdf.rect(left, top, cw, ch)
            text = _t(cell.text).strip()
            if text:
                pdf.set_font("Helvetica", "", min(12.0, max(7.0, ch / 3)))
                pdf.set_text_color(26, 26, 26)
                pdf.set_xy(left + 3, top + 3)
                try:
                    pdf.multi_cell(cw - 6, min(12.0, max(8.0, ch / 3)),
                                   text=text, align="L",
                                   new_x="LEFT", new_y="NEXT")
                except Exception:
                    pass
            left += cw
        top += heights[r] if r < len(heights) else 0
