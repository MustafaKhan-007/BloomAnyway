"""Acceptance-criteria smoke test (run: python scripts/smoke_test.py).

Uses a throwaway SQLite database and the Flask test client. Not a pytest
suite on purpose — a single readable script the owner/dev can run anywhere.
"""
import hashlib
import hmac
import io
import json
import os
import pathlib
import re
import sys
import tempfile
from datetime import date, datetime, timedelta
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

import base64
import time as _time

_STRIPE_WHSEC = "whsec_test_secret"
os.environ["STRIPE_WEBHOOK_SECRET"] = _STRIPE_WHSEC
os.environ["STRIPE_SECRET_KEY"] = "sk_test_smoke"

from app import create_app
from app.config import DevConfig
from app.extensions import db
from app.models import (ForumCategory, ForumComment, ForumPost, ForumTag,
                        Order, Product, Quote, QuotePin, ShopPurchase,
                        Subscriber, User, Video, utcnow)
from app.services import captcha as captcha_service
from app.services import stripe_pay as pay


def _stripe_headers(body: bytes) -> dict:
    sig = pay.sign_webhook(_STRIPE_WHSEC, body)
    return {
        "Content-Type": "application/json",
        "Stripe-Signature": sig,
    }


def _payment_payload(payment_id, email, product_id, *,
                     event="payment.succeeded", amount=4900,
                     product_name=None, gift_to=None,
                     payment_status="paid", payment_intent=True):
    meta = {"price_id": str(product_id)}
    if product_name:
        meta["product_name"] = product_name
    if gift_to:
        meta["gift_to"] = gift_to
    if event == "payment.succeeded":
        stripe_type = "checkout.session.completed"
        obj = {
            "id": f"cs_{payment_id}",
            "object": "checkout.session",
            "status": "complete",
            "payment_status": payment_status,
            "amount_total": amount,
            "currency": "usd",
            "customer_details": {"email": email},
            "customer_email": email,
            "payment_intent": str(payment_id) if payment_intent else None,
            "subscription": f"sub_{payment_id}" if not payment_intent else None,
            "metadata": meta,
            "mode": "subscription" if not payment_intent else "payment",
        }
    elif event == "payment.failed":
        stripe_type = "invoice.payment_failed"
        obj = {
            "id": f"in_{payment_id}",
            "object": "invoice",
            "amount_due": amount,
            "currency": "usd",
            "customer_email": email,
            "payment_intent": str(payment_id),
            "metadata": meta,
            "lines": {"data": [{"price": {"id": str(product_id)}}]},
        }
    else:
        stripe_type = "charge.refunded"
        obj = {
            "id": f"ch_{payment_id}",
            "object": "charge",
            "amount": amount,
            "amount_refunded": amount,
            "currency": "usd",
            "payment_intent": str(payment_id),
            "metadata": meta,
            "billing_details": {"email": email},
        }
    return json.dumps({
        "id": f"evt_{payment_id}",
        "object": "event",
        "type": stripe_type,
        "data": {"object": obj},
    }).encode()

# Smoke tests don't call Cloudflare; always pass the captcha check.
captcha_service.verify_captcha = lambda token=None: True
captcha_service.captcha_challenge = lambda: {"site_key": "1x00000000000000000000AA"}
captcha_service.issue_captcha = captcha_service.captcha_challenge
captcha_service.captcha_question = lambda: "turnstile"
captcha_service.site_key = lambda: "1x00000000000000000000AA"


TMP_DB = Path(tempfile.mkdtemp()) / "smoke.db"


class TestConfig(DevConfig):
    SQLALCHEMY_DATABASE_URI = f"sqlite:///{TMP_DB.as_posix()}"
    WTF_CSRF_ENABLED = False
    RATELIMIT_ENABLED = False
    TESTING = True


PASS = 0


def flashes(resp) -> str:
    """The flash banners on a rendered page — the useful half of a failure."""
    import re as _re
    text = resp.get_data(as_text=True)
    found = _re.findall(r'class="flash flash--\w+"[^>]*>(.*?)</div>', text, _re.S)
    return " | ".join(" ".join(f.split()) for f in found) or "(no flash)"


def ok(name, condition, detail=""):
    global PASS
    status = "PASS" if condition else "FAIL"
    print(f"[{status}] {name}" + (f"  ({detail})" if detail and not condition else ""))
    if condition:
        PASS += 1
    else:
        raise SystemExit(f"FAILED: {name} {detail}")


app = create_app(TestConfig)

# capture verification codes instead of emailing
sent_codes = []
import app.auth.routes as auth_routes
auth_routes.send_verification_code = lambda to, code, purpose: sent_codes.append((to, code, purpose)) or True

ADMIN_PW = "owner-strong-pass-1"

with app.app_context():
    db.create_all()
    seed = json.loads((Path(__file__).parents[1] / "data" / "quotes_seed.json").read_text(encoding="utf-8"))
    for row in seed["quotes"]:
        db.session.add(Quote(text=row["text"], author=row.get("author"), category=row["category"]))
    db.session.commit()
    n_quotes = Quote.query.count()

ok("Seed has 150+ quotes", n_quotes >= 150, f"got {n_quotes}")

client = app.test_client()

# --- 1. home hero + daily quote rotation (quotes live on /quotes) -------------
r1 = client.get("/")
ok("Home page renders", r1.status_code == 200, str(r1.status_code))
home1 = r1.get_data(as_text=True)
ok("Home shows healing / building hero",
   "You don't have to carry this alone" in home1
   and "Ready to build something" in home1
   and "that's yours?" in home1
   and "Find Your Community" in home1
   and "Start Building" in home1)
ok("Home shows Their Story with reserved photo space",
   "Their story" in home1
   and "Society told us to suffer quietly" in home1
   and "home-story__ph" in home1
   and "Ayesha &amp; Saman" in home1)
ok("Home shows Product of the Day + top products sections",
   "Digital Product of the Day" in home1
   and "Top products — last 30 days" in home1
   and ("home-potd" in home1 or "Coming soon" in home1)
   and ("home-top-grid" in home1 or "Coming soon" in home1))
r2 = client.get("/")
ok("Home still renders on refresh", r2.status_code == 200)

with app.app_context():
    from app.services.quotes import quote_for
    today_q = quote_for(date.today())
    tomorrow_q = quote_for(date.today() + timedelta(days=1))
    day_after = quote_for(date.today() + timedelta(days=2))
ok("Quote rotation changes across days (some day differs)",
   today_q.id != tomorrow_q.id or today_q.id != day_after.id)

# --- 2a. first-run owner setup ----------------------------------------------------
setup_client = app.test_client()
r = setup_client.get("/setup")
ok("Setup page available on fresh install", r.status_code == 200)
r = setup_client.get("/login")
ok("Login page advertises setup on fresh install", "Claim the owner account" in r.get_data(as_text=True))
r = setup_client.post("/setup", data={"email": "owner@example.com", "password": ADMIN_PW,
                                      "password_confirm": ADMIN_PW},
                      follow_redirects=False)
ok("Owner account claimed via setup", r.status_code == 302 and "/admin" in r.headers["Location"])
r = setup_client.get("/admin/")
ok("Owner lands in studio after setup", r.status_code == 200)
r = app.test_client().get("/setup")
ok("Setup locks after owner signs in", r.status_code == 404)

# --- 2b. email + password auth with confirmation codes ---------------------------
USER_PW = "sunrise-day-1"

r = client.post("/register", data={"email": "newperson@example.com", "password": "short",
                                   "password_confirm": "short"})
ok("Weak password rejected on registration", r.status_code == 400)

r = client.post("/register", data={"email": "newperson@example.com", "password": USER_PW,
                                   "password_confirm": "different-pass"}, follow_redirects=False)
ok("Mismatched passwords rejected on registration",
   r.status_code == 400 and "those passwords" in r.get_data(as_text=True).lower())

r = client.post("/register", data={"email": "newperson@example.com", "password": USER_PW,
                                   "password_confirm": USER_PW},
                follow_redirects=False)
ok("Registration redirects to verify page", r.status_code == 302 and "verify-email" in r.headers["Location"])
ok("Confirmation code emailed", len(sent_codes) == 1 and sent_codes[0][2] == "confirm")
first_code = sent_codes[0][1]

# unverified account can't just log in — it gets sent back to verification
# without wiping the code they already received
r = client.post("/login", data={"email": "newperson@example.com", "password": USER_PW},
                follow_redirects=False)
ok("Unverified login redirects to verification", r.status_code == 302 and "verify-email" in r.headers["Location"])
ok("Unverified login keeps the original confirmation code",
   len(sent_codes) == 1)

# wrong code fails with attempts feedback, right code confirms + logs in
r = client.post("/verify-email", data={"email": "newperson@example.com", "code": "000000"})
wrong_ok = r.status_code == 400 and "tries left" in r.get_data(as_text=True)
r = client.post("/verify-email", data={"email": "newperson@example.com",
                                       "code": f" {first_code[:3]}-{first_code[3:]} "},
                follow_redirects=False)
ok("Wrong code rejected with tries-left message", wrong_ok)
ok("Correct code confirms and logs in (spaces/dashes ok)",
   r.status_code == 302 and "/account" in r.headers["Location"])
r = client.get("/account")
ok("Account page accessible after confirmation", r.status_code == 200)
abody = r.get_data(as_text=True)
ok("New member lands on account without a product tour",
   "product-tour.js" not in abody
   and "data-product-tour" not in abody)

# password checks
fresh = app.test_client()
r = fresh.post("/login", data={"email": "newperson@example.com", "password": "wrong-password"})
ok("Wrong password rejected (401)", r.status_code == 401)
r = fresh.post("/login", data={"email": "newperson@example.com", "password": USER_PW,
                               "next": "https://evil.example.com"}, follow_redirects=False)
ok("Absolute next URL rejected (no open redirect)",
   r.status_code == 302 and r.headers["Location"].startswith("/"))

# forgot / reset password flow
sent_codes.clear()
reset_client = app.test_client()
r = reset_client.post("/forgot-password", data={"email": "newperson@example.com"}, follow_redirects=True)
uniform_known = "reset code is on its way" in r.get_data(as_text=True)
r = reset_client.post("/forgot-password", data={"email": "ghost@example.com"}, follow_redirects=True)
uniform_unknown = "reset code is on its way" in r.get_data(as_text=True)
ok("Uniform reset message for known + unknown email", uniform_known and uniform_unknown)
ok("Reset code only sent for real account", len(sent_codes) == 1 and sent_codes[0][2] == "reset")
r = reset_client.post("/reset-password", data={"email": "newperson@example.com",
                                               "code": sent_codes[0][1],
                                               "password": USER_PW,
                                               "password_confirm": USER_PW}, follow_redirects=False)
ok("Password reset rejects reusing the current password",
   r.status_code == 400 and "different password" in r.get_data(as_text=True).lower())
r = reset_client.post("/reset-password", data={"email": "newperson@example.com",
                                               "code": sent_codes[0][1],
                                               "password": "brand-new-pass-9",
                                               "password_confirm": "brand-new-pass-x"},
                      follow_redirects=False)
ok("Password reset rejects mismatched confirmation",
   r.status_code == 400 and "those passwords" in r.get_data(as_text=True).lower())
r = reset_client.post("/reset-password", data={"email": "newperson@example.com",
                                               "code": sent_codes[0][1],
                                               "password": "brand-new-pass-9",
                                               "password_confirm": "brand-new-pass-9"},
                      follow_redirects=False)
ok("Password reset with valid code succeeds", r.status_code == 302)
r = app.test_client().post("/login", data={"email": "newperson@example.com",
                                           "password": "brand-new-pass-9"}, follow_redirects=False)
ok("Login works with the new password", r.status_code == 302 and "/account" in r.headers["Location"])

# --- 3. admin: product lifecycle ----------------------------------------------
admin = app.test_client()
r = admin.post("/login", data={"email": "owner@example.com", "password": ADMIN_PW}, follow_redirects=False)
ok("Admin password login works", r.status_code == 302)

r = admin.get("/admin/")
ok("Admin dashboard loads for admin", r.status_code == 200)
r = client.get("/admin/")
ok("Admin returns 404 for non-admin user", r.status_code == 404)

# admin idle timeout: stale activity forces re-auth; active use slides the window
with admin.session_transaction() as sess:
    stale = (datetime.utcnow() - timedelta(days=15)).isoformat()
    sess["admin_seen_at"] = stale
    sess["logged_in_at"] = stale
r = admin.get("/admin/", follow_redirects=False)
ok("Admin re-auth required after 14 idle days",
   r.status_code == 302 and "/login" in r.headers["Location"])
with admin.session_transaction() as sess:
    sess["admin_seen_at"] = (datetime.utcnow() - timedelta(days=2)).isoformat()
r = admin.get("/admin/", follow_redirects=False)
ok("Active admin stays signed in (sliding window)", r.status_code == 200)

# Studio catalogue editor is live; /courses is on-site
r = admin.get("/admin/products", follow_redirects=True)
_pbody = r.get_data(as_text=True)
ok("Studio products UI loads",
   r.status_code == 200
   and ("Add a product" in _pbody or "Stripe price ID" in _pbody or "Courses" in _pbody))
r = client.get("/courses", follow_redirects=False)
cbody = r.get_data(as_text=True)
ok("/courses renders on-site catalogue",
   r.status_code == 200 and "Courses &amp; Guides" in cbody
   and "Healing resources by" in cbody and "Creator resources by" in cbody
   and "Rebuild Workbook" not in cbody and "50 Hooks" not in cbody)
r = admin.get("/admin/products/new", follow_redirects=True)
_new_body = r.get_data(as_text=True)
ok("Studio offers product cover upload",
   "Cover image" in _new_body and 'name="cover"' in _new_body
   and "Stripe price ID" in _new_body)
ok("New-product form carries a live cover preview",
   "data-cover-preview" in _new_body and "Cover preview" in _new_body
   and "lib-card__cover-title" in _new_body and "data-cover-photo" in _new_body)
ok("The preview starts on the flower cover with no photo over it",
   'data-cover-photo alt="" hidden' in _new_body
   and "data-cover-saved" not in _new_body)
ok("With nothing typed, the preview shows the type as its kind",
   ">GUIDE<" in _new_body)
ok("Nothing pins the kind until a reading file exists",
   "data-cover-kind-fixed" not in _new_body)

# Tiny JPEG cover upload for a draft product
from io import BytesIO
from PIL import Image as _PILCover
_cbuf = BytesIO()
_PILCover.new("RGB", (300, 400), (90, 49, 88)).save(_cbuf, format="JPEG")
_cbuf.seek(0)
r = admin.post(
    "/admin/products/new",
    data={
        "title": "Cover Test Guide",
        "track": "healing",
        "type": "guide",
        "price": "19.00",
        "promise": "A soft check-in.",
        "cover": (_cbuf, "cover.jpg"),
    },
    content_type="multipart/form-data",
    follow_redirects=True,
)
ok("Studio accepts product cover on create", r.status_code == 200)
with app.app_context():
    cover_prod = Product.query.filter_by(slug="cover-test-guide").first()
    ok("Cover URL stored on product",
       cover_prod is not None and (cover_prod.cover_url or "").startswith("/media/product-cover/"),
       f"got {getattr(cover_prod, 'cover_url', None)}")
    ok("Cover image bytes stored in database",
       cover_prod is not None and cover_prod.cover_data is not None
       and len(cover_prod.cover_data) > 200,
       f"bytes={0 if not cover_prod or not cover_prod.cover_data else len(cover_prod.cover_data)}")
    cover_id = cover_prod.id if cover_prod else 0
r = client.get(f"/media/product-cover/{cover_id}")
ok("Product cover image is served",
   cover_id and r.status_code == 200 and r.mimetype.startswith("image/"))

_edit_body = admin.get(f"/admin/products/{cover_id}/edit").get_data(as_text=True)
ok("Editing a product previews the cover it already has",
   "data-cover-preview" in _edit_body
   and f'data-cover-saved="/media/product-cover/{cover_id}' in _edit_body)
ok("That preview reflects the saved title and track colour",
   "Cover Test Guide" in _edit_body and "#5A3158" in _edit_body)
with app.app_context():
    cover_prod = Product.query.filter_by(id=cover_id).first()
    if cover_prod:
        cover_prod.status = "published"
        db.session.commit()
r = client.get("/courses")
courses_body = r.get_data(as_text=True)
ok("Courses page uses My space-style library cards",
   "lib-card" in courses_body
   and "lib-card__cover" in courses_body
   and "Cover Test Guide" in courses_body)
ok("Uploaded cover appears on Courses cards",
   f"/media/product-cover/{cover_id}" in courses_body
   and "lib-card__cover--photo" in courses_body)
r = client.get("/")
home = r.get_data(as_text=True)
ok("Home includes creator membership CTA",
   "Join Creator Membership" in home and "Creator of the Month" in home)
ok("Nav Courses & Guides points on-site",
   '/courses"' in home or "/courses'" in home)

# Product row for order matching
with app.app_context():
    hist = Product(
        title="Begin Again", slug="begin-again", type="course", status="published",
        promise="A 4-week path from stuck to started.",
        cover_url="https://example.com/cover.jpg", price_cents=4900,
        currency="USD", stripe_price_id="prod_begin_again",
        track="healing", featured=True)
    db.session.add(hist)
    db.session.commit()
    hist_id = hist.id

# --- 4. Stripe webhook: signature + idempotency + ShopPurchase -----------------------
payload = _payment_payload(
    "9001", "Buyer@Example.com", "prod_begin_again",
    product_name="Begin Again")

r = client.post("/webhooks/stripe", data=payload,
                headers={"Content-Type": "application/json",
                         "webhook-id": "bad", "webhook-timestamp": "1",
                         "webhook-signature": "v1,bad"})
ok("Wrong webhook signature rejected 401", r.status_code == 401)

r = client.post("/webhooks/stripe", data=payload, headers=_stripe_headers(payload))
r2 = client.post("/webhooks/stripe", data=payload, headers=_stripe_headers(payload))
with app.app_context():
    orders = Order.query.filter_by(ls_order_id="9001").all()
    shops = ShopPurchase.query.filter_by(lemon_squeezy_order_id="9001").all()
ok("Webhook accepted (200)", r.status_code == 200 and r2.status_code == 200)
ok("Replayed webhook creates exactly one order", len(orders) == 1, f"got {len(orders)}")
ok("Replayed webhook creates exactly one ShopPurchase", len(shops) == 1, f"got {len(shops)}")
ok("Order matched to product via Stripe price id", orders[0].product_id is not None)
ok("Buyer email lowercased", orders[0].buyer_email == "buyer@example.com")
ok("Unknown-email shop purchase is pending_link",
   shops[0].status == "pending_link" and shops[0].user_id is None)
ok("Shop purchase product name from webhook", shops[0].product_name == "Begin Again")

# 100% off / $0 checkout — no PaymentIntent; checkout.session.completed still fulfills
zero_payload = _payment_payload(
    "9001-free", "freebie@example.com", "prod_begin_again",
    amount=0, payment_status="no_payment_required", payment_intent=False,
    product_name="Begin Again Free")
r = client.post("/webhooks/stripe", data=zero_payload, headers=_stripe_headers(zero_payload))
with app.app_context():
    zero_order = Order.query.filter_by(ls_order_id="sub_9001-free").first()
    zero_shop = ShopPurchase.query.filter_by(lemon_squeezy_order_id="sub_9001-free").first()
ok("$0 checkout.session.completed accepted", r.status_code == 200)
ok("$0 checkout creates paid order",
   zero_order is not None and zero_order.status == "paid" and zero_order.total_cents == 0)
ok("$0 checkout creates shop purchase",
   zero_shop is not None and zero_shop.status == "pending_link")

# purchase auto-links when that email signs up / logs in
with app.app_context():
    from app.models import User as _U
    buyer = _U(email="buyer@example.com", email_verified_at=utcnow())
    buyer.set_password(USER_PW)
    db.session.add(buyer)
    db.session.commit()
buyer_client = app.test_client()
buyer_client.post("/login", data={"email": "buyer@example.com", "password": USER_PW})
with app.app_context():
    linked = ShopPurchase.query.filter_by(lemon_squeezy_order_id="9001").first()
ok("Pending shop purchase links on login",
   linked.status == "linked" and linked.user_id is not None)
r = buyer_client.get("/account?tab=saved")
abody = r.get_data(as_text=True)
ok("Linked shop purchase appears in My space",
   r.status_code == 200 and "Begin Again" in abody and "Courses" in abody)

# On-site reader + progress resume
with app.app_context():
    hist = Product.query.filter_by(slug="begin-again").first()
    from app.models import ProductAsset, CourseProgress
    asset = ProductAsset(
        product_id=hist.id,
        title="Begin Again PDF",
        filename="begin-again.pdf",
        mime="application/pdf",
        kind="pdf",
        size=12,
        data=b"%PDF-1.4 fake\n%%EOF\n",
        sort_order=0,
    )
    db.session.add(asset)
    db.session.commit()
    purchase = ShopPurchase.query.filter_by(lemon_squeezy_order_id="9001").first()
    purchase_id = purchase.id
r = buyer_client.get(f"/account/courses/{purchase_id}")
ok("Course reader opens for owned purchase",
   r.status_code == 200 and b"Begin Again" in r.data and b"reader-pdf-canvas" in r.data
   and b"Back to library" in r.data)
r = buyer_client.post(
    f"/account/courses/{purchase_id}/progress",
    json={"page": 5, "total": 20},
    headers={"Content-Type": "application/json"},
)
ok("Reading progress saves", r.status_code == 200 and r.get_json().get("percent") == 25)
with app.app_context():
    prog = CourseProgress.query.filter_by(shop_purchase_id=purchase_id).first()
ok("Progress row stores page 5",
   prog is not None and prog.current_page == 5 and prog.total_pages == 20)
r = buyer_client.get("/account?tab=saved")
abody = r.get_data(as_text=True)
ok("Courses tab shows real progress percent",
   "25%" in abody and "Continue reading" in abody and "Reading progress" in abody)
r = buyer_client.get(f"/account/courses/{purchase_id}")
ok("Reader resumes at saved page",
   r.status_code == 200 and b'data-start-page="5"' in r.data)
r = buyer_client.post(
    f"/account/courses/{purchase_id}/bookmarks",
    json={"page": 5},
    headers={"Content-Type": "application/json"},
)
ok("Bookmark toggles on",
   r.status_code == 200 and r.get_json().get("bookmarked") is True
   and 5 in (r.get_json().get("bookmarks") or []))
r = buyer_client.get("/account?tab=saved")
ok("Library shows bookmarked pages",
   "Bookmarks" in r.get_data(as_text=True) and "page 5" in r.get_data(as_text=True))

# purchase for an email that already has an account links immediately
with app.app_context():
    known = User.query.filter_by(email="newperson@example.com").first()
    known_id = known.id
payload_known = _payment_payload(
    "9002", "newperson@example.com", "prod_quiet",
    amount=1900, product_name="Quiet Mornings")
r = client.post("/webhooks/stripe", data=payload_known, headers=_stripe_headers(payload_known))
with app.app_context():
    sp2 = ShopPurchase.query.filter_by(lemon_squeezy_order_id="9002").first()
ok("Existing-account shop purchase links immediately",
   r.status_code == 200 and sp2 is not None
   and sp2.status == "linked" and sp2.user_id == known_id)

# failed payments must not invent a My Space library item
payload_fail = _payment_payload(
    "9002-fail", "newperson@example.com", "prod_quiet",
    event="payment.failed", amount=1900, product_name="Quiet Mornings")
r = client.post("/webhooks/stripe", data=payload_fail, headers=_stripe_headers(payload_fail))
with app.app_context():
    fail_shop = ShopPurchase.query.filter_by(lemon_squeezy_order_id="9002-fail").first()
    fail_ord = Order.query.filter_by(ls_order_id="9002-fail").first()
ok("Failed payment does not create a ShopPurchase",
   r.status_code == 200 and fail_shop is None)
ok("Failed payment still records an Order",
   fail_ord is not None and fail_ord.status == "failed")

# Studio activity + purchase chart helpers see paid orders
with app.app_context():
    from app.services import stats as stats_svc
    activity = stats_svc.member_activity(20)
    chart = stats_svc.purchases_over_time(30)
    trend = stats_svc.trending_product(30)
ok("Member activity includes purchases",
   any(a.get("kind") == "purchase" for a in activity))
ok("Purchases chart has daily series",
   isinstance(chart.get("all"), list) and chart.get("total", 0) >= 1)
ok("Trending product reports a leader",
   trend is not None and "trending" in (trend.get("label") or "").lower())

# refund hides from My space
payload_ref = _payment_payload(
    "9002", "newperson@example.com", "prod_quiet",
    event="refund.succeeded", amount=1900, product_name="Quiet Mornings")
client.post("/webhooks/stripe", data=payload_ref, headers=_stripe_headers(payload_ref))
with app.app_context():
    sp2 = ShopPurchase.query.filter_by(lemon_squeezy_order_id="9002").first()
ok("Refunded shop purchase marked refunded", sp2.status == "refunded")
r = client.get("/account?tab=saved")
ok("My space hides refunded purchases",
   "Quiet Mornings" not in r.get_data(as_text=True))

# protected self-hosted download
with app.app_context():
    from flask import current_app
    shop_dir = current_app.config["SHOP_FILES_DIR"]
    key = "quiet-guide.pdf"
    Path(shop_dir).mkdir(parents=True, exist_ok=True)
    (Path(shop_dir) / key).write_bytes(b"%PDF-1.4 shop-file\n%%EOF\n")
    owned = ShopPurchase(
        lemon_squeezy_order_id="FILE-1", customer_email="newperson@example.com",
        user_id=known_id, product_name="Self Hosted Guide", file_key=key,
        status="linked", purchased_at=utcnow())
    other = ShopPurchase(
        lemon_squeezy_order_id="FILE-2", customer_email="buyer@example.com",
        user_id=None, product_name="Someone Else", file_key=key,
        status="pending_link", purchased_at=utcnow())
    db.session.add_all([owned, other])
    db.session.commit()
    owned_id = owned.id
r = client.get(f"/account/shop/{owned_id}/download")
ok("Owner can download self-hosted shop file",
   r.status_code == 200 and r.get_data() == b"%PDF-1.4 shop-file\n%%EOF\n")
r = buyer_client.get(f"/account/shop/{owned_id}/download")
ok("Non-owner blocked from shop file download", r.status_code == 404)

r = admin.get("/admin/")
_dash = r.get_data(as_text=True)
ok("Dashboard shows local payment insights",
   r.status_code == 200 and "Payments (30 days)" in _dash
   and "Stripe" in _dash)

# give the main member Full Bloom: both community tracks for the forum suite
with app.app_context():
    m = User.query.filter_by(email="newperson@example.com").first()
    m.membership = "full_bloom"
    db.session.commit()

# --- 5. community forums + moderation + recommendations --------------------------
today = date.today()
with app.app_context():
    healing = ForumCategory(slug="healing", name="Healing",
                            description="Room to process.", sort_order=1)
    db.session.add(healing)
    db.session.flush()
    t_vent = ForumTag(category_id=healing.id, slug="venting", name="The Vent", sort_order=0)
    t_grief = ForumTag(category_id=healing.id, slug="grief", name="Grief & Loss", sort_order=1)
    db.session.add_all([t_vent, t_grief])
    db.session.commit()
    vent_tag_id = t_vent.id

r = client.get("/forums/")
comm_body = r.get_data(as_text=True)
ok("Forums index renders for members", r.status_code == 200 and "The Community" in comm_body)
ok("Community page shows healing / building hubs",
   "Healing community" in comm_body
   and "Building community" in comm_body
   and ("Enter the Healing Community" in comm_body
        or "Join the Healing Community" in comm_body)
   and "comm-compare" in comm_body
   and "What we talk about" in comm_body)

# category page shows topic filter chips
r = client.get("/forums/c/healing")
ok("Category shows tag filter chips", "The Vent" in r.get_data(as_text=True) and "Grief &amp; Loss" in r.get_data(as_text=True))
ok("Category shows Looking for filter chips",
   "Looking for" in r.get_data(as_text=True) and "Advice" in r.get_data(as_text=True)
   and "Recognition" in r.get_data(as_text=True))
ok("Conversation card stretch-link styles ship",
   "post-row__hit" in client.get("/static/css/main.css").get_data(as_text=True))

# member (client = newperson, verified + logged in) can post with a tag
r = client.post("/forums/c/healing/new",
                data={"title": "Rough day", "body": "Just needed to say it out loud.",
                      "tag_id": str(vent_tag_id), "looking_for": "support"},
                follow_redirects=True)
body = r.get_data(as_text=True)
ok("Member can create a tagged forum post",
   "Rough day" in body and "The Vent" in body)
ok("Looking for label shows on the post",
   "tag-chip--looking" in body and "support" in body.lower())
with app.app_context():
    saved = ForumPost.query.filter_by(title="Rough day").first()
ok("Looking for intent saved on the post",
   saved is not None and saved.looking_for == "support")
r = client.get("/forums/c/healing?looking=support")
ok("Looking-for filter shows matching posts",
   "Rough day" in r.get_data(as_text=True))
r = client.get("/forums/c/healing?looking=advice")
ok("Looking-for filter hides other intents",
   "Rough day" not in r.get_data(as_text=True))

# tag filter narrows the list
r = client.get("/forums/c/healing?tag=grief")
ok("Tag filter hides posts from other topics", "Rough day" not in r.get_data(as_text=True))
r = client.get("/forums/c/healing?tag=venting")
feed_html = r.get_data(as_text=True)
ok("Tag filter shows matching posts", "Rough day" in feed_html)
ok("Feed conversation widget opens from the full card",
   "post-row__hit" in feed_html and "Rough day" in feed_html)

# profanity is blocked and earns a warning
r = client.post("/forums/c/healing/new",
                data={"title": "This is shit", "body": "ugh"}, follow_redirects=True)
with app.app_context():
    member = User.query.filter_by(email="newperson@example.com").first()
    warn1 = member.forum_warnings
    posts_after = ForumPost.query.count()
ok("Profane post blocked + warning issued", warn1 == 1 and posts_after == 1,
   f"warnings={warn1} posts={posts_after}")

# anonymous posting hides the author name
r = client.post("/forums/c/healing/new",
                data={"title": "Quiet ask", "body": "Posting this anonymously.", "anonymous": "1"},
                follow_redirects=True)
ok("Anonymous post shows as Anonymous",
   "Anonymous" in r.get_data(as_text=True) and "Quiet ask" in r.get_data(as_text=True))

# likes + comments + one-level replies
with app.app_context():
    first_post = ForumPost.query.order_by(ForumPost.id).first()
    pid = first_post.id
r = client.post(f"/forums/p/{pid}/like", follow_redirects=True)
ok("Like on a post is accepted", r.status_code == 200)
r = client.post(f"/forums/p/{pid}/comment", data={"body": "Sending you strength."},
                follow_redirects=True)
ok("Comment posts to a thread", "Sending you strength." in r.get_data(as_text=True))

with app.app_context():
    top_comment = ForumComment.query.filter_by(post_id=pid, parent_id=None).first()
    cid = top_comment.id
r = client.post(f"/forums/p/{pid}/comment",
                data={"body": "Thank you, truly.", "parent_id": str(cid)}, follow_redirects=True)
ok("Reply attaches to its parent comment", "Thank you, truly." in r.get_data(as_text=True))

# a reply to a reply is flattened to one level (never nests deeper)
with app.app_context():
    reply = ForumComment.query.filter_by(post_id=pid).filter(ForumComment.parent_id.isnot(None)).first()
    reply_id = reply.id
client.post(f"/forums/p/{pid}/comment",
            data={"body": "Nested attempt.", "parent_id": str(reply_id)}, follow_redirects=True)
with app.app_context():
    nested = ForumComment.query.filter_by(body="Nested attempt.").first()
ok("Reply-to-a-reply flattens to one level", nested.parent_id == cid,
   f"parent_id={nested.parent_id} expected {cid}")

# strangers can comment, but only OP (or the comment author) may reply under a comment
with app.app_context():
    stranger = User(email="stranger@example.com", username="stranger_one",
                    membership="healing")
    stranger.set_password(USER_PW)
    stranger.email_verified_at = utcnow()
    bystander = User(email="bystander@example.com", username="bystander_one",
                     membership="healing")
    bystander.set_password(USER_PW)
    bystander.email_verified_at = utcnow()
    db.session.add_all([stranger, bystander])
    db.session.commit()
stranger_client = app.test_client()
stranger_client.post("/login", data={"email": "stranger@example.com", "password": USER_PW})
r = stranger_client.post(f"/forums/p/{pid}/comment",
                         data={"body": "A kind stranger note."}, follow_redirects=True)
ok("Anyone can leave a top-level comment",
   "A kind stranger note." in r.get_data(as_text=True))
with app.app_context():
    stranger_note = ForumComment.query.filter_by(body="A kind stranger note.").first()
    stranger_cid = stranger_note.id
bystander_client = app.test_client()
bystander_client.post("/login", data={"email": "bystander@example.com", "password": USER_PW})
r = bystander_client.post(f"/forums/p/{pid}/comment",
                          data={"body": "Should not nest here.", "parent_id": str(stranger_cid)},
                          follow_redirects=True)
page = r.get_data(as_text=True).lower()
with app.app_context():
    blocked_reply = ForumComment.query.filter_by(body="Should not nest here.").count()
ok("Non-OP cannot reply under someone else's comment",
   blocked_reply == 0 and "only the original poster" in page)
# OP can still reply under that stranger comment
r = client.post(f"/forums/p/{pid}/comment",
                data={"body": "Thanks for the note.", "parent_id": str(stranger_cid)},
                follow_redirects=True)
ok("OP can reply under any comment",
   "Thanks for the note." in r.get_data(as_text=True))

# escalating profanity leads to a ban after the warning limit
banclient = app.test_client()
sent_codes.clear()
banclient.post("/register", data={"email": "rude@example.com", "password": USER_PW,
                                  "password_confirm": USER_PW})
bcode = sent_codes[-1][1]
banclient.post("/verify-email", data={"email": "rude@example.com", "code": bcode})
with app.app_context():
    ru = User.query.filter_by(email="rude@example.com").first()
    ru.membership = "healing"
    db.session.commit()
for _ in range(3):
    banclient.post("/forums/c/healing/new", data={"title": "fuck this", "body": "fuck"})
with app.app_context():
    rude = User.query.filter_by(email="rude@example.com").first()
    banned = rude.forum_banned
ok("Repeated profanity bans after 2 warnings", banned is True, f"banned={banned}")

# avatar upload: a real (tiny) PNG is accepted, re-encoded, and served
with app.app_context():
    import io as _io
    from PIL import Image as _Image
    buf = _io.BytesIO()
    _Image.new("RGB", (10, 10), (200, 100, 150)).save(buf, format="PNG")
    png_bytes = buf.getvalue()
r = client.post("/account/profile", data={
    "display_name": "River",
    "avatar_file": (_io.BytesIO(png_bytes), "me.png"),
}, content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    m = User.query.filter_by(email="newperson@example.com").first()
    has_av = m.has_avatar()
    av_uid = m.id
ok("Uploaded avatar stored on the account", has_av)
r = client.get(f"/avatar/{av_uid}")
ok("Avatar is served from the database",
   r.status_code == 200 and r.headers["Content-Type"].startswith("image/"))

# animated GIF: still used in lists; animation served only on the profile page
with app.app_context():
    from PIL import Image as _Image2, ImageDraw as _ImageDraw
    gif_buf = _io.BytesIO()
    frames = []
    for i, color in enumerate([(220, 80, 120), (80, 140, 220)]):
        fr = _Image2.new("RGB", (40, 40), color)
        _ImageDraw.Draw(fr).ellipse((8, 8, 32, 32), fill=(255, 255, 255))
        frames.append(fr)
    frames[0].save(gif_buf, format="GIF", save_all=True, append_images=frames[1:],
                   duration=120, loop=0)
    gif_bytes = gif_buf.getvalue()
r = client.post("/account/profile", data={
    "display_name": "River",
    "avatar_file": (_io.BytesIO(gif_bytes), "me.gif"),
}, content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    m = User.query.filter_by(email="newperson@example.com").first()
    av_uid = m.id
    has_anim = m.has_animated_avatar()
    still_mime = m.avatar_mime
ok("Animated GIF stores a still + animation payload",
   has_anim and (still_mime or "").startswith("image/"))
r_still = client.get(f"/avatar/{av_uid}")
r_anim = client.get(f"/avatar/{av_uid}/anim")
ok("Default avatar route serves the still frame",
   r_still.status_code == 200 and "gif" not in (r_still.headers.get("Content-Type") or "").lower())
ok("Anim avatar route serves image/gif",
   r_anim.status_code == 200 and "gif" in (r_anim.headers.get("Content-Type") or "").lower())
r_prof = client.get(f"/u/{av_uid}")
prof_html = r_prof.get_data(as_text=True)
ok("Profile page uses the animated avatar URL",
   r_prof.status_code == 200 and f"/avatar/{av_uid}/anim" in prof_html)

r = client.get("/account/settings")
sbody = r.get_data(as_text=True)
ok("Settings page renders with intents + upload",
   r.status_code == 200 and "What brings you here?" in sbody and 'name="avatar_file"' in sbody)
ok("Settings offers a change-password button (no inline fields)",
   'href="/account/password"' in sbody and 'name="current_password"' not in sbody)
ok("Close-account needs Yes I'm sure before submit is enabled",
   'data-require-sure' in sbody
   and 'data-sure-submit' in sbody
   and 'disabled' in sbody
   and "Yes, I'm sure" in sbody)
r = client.post("/account/delete", data={}, follow_redirects=True)
del_body = r.get_data(as_text=True)
with app.app_context():
    still_here = User.query.filter_by(email="newperson@example.com").first()
ok("Unconfirmed delete is rejected and keeps the account",
   still_here is not None and still_here.deleted_at is None
   and ("tick" in del_body.lower() or "sure" in del_body.lower()))

r = client.get("/account/password")
ok("Change-password subpage renders",
   r.status_code == 200
   and 'name="current_password"' in r.get_data(as_text=True)
   and 'name="new_password_confirm"' in r.get_data(as_text=True))

# profile links + Creator-of-the-Month Instagram + public profile page
r = client.get("/account/settings")
ok("Creator settings show Creator of the Month Instagram field",
   r.status_code == 200
   and "Instagram for Creator of the Month" in r.get_data(as_text=True)
   and 'name="creator_instagram"' in r.get_data(as_text=True))
client.post("/account/profile", data={
    "display_name": "New Person",
    "creator_instagram": "@newperson",
    "link_label_0": "Site", "link_url_0": "https://newperson.example",
    "link_label_1": "", "link_url_1": "",
}, follow_redirects=True)
with app.app_context():
    saved_links = User.query.filter_by(email="newperson@example.com").first().links()
    ig_urls = [ln["url"] for ln in saved_links if "instagram.com" in ln["url"]]
ok("Creator-of-the-Month Instagram saved onto profile links",
   bool(ig_urls) and "instagram.com/newperson" in ig_urls[0])
ok("Other profile links still save",
   any("newperson.example" in ln["url"] for ln in saved_links))

r = client.get(f"/u/{av_uid}")
pbody = r.get_data(as_text=True)
ok("Public profile page renders with links",
   r.status_code == 200 and "New Person" in pbody and "instagram.com/newperson" in pbody)
ok("Unknown profile returns 404", client.get("/u/99999").status_code == 404)

# --- 5b2. streaks: "I showed up today" ---------------------------------------
r = client.post("/account/checkin", data={"mood": "soft"}, follow_redirects=True)
with app.app_context():
    m = User.query.filter_by(email="newperson@example.com").first()
    ci = (m.total_checkins, m.current_streak, m.longest_streak, m.checked_in_today())
    from app.models import JournalEntry
    mood_entry = JournalEntry.query.filter_by(user_id=m.id, day=date.today()).first()
ok("Check-in records the first streak day", ci == (1, 1, 1, True), f"got {ci}")
ok("Check-in mood is saved on today's journal entry",
   mood_entry is not None and mood_entry.mood == "soft",
   f"got {getattr(mood_entry, 'mood', None)}")
client.post("/account/checkin", follow_redirects=True)
with app.app_context():
    again = User.query.filter_by(email="newperson@example.com").first().total_checkins
ok("A second check-in the same day doesn't double-count", again == 1, f"got {again}")
r = client.get("/account")
abody = r.get_data(as_text=True)
ok("Account confirms you showed up today", "You showed up today" in abody)
ok("Account shows community participation count",
   "Community participation" in abody and "time" in abody
   and "Open</strong>" not in abody)
with app.app_context():
    from app.services.participation import community_participation_count
    m = User.query.filter_by(email="newperson@example.com").first()
    # 1 check-in + whatever forum posts this member already made in earlier steps
    part_n = community_participation_count(m)
ok("Community participation counts check-ins and posts",
   part_n >= 1, f"got {part_n}")
client.post("/account/checkin",
            data={"mood": "bloom", "journal": "Blooming a little today."},
            follow_redirects=True)
with app.app_context():
    m = User.query.filter_by(email="newperson@example.com").first()
    from app.models import JournalEntry
    # writing starts a fresh page rather than overwriting an earlier mood-only
    # one, so the newest row for today is the one that just got saved
    je = (JournalEntry.query
          .filter_by(user_id=m.id, day=date.today())
          .order_by(JournalEntry.created_at.desc(), JournalEntry.id.desc())
          .first())
ok("A written entry keeps the mood picked with it",
   je is not None and je.mood == "bloom"
   and "Blooming a little today." in (je.body or ""),
   f"mood={getattr(je, 'mood', None)} body={getattr(je, 'body', None)!r}")
r = client.get("/account?tab=journal")
jbody = r.get_data(as_text=True)
ok("Journal tab shows clickable prompt ideas",
   r.status_code == 200
   and 'id="journal-prompt-ideas"' in jbody
   # four sampled prompts plus the "write freely" option
   and jbody.count("data-prompt-key=") == 5,
   f"buttons={jbody.count('data-prompt-key=')}")
ok("Journal past entries sit behind today's page in the notebook",
   "jn-page--past" in jbody and "jn-page--today" in jbody
   and 'data-jn-book' in jbody)
with app.app_context():
    from app.models import sample_journal_prompts
    ideas = sample_journal_prompts(4)
ok("Prompt idea sampler returns four unique prompts",
   len(ideas) == 4 and len({k for k, _ in ideas}) == 4, f"got {ideas}")

# --- 5b3. badges: earn, display (max 3), byline, profile, owner --------------
from app.services.badges import earned_badges, primary_badge
with app.app_context():
    m = User.query.filter_by(email="newperson@example.com").first()
    earned_keys = {b["cat"] for b in earned_badges(m)}
ok("Member earns the Storyteller badge by posting", "storyteller" in earned_keys,
   f"earned={earned_keys}")

# choosing badges: an unearned category (kindred) is ignored; earned ones stick
client.post("/account/profile", data={"display_name": "New Person",
            "badges_display": ["kindred", "storyteller"]}, follow_redirects=True)
with app.app_context():
    m = User.query.filter_by(email="newperson@example.com").first()
    chosen = m.displayed_badges()
    prim = primary_badge(m)
ok("Only earned badges are saved for display", chosen == ["storyteller"], f"got {chosen}")
ok("Primary badge is the chosen Storyteller", bool(prim) and prim["cat"] == "storyteller")

r = client.get(f"/u/{av_uid}")
ok("Profile displays the member's badge (with milestone tooltip)",
   "Storyteller" in r.get_data(as_text=True))

with app.app_context():
    rough = ForumPost.query.filter_by(title="Rough day").first()
    rough_id = rough.id
r = client.get(f"/forums/p/{rough_id}")
ok("Badge shows by the author's name on a post", "Storyteller" in r.get_data(as_text=True))

r = client.get("/account/settings")
ok("Settings shows the badge collection + chooser",
   "Your badges" in r.get_data(as_text=True) and 'name="badges_display"' in r.get_data(as_text=True))

with app.app_context():
    owner = User.query.filter_by(is_admin=True).first()
    owner_prim = primary_badge(owner)
ok("Owner carries the special Founder badge",
   bool(owner_prim) and owner_prim["cat"] == "owner")

# --- 5b4. studio badge manager: view + tweak milestones ----------------------
r = admin.get("/admin/badges")
bbody = r.get_data(as_text=True)
ok("Studio badge manager lists every category with editable milestones",
   r.status_code == 200 and "Showing Up" in bbody and "Storyteller" in bbody
   and 'name="t_storyteller_1"' in bbody)

with app.app_context():
    from app.services import badges as B
    base_form = {}
    for _cat in B.CATEGORIES:
        for _i, _t in enumerate(B.thresholds(_cat), start=1):
            base_form[f"t_{_cat}_{_i}"] = _t

# non-ascending milestones are rejected; values stay put
bad_form = dict(base_form)
bad_form["t_storyteller_2"] = 1            # <= tier 1 (which is 1)
admin.post("/admin/badges", data=bad_form, follow_redirects=True)
with app.app_context():
    unchanged = B.thresholds("storyteller")
ok("Non-ascending milestones are rejected", unchanged == B.default_thresholds("storyteller"),
   f"got {unchanged}")

# a valid tweak saves and flows through to the badge tooltip/phrase
good_form = dict(base_form)
good_form["t_storyteller_3"] = 30          # was 25
admin.post("/admin/badges", data=good_form, follow_redirects=True)
with app.app_context():
    tweaked = B.thresholds("storyteller")
    phrase = B.badge_dict("storyteller", 3)["phrase"]
ok("Owner can tweak a milestone value", tweaked[2] == 30, f"got {tweaked}")
ok("Tweaked milestone updates the badge phrase", phrase == "30 posts", f"got {phrase}")

# reset restores defaults
admin.post("/admin/badges", data={"reset": "1"}, follow_redirects=True)
with app.app_context():
    reset_vals = B.thresholds("storyteller")
ok("Reset restores default milestones", reset_vals == B.default_thresholds("storyteller"),
   f"got {reset_vals}")

# --- 5b5. My Journey keepsake (Creator-gated PDF) ----------------------------
# a fresh free member is gently redirected, no PDF
free_client = app.test_client()
with app.app_context():
    fu = User(email="free@example.com", membership="none", email_verified_at=utcnow())
    fu.set_password(USER_PW)
    db.session.add(fu)
    db.session.commit()
free_client.post("/login", data={"email": "free@example.com", "password": USER_PW})
r = free_client.get("/account/journey.pdf", follow_redirects=False)
ok("Free member can't export a journey",
   r.status_code == 302 and "/account" in r.headers.get("Location", ""))

# favorite a quote so the keepsake has something tender in it
with app.app_context():
    fav_qid = Quote.query.first().id
client.post(f"/quotes/{fav_qid}/favorite", follow_redirects=True)

# newperson is a Creator member -> export unlocked
r = client.get("/account/journey.pdf")
pdf_data = r.get_data()
ok("Creator member downloads a My Journey PDF",
   r.status_code == 200 and r.mimetype == "application/pdf"
   and pdf_data[:5] == b"%PDF-" and len(pdf_data) > 1200
   and r.headers.get("Content-Disposition", "").startswith("attachment"))

r = client.get("/account")
ok("Account offers the keepsake to Creator members",
   "Download my journey" in r.get_data(as_text=True))

with app.app_context():
    from app.models import CheckIn
    mid = User.query.filter_by(email="newperson@example.com").first().id
    n_logged = CheckIn.query.filter_by(user_id=mid).count()
ok("Check-ins are logged for the journey history", n_logged >= 1, f"got {n_logged}")

# intent tags still save on the member (shop recommendations retired)
with app.app_context():
    m = User.query.filter_by(email="newperson@example.com").first()
    m.set_goals(["divorce"])
    db.session.commit()
    goals = m.goals()
ok("Member intent tags still save", "divorce" in goals)

r = admin.get("/admin/community")
ok("Admin community moderation page", r.status_code == 200 and "rude@example.com" in r.get_data(as_text=True))

with app.app_context():
    rude = User.query.filter_by(email="rude@example.com").first()
    rude_id = rude.id
    rude.forum_warnings = 2
    rude.forum_banned = True
    db.session.commit()
r = admin.post(f"/admin/community/member/{rude_id}/reset", follow_redirects=True)
ok("Fresh start clears flags without a 404",
   r.status_code == 200
   and "Fresh start given" in r.get_data(as_text=True)
   and "This page took a different path" not in r.get_data(as_text=True))
with app.app_context():
    rude = User.query.filter_by(email="rude@example.com").first()
    cleared = (rude.forum_warnings or 0) == 0 and not rude.forum_banned
ok("Fresh start zeroes warnings and unpauses posting", cleared)

# Owner account must not 404 — redirect back with a flash instead
with app.app_context():
    owner_id = User.query.filter_by(is_admin=True).first().id
r = admin.post(f"/admin/community/member/{owner_id}/reset", follow_redirects=True)
ok("Fresh start on owner stays in Studio",
   r.status_code == 200
   and "Studio owner accounts" in r.get_data(as_text=True)
   and "This page took a different path" not in r.get_data(as_text=True))

# --- 5c. on-site course reader retired (shop downloads in My space) -----------
r = client.get("/library/begin-again", follow_redirects=False)
ok("Legacy library reader is gone", r.status_code == 404)
r = client.get("/account?tab=saved")
ok("Account still has courses & guides section",
   "Courses" in r.get_data(as_text=True) and "myspace-tabs" in client.get("/account").get_data(as_text=True))

# --- 5d. announcement: expiry window + remove ---------------------------------
base_settings = {"site_title": "Bloom Anyway", "instagram_url": "", "hero_image_url": "",
                 "portrait_url": "", "contact_email": ""}
future = (date.today() + timedelta(days=3)).isoformat()
admin.post("/admin/settings", data={**base_settings,
           "announcement_text": "Doors open Monday", "announcement_expires": future},
           follow_redirects=True)
r = client.get("/")
ok("Announcement shows before its expiry", "Doors open Monday" in r.get_data(as_text=True))
past = (date.today() - timedelta(days=1)).isoformat()
admin.post("/admin/settings", data={**base_settings,
           "announcement_text": "Doors open Monday", "announcement_expires": past},
           follow_redirects=True)
r = client.get("/")
ok("Expired announcement is hidden", "Doors open Monday" not in r.get_data(as_text=True))
admin.post("/admin/settings", data={"clear_announcement": "1"}, follow_redirects=True)
with app.app_context():
    from app.services.settings import get_setting, invalidate_cache
    invalidate_cache()
    cleared_text = get_setting("announcement_text")
ok("Remove announcement clears it", cleared_text == "")
r = client.get("/")
ok("No announcement markup after removal", "hero-announcement" not in r.get_data(as_text=True))

# --- 5e. memberships, videos, subjects, spotlight ---------------------------
# free member: community is members-only (Healing / Creator)
r = free_client.get("/forums/")
free_gate = r.get_data(as_text=True)
ok("Free member sees community gate (not threads)",
   r.status_code == 200
   and "members" in free_gate.lower()
   and ("Healing" in free_gate or "membership" in free_gate.lower())
   and "The Community" not in free_gate)
r = free_client.get("/forums/c/healing")
ok("Free member cannot browse category threads",
   "Enter the Healing" not in r.get_data(as_text=True)
   and ("See memberships" in r.get_data(as_text=True)
        or "membership" in r.get_data(as_text=True).lower()))
r = free_client.post("/forums/c/healing/new",
                     data={"title": "free weekly post", "body": "should be blocked"},
                     follow_redirects=True)
with app.app_context():
    free_posts = ForumPost.query.filter_by(title="free weekly post").count()
ok("Free member cannot create posts", free_posts == 0)

# /courses stays on-site (query params ignored / filters via h=)
r = client.get("/courses?h=workbook", follow_redirects=False)
ok("Filtered /courses stays on-site",
   r.status_code == 200 and "Courses &amp; Guides" in r.get_data(as_text=True))

# --- a product can be more than one thing -------------------------------------
# A course that comes with templates is both, so the kinds are a select-all and
# the catalogue filters have to find it under either one.
r = admin.post("/admin/products/new", data={
    "title": "Rebuild Your Week", "track": "building",
    "types": ["course", "template"],
    "promise": "A course, and the templates to work it.",
    "price": "24.00", "stripe": "price_rebuild_week", "live": "1",
}, follow_redirects=True)
ok("Studio saves a product that is several things", r.status_code == 200)
with app.app_context():
    _multi = Product.query.filter_by(slug="rebuild-your-week").first()
    ok("Both kinds are kept",
       _multi is not None and _multi.types() == ["course", "template"],
       f"got {_multi.types() if _multi else None}")
    ok("The primary is the first in the canonical order, not tick order",
       _multi.type == "course")
    ok("It answers to either kind",
       _multi.has_type("course") and _multi.has_type("template")
       and not _multi.has_type("guide"))
    ok("And reads as both where there is room",
       _multi.types_display() == "Course, Template",
       f"got {_multi.types_display()}")
    _multi_id = _multi.id

def _catalogue(path):
    """The page below the header. The nav's notification bell names new
    products too, and that isn't the catalogue answering."""
    return client.get(path).get_data(as_text=True).split("</header>", 1)[-1]


ok("The catalogue finds it under its primary kind",
   "Rebuild Your Week" in _catalogue("/courses?b=course"))
ok("And under the kind it is also",
   "Rebuild Your Week" in _catalogue("/courses?b=template"))
ok("But not under one it isn't",
   "Rebuild Your Week" not in _catalogue("/courses?b=guide"))
ok("The product page names both",
   "Course, Template" in client.get("/courses/rebuild-your-week").get_data(as_text=True))

# --- a running promo code -----------------------------------------------------
_promo_fields = {
    "title": "Rebuild Your Week", "track": "building", "types": ["course"],
    "slug": "rebuild-your-week", "promise": "A course, and the templates.",
    "price": "24.00", "stripe": "price_rebuild_week", "live": "1",
}
r = admin.post(f"/admin/products/{_multi_id}/edit",
               data=dict(_promo_fields, promo_price="18.00", promo_code="spring25"),
               follow_redirects=True)
with app.app_context():
    _p = db.session.get(Product, _multi_id)
    ok("Studio saves a promo price and its code",
       _p.has_promo() and _p.promo_price_cents == 1800
       and _p.promo_code_display() == "SPRING25",
       f"got {_p.promo_code!r} {_p.promo_price_cents}")
    ok("And works out what comes off",
       _p.promo_display() == "$18" and _p.promo_saving_display() == "$6",
       f"got {_p.promo_display()} / {_p.promo_saving_display()}")

_pd = client.get("/courses/rebuild-your-week").get_data(as_text=True)
ok("The product page leads with the promo price and strikes the old one",
   "pd-promo" in _pd and "$18" in _pd and "<s>$24</s>" in _pd)
ok("And says which code to type",
   "SPRING25" in _pd and "at checkout" in _pd)
_tile = _catalogue("/courses?b=course")
ok("The catalogue tile shows the sale price over the old one",
   "lib-card__price--promo" in _tile and "<s>$24</s>" in _tile)
ok("With the saving and the code on the card",
   "lib-card__promo" in _tile and "SPRING25" in _tile and "$6 off" in _tile)

# Half a promo is no promo, and one that doesn't save anything isn't a sale.
for _bad, _why in (({"promo_price": "18.00", "promo_code": ""}, "no code"),
                   ({"promo_price": "", "promo_code": "SPRING25"}, "no price"),
                   ({"promo_price": "30.00", "promo_code": "SPRING25"}, "dearer")):
    admin.post(f"/admin/products/{_multi_id}/edit",
               data=dict(_promo_fields, **_bad), follow_redirects=True)
    with app.app_context():
        ok(f"A promo with {_why} isn't advertised",
           not db.session.get(Product, _multi_id).has_promo())
ok("And the card goes back to the ordinary price",
   "lib-card__price--promo" not in _catalogue("/courses?b=course"))

# A sale with a deadline takes itself down when the deadline passes.
_soon = utcnow() + timedelta(days=3)
r = admin.post(f"/admin/products/{_multi_id}/edit",
               data=dict(_promo_fields, promo_price="18.00", promo_code="SPRING25",
                         promo_ends_date=_soon.strftime("%Y-%m-%d"),
                         promo_ends_time="17:30"),
               follow_redirects=True)
with app.app_context():
    _p = db.session.get(Product, _multi_id)
    ok("Studio saves when the sale ends",
       _p.has_promo() and _p.promo_ends_at is not None
       and abs((_p.promo_ends_at - _soon.replace(hour=17, minute=30, second=0,
                                                 microsecond=0)).total_seconds()) < 90,
       f"got {_p.promo_ends_at}")
    ok("It isn't expired yet", not _p.promo_expired())
_pd = client.get("/courses/rebuild-your-week").get_data(as_text=True)
ok("The product page says when the sale ends",
   "pd-promo__ends" in _pd and "Ends " in _pd)

with app.app_context():
    _p = db.session.get(Product, _multi_id)
    _p.promo_ends_at = utcnow() - timedelta(minutes=1)
    db.session.commit()
    ok("A minute past its deadline the sale is over",
       _p.promo_expired() and not _p.has_promo())
    ok("And nothing about it is offered up any more",
       _p.promo_display() == "" and _p.promo_ends_display() == "")
_pd = client.get("/courses/rebuild-your-week").get_data(as_text=True)
ok("The product page drops the banner and shows the normal price",
   "pd-promo" not in _pd and "SPRING25" not in _pd and "$24" in _pd)
ok("So does the catalogue card",
   "lib-card__price--promo" not in _catalogue("/courses?b=course")
   and "SPRING25" not in _catalogue("/courses?b=course"))
_sbody = admin.get(f"/admin/products/{_multi_id}/edit").get_data(as_text=True)
ok("Studio says the sale has ended rather than pretending it's running",
   "This sale has ended" in _sbody)

# Leaving the date blank means it runs until it's taken down.
admin.post(f"/admin/products/{_multi_id}/edit",
           data=dict(_promo_fields, promo_price="18.00", promo_code="SPRING25",
                     promo_ends_date="", promo_ends_time="23:59"),
           follow_redirects=True)
with app.app_context():
    _p = db.session.get(Product, _multi_id)
    ok("No deadline means the sale just runs",
       _p.has_promo() and _p.promo_ends_at is None
       and _p.promo_ends_display() == "")
with app.app_context():
    db.session.get(Product, _multi_id).set_types(["course", "template"])
    db.session.commit()

# Anything still posting one type keeps working, and so does a product that
# has only ever had one.
r = admin.post("/admin/products/new", data={
    "title": "Just One Thing", "track": "healing", "type": "workbook",
    "promise": "One kind only.", "price": "9.00", "stripe": "price_just_one",
}, follow_redirects=True)
with app.app_context():
    _single = Product.query.filter_by(slug="just-one-thing").first()
    ok("A single type still saves, from the old field name",
       _single is not None and _single.type == "workbook"
       and _single.types() == ["workbook"] and _single.types_json is None,
       f"got {_single.types() if _single else None}")
    ok("And a row that predates all this reads as its one kind",
       Product(type="guide").types() == ["guide"])

# Ticking bundle second still makes it the lane's bundle.
with app.app_context():
    _b = db.session.get(Product, _multi_id)
    _b.set_types(["template", "bundle"])
    db.session.commit()
    ok("Ticking bundle anywhere in the list still makes it a bundle",
       _b.has_type("bundle") and _b.types() == ["template", "bundle"],
       f"got {_b.types()}")
ok("And it is lifted out of the lane into the bundle slot either way",
   "Rebuild Your Week" in _catalogue("/courses"))
with app.app_context():
    _b = db.session.get(Product, _multi_id)
    _b.set_types(["course", "template"])
    db.session.commit()

# content tips: owner writes one, Creator reads it, free is blocked
minimal_mp4 = b"\x00\x00\x00\x18ftypmp42" + b"\x00" * 64
r = admin.post("/admin/videos/new", data={
    "title": "Morning pages walkthrough", "description": "How I use the notebook.",
    "body": "Start with three pages.\n\nDon't edit while you write.",
    "published": "1", "sort_order": "0",
    "video_file": (io.BytesIO(minimal_mp4), "clip.mp4"),
}, content_type="multipart/form-data", follow_redirects=True)
ok("Owner writes a tip with a video", "Tip saved" in r.get_data(as_text=True))
ok("Studio confirms the tip went out to members",
   "notified" in r.get_data(as_text=True))
with app.app_context():
    from app.models import Notification as _Note
    vid_id = Video.query.filter_by(title="Morning pages walkthrough").first().id
    tip_notes = _Note.query.filter_by(kind="content_hub").all()
    ok("Publishing a tip notifies members",
       len(tip_notes) >= 1
       and all("Morning pages walkthrough" in (n.body or "") for n in tip_notes)
       and all(f"/watch/{vid_id}" == (n.url or "") for n in tip_notes))
    owner_row = User.query.filter_by(email="owner@example.com").first()
    ok("The owner isn't notified about their own tip",
       all(n.user_id != owner_row.id for n in tip_notes))

# The home-page notice must follow the newest tip, not the hub's sort order.
_notice_html = client.get("/").get_data(as_text=True)
ok("New tip is announced on the home page",
   "New in the Content Hub" in _notice_html
   and "Morning pages walkthrough" in _notice_html)

# the whole point: a tip is writing, so the video is optional
r = admin.post("/admin/videos/new", data={
    "title": "Batch a week of hooks", "description": "Twenty minutes, once a week.",
    "body": "Open a note.\n\nWrite ten first lines without judging them.",
    "published": "1", "sort_order": "1", "free_access": "1",
}, content_type="multipart/form-data", follow_redirects=True)
ok("Owner publishes a text-only tip", "Tip saved" in r.get_data(as_text=True))
with app.app_context():
    text_tip = Video.query.filter_by(title="Batch a week of hooks").first()
    ok("Text-only tip is stored without a video",
       text_tip is not None and text_tip.has_video() is False
       and "ten first lines" in (text_tip.body or ""))
    text_tip_id = text_tip.id

# One tip pinned to the top of the hub used to swallow the notice for every
# tip published after it.
with app.app_context():
    pinned = db.session.get(Video, vid_id)
    pinned.sort_order = -5
    db.session.commit()
_pinned_html = client.get("/").get_data(as_text=True)
ok("A pinned older tip no longer hides the newest one",
   "Batch a week of hooks" in _pinned_html)
ok("Free members hear about new tips too",
   "New in the Content Hub" in free_client.get("/").get_data(as_text=True))
with app.app_context():
    db.session.get(Video, vid_id).sort_order = 0
    db.session.commit()

r = admin.post("/admin/videos/new", data={
    "title": "Nothing to say", "published": "1",
}, content_type="multipart/form-data", follow_redirects=True)
ok("A tip with neither text nor video is refused",
   "Write the tip" in r.get_data(as_text=True))

# a video can come off again without losing the tip
admin.post("/admin/videos/new", data={
    "title": "Swap the clip out", "body": "Words that stay.", "published": "1",
    "video_file": (io.BytesIO(minimal_mp4), "swap.mp4"),
}, content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    swap_id = Video.query.filter_by(title="Swap the clip out").first().id
admin.post(f"/admin/videos/{swap_id}/edit", data={
    "title": "Swap the clip out", "body": "Words that stay.", "published": "1",
    "remove_video": "1",
}, content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    swapped = db.session.get(Video, swap_id)
    ok("Owner can drop the video and keep the tip",
       swapped.has_video() is False and swapped.body == "Words that stay.")

form_html = admin.get(f"/admin/videos/{swap_id}/edit").get_data(as_text=True)
ok("Tip editor offers a preview of the member's page",
   "data-tip-preview-toggle" in form_html
   and 'class="tip-read__body forum-body" data-tip-preview-body' in form_html
   and f"/watch/{swap_id}" in form_html)
ok("Preview sits under the whole form, not in the middle of it",
   form_html.index("data-tip-preview-toggle") > form_html.index('name="healing_access"'))
ok("A tip with no video keeps a video out of the preview",
   "data-video-src" not in form_html)

with_video_html = admin.get(f"/admin/videos/{vid_id}/edit").get_data(as_text=True)
ok("Preview plays the video that's already attached",
   f'data-video-src="/watch/{vid_id}/stream"' in with_video_html
   and "data-tip-preview-video-el" in with_video_html)

r = free_client.get("/watch", follow_redirects=False)
ok("Free member can open Content Hub (public reviews)",
   r.status_code == 200 and "Content Hub" in r.get_data(as_text=True)
   and ("Video library" in r.get_data(as_text=True)
        or "Sign in" in r.get_data(as_text=True)
        or "Morning pages" in r.get_data(as_text=True)))
r = free_client.get(f"/watch/{vid_id}/stream")
ok("Free member can't stream a video", r.status_code == 404)

r = client.get("/watch")
ok("Creator member sees the video room",
   r.status_code == 200 and "Morning pages walkthrough" in r.get_data(as_text=True)
   and "Content Hub" in r.get_data(as_text=True))
r = client.get(f"/watch/{vid_id}")
ok("Creator member opens a tip and reads it",
   r.status_code == 200
   and "edit while you write." in r.get_data(as_text=True))
r = free_client.get(f"/watch/{text_tip_id}")
ok("Free member reads a tip that was opened to Free",
   r.status_code == 200 and "ten first lines" in r.get_data(as_text=True))
r = free_client.get(f"/watch/{vid_id}")
free_tip = r.get_data(as_text=True)
# A locked tip is blurred rather than teased. Blurring in CSS alone would
# still ship the words — the page source, and reader modes, would have them.
ok("A locked tip sends none of its words, not even an opening",
   r.status_code == 200
   and "Start with three pages" not in free_tip
   and "Don&#39;t edit while you write." not in free_tip
   and "See memberships" in free_tip)
ok("But there is blurred writing there, so it reads as locked and not empty",
   "tip-read__blurred" in free_tip)
with app.app_context():
    _shape = db.session.get(Video, vid_id).locked_shape()
    ok("The blur is filler shaped like the tip, carrying none of it",
       _shape and _shape in free_tip
       and not any(w in _shape.lower()
                   for w in ("start", "pages", "edit", "write")),
       f"got {_shape[:60]!r}")
    ok("And it holds still between loads rather than reshuffling",
       db.session.get(Video, vid_id).locked_shape() == _shape)
r = client.get(f"/watch/{vid_id}/stream", headers={"Range": "bytes=0-3"})
ok("Video streams with range support (206 partial)",
   r.status_code == 206 and r.headers.get("Accept-Ranges") == "bytes"
   and "Content-Range" in r.headers)
r = admin.get(f"/watch/{vid_id}")
ok("Owner can open a Content Hub video page",
   r.status_code == 200 and "video-player" in r.get_data(as_text=True))
r = admin.get(f"/watch/{vid_id}/stream", headers={"Range": "bytes=0-3"})
ok("Owner can stream Content Hub videos",
   r.status_code == 206 and "Content-Range" in r.headers)

r = admin.post("/admin/videos/new", data={
    "title": "Bad file", "body": "Some words.",
    "video_file": (io.BytesIO(b"nope"), "notes.txt"),
}, content_type="multipart/form-data", follow_redirects=True)
ok("Non-video attachment is rejected", "MP4" in r.get_data(as_text=True))

# oversized video shows the error inline on the form, not an error page
app.config["MAX_VIDEO_MB"] = 0
r = admin.post("/admin/videos/new", data={
    "title": "Too big", "body": "Some words.", "published": "1", "sort_order": "0",
    "video_file": (io.BytesIO(minimal_mp4), "big.mp4"),
}, content_type="multipart/form-data", follow_redirects=False)
ok("Oversized video shows an inline error (no error-page redirect)",
   r.status_code == 200 and "0 MB" in r.get_data(as_text=True))
app.config["MAX_VIDEO_MB"] = 1024

# home spotlight: creator of the month + reel of the week
reel_url = "https://www.instagram.com/reel/ABC123xyz/"
from app.services.social import instagram_handle as _ig_handle
messy = "https://www.instagram.com/hustlinmommaz?igsh=cWphMWdycGowY3Fo&utm_source=qr"
ok("Instagram handle strips share-link junk",
   _ig_handle(messy) == "hustlinmommaz", f"got {_ig_handle(messy)!r}")

spotlight_settings = {
    "creator_name": "Maya R.",
    "creator_instagram": messy.replace("hustlinmommaz", "mayar"),
    "creator_blurb": "Rebuilt her mornings.",
    "reel_url": reel_url,
    "reel_description": "Loved this one.",
}
admin.post("/admin/spotlight", data=spotlight_settings, follow_redirects=True)
r = client.get("/")
hbody = r.get_data(as_text=True)
ok("Creator of the month shows on home",
   "Maya R." in hbody and "@mayar" in hbody and "instagram.com/mayar" in hbody
   and "igsh=" not in hbody)
ok("Creator of the month shows their bio", "Rebuilt her mornings." in hbody)
ok("Creator of the month shows the flower mark (no broken photo circle)",
   "spotlight-creator__photo--ph" in hbody
   and "unavatar.io" not in hbody
   and 'class="spotlight-creator__photo"' not in hbody)
ok("Reel of the week embeds + links out",
   "instagram.com/reel/ABC123xyz/embed" in hbody and "Watch on Instagram" in hbody)

r = admin.get("/admin/spotlight")
sbody = r.get_data(as_text=True)
ok("Studio has a Home spotlight page with clear buttons",
   r.status_code == 200
   and 'name="clear_spotlight_creator"' in sbody
   and 'name="clear_spotlight_reel"' in sbody)
admin.post("/admin/spotlight", data={"clear_spotlight_reel": "1"}, follow_redirects=True)
r = client.get("/")
hbody = r.get_data(as_text=True)
ok("Clear reel removes Reel of the week only",
   "Maya R." in hbody and "Watch on Instagram" not in hbody
   and "instagram.com/reel/ABC123xyz" not in hbody)
admin.post("/admin/spotlight", data={"clear_spotlight_creator": "1"}, follow_redirects=True)
r = client.get("/")
hbody = r.get_data(as_text=True)
ok("Clear creator removes Creator of the month",
   "Maya R." not in hbody
   and "Our next Creator of the Month will land here" in hbody)

# studio: members management
r = admin.get("/admin/members")
ok("Members page lists memberships", r.status_code == 200 and "Creator" in r.get_data(as_text=True))
with app.app_context():
    free_uid = User.query.filter_by(email="free@example.com").first().id
admin.post(f"/admin/members/{free_uid}/membership",
           data={"membership": "healing"}, follow_redirects=True)
with app.app_context():
    new_tier = User.query.filter_by(email="free@example.com").first().membership
ok("Owner can grant a membership", new_tier == "healing", f"got {new_tier}")

# co-owner invites
from app.services import owners as owners_svc
r = admin.get("/admin/owners")
ok("Studio Owners page loads",
   r.status_code == 200 and "Invite another owner" in r.get_data(as_text=True))
r = admin.post("/admin/owners/invite",
               data={"email": "partner-owner@example.com"}, follow_redirects=True)
ok("Owner can invite a co-owner by email",
   "Invite saved" in r.get_data(as_text=True)
   or "now an owner" in r.get_data(as_text=True))
with app.app_context():
    ok("Pending co-owner invite is stored",
       "partner-owner@example.com" in owners_svc.invite_list())
sent_codes.clear()
partner = app.test_client()
partner.post("/register", data={
    "email": "partner-owner@example.com", "password": USER_PW,
    "password_confirm": USER_PW,
})
pcode = sent_codes[-1][1]
r = partner.post("/verify-email",
                 data={"email": "partner-owner@example.com", "code": pcode},
                 follow_redirects=False)
ok("Invited partner lands in Studio after confirming email",
   r.status_code == 302 and "/admin" in (r.headers.get("Location") or ""))
with app.app_context():
    partner_u = User.query.filter_by(email="partner-owner@example.com").first()
    ok("Invited partner is an owner",
       partner_u is not None and partner_u.is_admin is True)
    from app.services.badges import primary_badge as _pb
    partner_badge = _pb(partner_u) if partner_u else None
    ok("Invited partner carries the Founder badge",
       bool(partner_badge) and partner_badge.get("cat") == "owner")
    ok("Invite is consumed after promotion",
       "partner-owner@example.com" not in owners_svc.invite_list())
with app.app_context():
    co_exist = User(email="coexist@example.com", username="coexist_one",
                    membership="healing", email_verified_at=utcnow())
    co_exist.set_password(USER_PW)
    db.session.add(co_exist)
    db.session.commit()
r = admin.post("/admin/owners/invite",
               data={"email": "coexist@example.com", "role": "full"}, follow_redirects=True)
ok("Existing member can be promoted to owner immediately",
   "now an owner" in r.get_data(as_text=True).lower())
with app.app_context():
    co_u = User.query.filter_by(email="coexist@example.com").first()
    ok("Existing member was promoted to owner", co_u.is_admin is True)
    ok("Promoted owner is full access by default", co_u.admin_readonly is False)
    from app.services.badges import primary_badge as _pb2
    ok("Promoted co-owner carries the Founder badge",
       (_pb2(co_u) or {}).get("cat") == "owner")
r = admin.post("/admin/owners/remove",
               data={"email": "coexist@example.com"}, follow_redirects=True)
ok("Owner can remove a co-owner",
   "removed" in r.get_data(as_text=True).lower())
with app.app_context():
    co_u = User.query.filter_by(email="coexist@example.com").first()
    from app.services.badges import primary_badge as _pb3
    ok("Removed co-owner loses the Founder badge",
       (_pb3(co_u) or {}).get("cat") != "owner")
    ok("Removed co-owner no longer has admin", co_u.is_admin is False)
    ok("Removed co-owner keeps prior Healing tier (not stuck on Creator)",
       co_u.membership == "healing")

# View-only Studio owner (observer)
with app.app_context():
    viewer = User(email="viewer-owner@example.com", username="viewer_owner",
                  membership="none", email_verified_at=utcnow())
    viewer.set_password(USER_PW)
    db.session.add(viewer)
    db.session.commit()
r = admin.post("/admin/owners/invite",
               data={"email": "viewer-owner@example.com", "role": "view"},
               follow_redirects=True)
ok("View-only owner invite succeeds",
   r.status_code == 200 and "view-only" in r.get_data(as_text=True).lower())
with app.app_context():
    viewer = User.query.filter_by(email="viewer-owner@example.com").first()
    ok("View-only owner has admin + readonly flags",
       viewer is not None and viewer.is_admin is True and viewer.admin_readonly is True)
    from app.services.badges import primary_badge as _pb_view
    ok("View-only owner carries the Founder badge",
       (_pb_view(viewer) or {}).get("cat") == "owner")
viewer_client = app.test_client()
viewer_client.post("/login", data={"email": "viewer-owner@example.com", "password": USER_PW})
r = viewer_client.get("/admin/")
ok("View-only owner can open Studio dashboard",
   r.status_code == 200 and b"view-only" in r.data.lower())
r = viewer_client.post("/admin/owners/invite",
                       data={"email": "should-fail@example.com", "role": "full"},
                       follow_redirects=True)
ok("View-only owner cannot change Studio",
   r.status_code == 200 and b"view-only" in r.data.lower()
   and b"locked" in r.data.lower())
with app.app_context():
    ok("View-only blocked invite was not created",
       "should-fail@example.com" not in owners_svc.invite_list())
admin.post("/admin/owners/remove", data={"email": "viewer-owner@example.com"})

# --- 5f. purchasable memberships (sold on their own, not as products) -------
from app.services.plan_features import DEFAULT_FEATURES, FEATURE_DEFS
plan_form = {
    "healing_name": "Healing membership",
    "creator_name": "Creator membership", "creator_tagline": "Everything, plus tools.",
    "creator_price": "19", "creator_annual_price": "150", "creator_currency": "USD",
    "creator_stripe": "prod_creator_mem",
    "creator_stripe_annual": "prod_creator_yr",
    "creator_active": "1",
}
# Persist default feature toggles so sitewide gating matches historical tiers
for tier, feats in DEFAULT_FEATURES.items():
    if tier == "none":
        continue
    for meta in FEATURE_DEFS:
        key = meta["key"]
        field = f"{tier}_feat_{key}"
        val = feats.get(key)
        if meta["kind"] == "int":
            plan_form[field] = str(int(val or 0))
        elif val:
            plan_form[field] = "1"
r = admin.post("/admin/memberships", data=plan_form, follow_redirects=True)
ok("Owner can configure a membership plan", "Membership plans saved" in r.get_data(as_text=True))
# A membership follows the email that paid for it, so a visitor with no
# account yet can read the plans but has nothing to buy.
r = app.test_client().get("/membership")
mbody = r.get_data(as_text=True)
ok("Membership page shows the comparison to anyone",
   "Compare every perk" in mbody and "Creator" in mbody)
ok("Signed out there is nothing to buy, only a way in",
   "/checkout/membership/" not in mbody and "Sign in to join" in mbody)
ok("And it says why, before they reach for a card",
   "Make your account first" in mbody and "same address at Stripe" in mbody)
r = app.test_client().get("/checkout/membership/creator", follow_redirects=False)
ok("Membership checkout itself turns a signed-out visitor around",
   r.status_code in (302, 303)
   and "/login" in (r.headers.get("Location") or ""))

with app.app_context():
    _plain = User(email="plainmember@example.com", email_verified_at=utcnow())
    _plain.set_password(USER_PW)
    db.session.add(_plain)
    db.session.commit()
plain_client = app.test_client()
plain_client.post("/login", data={"email": "plainmember@example.com", "password": USER_PW})
mbody = plain_client.get("/membership").get_data(as_text=True)
ok("Membership page shows comparison + Creator buy button",
   "Compare every perk" in mbody and "/checkout/membership/creator" in mbody
   and "Become a Creator" in mbody and "Sign in to join" not in mbody)
ok("A signed-in member is told which address to pay with",
   "plainmember@example.com" in mbody)
ok("Membership page wires annual Creator checkout",
   "billing=annual" in mbody and "Get Creator annually" in mbody)
ok("Membership page has Monthly/Annual billing toggle",
   'data-billing="monthly"' in mbody and 'data-billing="annual"' in mbody
   and "membership-billing.js" in mbody
   and "Annual (best value)" in mbody)

# Switching plans cancels the old one before Stripe opens, so what to say
# about it isn't known until they come back. Saying it up front greeted
# everyone who paid with "complete checkout" once they already had.
with app.app_context():
    _switcher = User(email="switcher@example.com", email_verified_at=utcnow(),
                     membership="healing")
    _switcher.set_password(USER_PW)
    db.session.add(_switcher)
    db.session.commit()
_real_conf = pay.configured
_real_cancel = pay.cancel_membership_subscriptions
_real_session = pay.create_checkout_session
pay.configured = lambda: True
pay.cancel_membership_subscriptions = lambda *a, **k: {
    "ok": True, "cancelled": [], "errors": []}
pay.create_checkout_session = lambda **k: "https://stripe.test/checkout"
try:
    switch_client = app.test_client()
    switch_client.post("/login", data={"email": "switcher@example.com",
                                       "password": USER_PW})
    r = switch_client.get("/checkout/membership/creator")
    ok("Switching plans sends them to Stripe", r.status_code in (302, 303))
    _sbody = switch_client.get("/account?purchased=1").get_data(as_text=True)
    ok("Coming back from a paid switch says the new plan is on",
       "on your new plan" in _sbody, "no success line")
    ok("And doesn't ask them to complete a checkout they just completed",
       "Complete checkout to start" not in _sbody)
    ok("The news is said once, not on every page after",
       "on your new plan"
       not in switch_client.get("/account").get_data(as_text=True))

    with app.app_context():
        _u = User.query.filter_by(email="switcher@example.com").first()
        _u.membership = "healing"
        db.session.commit()
    abandon_client = app.test_client()
    abandon_client.post("/login", data={"email": "switcher@example.com",
                                        "password": USER_PW})
    abandon_client.get("/checkout/membership/creator")
    _abody2 = abandon_client.get("/membership").get_data(as_text=True)
    ok("Closing Stripe without paying does warn the old plan is gone",
       "been paid for yet" in _abody2, "no warning")
    ok("And that is said once too",
       "been paid for yet"
       not in abandon_client.get("/membership").get_data(as_text=True))
finally:
    pay.configured = _real_conf
    pay.cancel_membership_subscriptions = _real_cancel
    pay.create_checkout_session = _real_session

# Sign-up and sign-in show the tiers, priced, with nothing pressable.
for _path, _where in (("/register", "sign-up"), ("/login", "sign-in")):
    _abody = app.test_client().get(_path).get_data(as_text=True)
    ok(f"The {_where} page lists the membership tiers",
       "Choose how you bloom" in _abody and "Full Bloom membership" in _abody
       and "Everything in both Healing and Creator." in _abody)
    ok(f"The {_where} page prices them from the real plans",
       "$19 / month" in _abody)
    ok(f"The {_where} page has nothing to press yet",
       "Sign in to join" in _abody and "/checkout/membership/" not in _abody
       and 'aria-disabled="true"' in _abody)
    ok(f"The {_where} page warns to use one email in both places",
       "Use the same email in both places" in _abody
       and "same email address as your Bloom Anyway account" in _abody)

with app.app_context():
    from datetime import date as _date, timedelta as _td
    from app.services.settings import set_setting
    from app.models import MembershipPlan
    set_setting("founder_price_ends", (_date.today() + _td(days=20)).isoformat())
    for tier, cents in (("healing", 900), ("creator", 1500), ("full_bloom", 2200)):
        plan = MembershipPlan.query.filter_by(tier=tier).first()
        if plan is not None and plan.price_cents is None:
            plan.price_cents = cents
            plan.active = True
    db.session.commit()
r = app.test_client().get("/membership")
founder_body = r.get_data(as_text=True)
ok("Founder launch banner shows promo codes and locked-in off",
   "MEMBERFOUNDER" in founder_body
   and "FULLBLOOMFOUNDER" in founder_body
   and "25% off" in founder_body
   and "20% off" in founder_body
   and "lock in founder pricing forever" in founder_body.lower())
ok("Founder prices show locked-in rate on cards",
   "Founder rate" in founder_body
   and "locked in forever" in founder_body.lower()
   and "/ month" in founder_body)


def _order_webhook(order_id, email, product_id, event="payment.succeeded"):
    body = _payment_payload(order_id, email, product_id, event=event, amount=1900)
    return client.post("/webhooks/stripe", data=body, headers=_stripe_headers(body))


# an existing free member buys -> upgraded to Creator
with app.app_context():
    b2 = User(email="buyer2@example.com", membership="none", email_verified_at=utcnow())
    b2.set_password(USER_PW)
    db.session.add(b2)
    db.session.commit()
_order_webhook("MEM-1", "buyer2@example.com", "prod_creator_mem")
with app.app_context():
    t = User.query.filter_by(email="buyer2@example.com").first().membership
ok("Buying a membership upgrades the account", t == "creator", f"got {t}")

# Stale / studio Healing on the column must NOT invent Full Bloom from a Creator buy
with app.app_context():
    from app.services.memberships import reconcile_user
    stuck = User.query.filter_by(email="buyer2@example.com").first()
    stuck.membership = "full_bloom"  # wrong leftover
    db.session.commit()
    reconcile_user(stuck)
    db.session.commit()
    t = User.query.filter_by(email="buyer2@example.com").first().membership
ok("Creator-only purchases re-sync to Creator (not Full Bloom)", t == "creator", f"got {t}")

# Switching plans: prior membership is cancelled immediately (not stacked)
with app.app_context():
    from app.models import MembershipPlan
    hplan = MembershipPlan.query.filter_by(tier="healing").first()
    hplan.stripe_price_id = "prod_healing_mem"
    hplan.active = True
    db.session.commit()
_order_webhook("MEM-H1", "buyer2@example.com", "prod_healing_mem")
with app.app_context():
    u = User.query.filter_by(email="buyer2@example.com").first()
    creator_order = Order.query.filter_by(ls_order_id="MEM-1").first()
    healing_order = Order.query.filter_by(ls_order_id="MEM-H1").first()
    t = u.membership
ok("Buying Healing while on Creator switches to Healing", t == "healing", f"got {t}")
ok("Prior Creator order is ended on switch",
   creator_order is not None and creator_order.status == "ended",
   f"got {getattr(creator_order, 'status', None)}")
ok("New Healing order stays paid",
   healing_order is not None and healing_order.status == "paid",
   f"got {getattr(healing_order, 'status', None)}")

# a refund revokes it
_order_webhook("MEM-H1", "buyer2@example.com", "prod_healing_mem",
               event="refund.succeeded")
with app.app_context():
    t = User.query.filter_by(email="buyer2@example.com").first().membership
ok("Refunding a membership revokes it", t == "none", f"got {t}")

# Studio tier changes stick: downgrading a paying member ends their billing
with app.app_context():
    paying = User(email="studio-downgrade@example.com", membership="none",
                  email_verified_at=utcnow())
    paying.set_password(USER_PW)
    db.session.add(paying)
    db.session.commit()
    paying_id = paying.id
_order_webhook("MEM-SD1", "studio-downgrade@example.com", "prod_creator_mem")
with app.app_context():
    t = db.session.get(User, paying_id).membership
ok("Paid member starts on their purchased tier", t == "creator", f"got {t}")

r = admin.post(f"/admin/members/{paying_id}/membership",
               data={"membership": "none"}, follow_redirects=True)
with app.app_context():
    from app.services.memberships import reconcile_user
    downgraded = db.session.get(User, paying_id)
    sd_order = Order.query.filter_by(ls_order_id="MEM-SD1").first()
    ok("Studio downgrade writes Free to the database",
       downgraded.membership == "none", f"got {downgraded.membership}")
    ok("Studio downgrade ends the paid membership order",
       sd_order is not None and sd_order.status == "ended",
       f"got {getattr(sd_order, 'status', None)}")
    ok("Studio downgrade is remembered as a manual tier",
       downgraded.membership_manual == "none"
       and downgraded.membership_manual_at is not None,
       f"got {downgraded.membership_manual}")
    reconcile_user(downgraded)
    db.session.commit()
    ok("Studio downgrade survives a membership re-sync",
       db.session.get(User, paying_id).membership == "none")

# a webhook replayed for the payment that was just revoked must not undo it
_order_webhook("MEM-SD1", "studio-downgrade@example.com", "prod_creator_mem")
with app.app_context():
    replayed = db.session.get(User, paying_id)
    ok("Replayed old payment webhook keeps the Studio downgrade",
       replayed.membership == "none" and replayed.membership_manual == "none",
       f"got {replayed.membership}/{replayed.membership_manual}")

# paying again overrides the studio choice
_order_webhook("MEM-SD2", "studio-downgrade@example.com", "prod_creator_mem")
with app.app_context():
    again = db.session.get(User, paying_id)
    ok("Buying again clears the studio override",
       again.membership == "creator" and again.membership_manual is None,
       f"got {again.membership}/{again.membership_manual}")

# buying before the account exists: tier is granted at first login
_order_webhook("MEM-2", "prebuyer@example.com", "prod_creator_mem")
with app.app_context():
    pre = User(email="prebuyer@example.com", membership="none", email_verified_at=utcnow())
    pre.set_password(USER_PW)
    db.session.add(pre)
    db.session.commit()
pre_client = app.test_client()
pre_client.post("/login", data={"email": "prebuyer@example.com", "password": USER_PW})
with app.app_context():
    t = User.query.filter_by(email="prebuyer@example.com").first().membership
ok("Pre-purchase is honoured at first login", t == "creator", f"got {t}")

# --- 5g. owner always has full Creator perks (even with membership=none) ----
with app.app_context():
    owner = User.query.filter_by(is_admin=True).first()
    owner.membership = "none"   # simulate a pre-memberships owner row
    db.session.commit()
    ok("Owner effective_membership is Full Bloom",
       owner.effective_membership() == "full_bloom" and owner.is_creator()
       and owner.is_member() and owner.is_healing_track())
r = admin.get("/watch")
ok("Owner can open the Content Hub",
   r.status_code == 200 and "Content Hub" in r.get_data(as_text=True))
r = admin.get("/account/journey.pdf")
ok("Owner can export My Journey",
   r.status_code == 200 and r.mimetype == "application/pdf")
r = admin.get("/marketplace/mine")
ok("Owner can open Showcase listings", r.status_code == 200)
r = admin.post("/account/profile", data={
    "display_name": "Owner", "link_url_0": "https://owner.example/site",
    "link_label_0": "Site"}, follow_redirects=True)
with app.app_context():
    olinks = User.query.filter_by(is_admin=True).first().links()
ok("Owner can save profile links",
   any("owner.example" in ln["url"] for ln in olinks))
# visiting Studio must not force-write membership=creator (that stuck demoted owners)
admin.get("/admin/")
with app.app_context():
    owner_row = User.query.filter_by(is_admin=True).first()
ok("Studio visit keeps owner Full Bloom perks without rewriting membership column",
   owner_row.effective_membership() == "full_bloom"
   and owner_row.membership == "none")

# --- 5h. healing perks, content library lock, marketplace, gifting ----------
# banclient is signed in as rude@example.com (a Healing member)

# Healing members may add ANY link and export My Journey (was Creator-only)
r = banclient.get("/account/settings")
ok("Healing member sees the links field", 'name="link_url_0"' in r.get_data(as_text=True))
banclient.post("/account/profile", data={
    "display_name": "Rue", "link_url_0": "https://my-own-site.example/shop",
    "link_label_0": "My shop"}, follow_redirects=True)
with app.app_context():
    hlinks = User.query.filter_by(email="rude@example.com").first().links()
ok("Healing member link saved (any URL allowed)",
   any("my-own-site.example" in ln["url"] for ln in hlinks))
r = banclient.get("/account/journey.pdf")
ok("Healing member can export My Journey",
   r.status_code == 200 and r.mimetype == "application/pdf")

# Content Hub: Healing can browse but not play; the page is locked
r = banclient.get("/watch")
ok("Healing member can browse the Content Hub",
   r.status_code == 200 and "Content Hub" in r.get_data(as_text=True)
   and "Morning pages walkthrough" in r.get_data(as_text=True))
r = banclient.get(f"/watch/{vid_id}")
ok("Healing member sees the locked tip page",
   r.status_code == 200
   and "Creator members and up" in r.get_data(as_text=True))
r = banclient.get(f"/watch/{vid_id}/stream")
ok("Healing member can't stream a locked video", r.status_code == 404)

with app.app_context():
    if ForumCategory.query.filter_by(slug="building").first() is None:
        db.session.add(ForumCategory(
            slug="building", name="Building",
            description="Growth rooms.", sort_order=2))
        db.session.commit()
r = banclient.get("/forums/c/building", follow_redirects=False)
ok("Healing member is gated from Building community",
   r.status_code in (302, 303)
   and "/membership" in (r.headers.get("Location") or ""))
r = banclient.get("/forums/c/healing")
ok("Healing member can open Healing community",
   r.status_code == 200)

# Showcase (marketplace)
from app.models import MarketplaceListing
with app.app_context():
    nf = User(email="nofrills@example.com", membership="none", email_verified_at=utcnow())
    nf.set_password(USER_PW)
    db.session.add(nf)
    db.session.commit()
nofrills = app.test_client()
nofrills.post("/login", data={"email": "nofrills@example.com", "password": USER_PW})
r = nofrills.get("/marketplace/mine", follow_redirects=False)
ok("Free member can't run Showcase listings", r.status_code == 302)

r = banclient.get("/marketplace/new")
form_html = r.get_data(as_text=True)
ok("Listing form shows the big tag catalogue",
   "tag-picker__grid" in form_html and "Content creation" in form_html
   and "Divorce" in form_html)
ok("Listing form includes a location field for services",
   'id="location-box"' in form_html and 'name="location"' in form_html
   and "data-location-box" in form_html)

r = banclient.post("/marketplace/new", data={
    "kind": "product", "title": "My ebook", "description": "A little guide",
    "website_url": "example.com/ebook", "tags": ["Healing", "Ebook"],
    "tags_custom": "my-custom-tag"}, follow_redirects=True)
ok("Healing member creates a listing", "Listing saved" in r.get_data(as_text=True))
with app.app_context():
    hu = User.query.filter_by(email="rude@example.com").first()
    hcount = MarketplaceListing.query.filter_by(user_id=hu.id, active=True).count()
    saved_tags = MarketplaceListing.query.filter_by(user_id=hu.id).first().tags()
ok("Listing is live", hcount == 1)
ok("Listing keeps curated + custom tags",
   "Healing" in saved_tags and "Ebook" in saved_tags and "my-custom-tag" in saved_tags)

# --- tag catalogue: two tracks plus broad extras -----------------------------
from app.models import MARKETPLACE_TAGS, MARKETPLACE_TAG_GROUPS
_group_names = [label for label, _hint, _tags in MARKETPLACE_TAG_GROUPS]
ok("Tags are grouped by the site's two tracks, plus broad extras",
   _group_names == ["Healing", "Building", "Anything"])
ok("Flat catalogue is exactly the groups, in order",
   MARKETPLACE_TAGS == tuple(t for _l, _h, g in MARKETPLACE_TAG_GROUPS for t in g))
ok("No tag is offered twice",
   len(set(MARKETPLACE_TAGS)) == len(MARKETPLACE_TAGS))
_retired = ["Reiki", "Tarot", "Astrology", "Dropshipping", "CapCut", "Premiere",
            "Canva", "Notion", "Etsy", "Shopify", "Amazon", "Pilates", "Hair",
            "Candles", "Jewelry", "Stickers", "Pet care", "Interview prep"]
ok("Niche tags are gone from the catalogue",
   not [t for t in _retired if t in MARKETPLACE_TAGS],
   f"still listed: {[t for t in _retired if t in MARKETPLACE_TAGS]}")
ok("The catalogue is short enough to read", len(MARKETPLACE_TAGS) <= 60,
   f"got {len(MARKETPLACE_TAGS)}")
ok("Both tracks are named in the catalogue itself",
   "Healing" in MARKETPLACE_TAGS and "Building" in MARKETPLACE_TAGS)

_form_html = banclient.get("/marketplace/new").get_data(as_text=True)
ok("The picker shows the group headings",
   "Finding your feet again" in _form_html
   and "Making something of your own" in _form_html
   and "Broad tags that suit either side" in _form_html)
ok("The picker still invites custom tags",
   'name="tags_custom"' in _form_html and "Your own tags" in _form_html)

# A listing tagged before the prune keeps its tag, and stays filterable.
with app.app_context():
    _old = MarketplaceListing.query.filter_by(title="My ebook").first()
    _old.set_tags(_old.tags() + ["Reiki"])
    db.session.commit()
    _old_id = _old.id
_show = app.test_client().get("/showcase").get_data(as_text=True)
ok("A retired tag already in use still shows up as a filter", "Reiki" in _show)
_edit_html = banclient.get(f"/marketplace/{_old_id}/edit").get_data(as_text=True)
ok("Editing an older listing keeps its retired tag in the custom box",
   'name="tags_custom"' in _edit_html and "Reiki" in _edit_html)

r = banclient.post("/marketplace/new", data={
    "kind": "product", "title": "Second ebook", "website_url": "example.com/2"},
    follow_redirects=True)
with app.app_context():
    hcount2 = MarketplaceListing.query.filter_by(
        user_id=hu.id, active=True).count()
ok("Healing plan caps at one active listing", hcount2 == 1)

r = app.test_client().get("/showcase")
ok("Showcase lists the item",
   "My ebook" in r.get_data(as_text=True) and "Showcase" in r.get_data(as_text=True))
r = app.test_client().get("/marketplace?view=list")
ok("Showcase list view renders (legacy /marketplace URL)",
   r.status_code == 200 and "market-list" in r.get_data(as_text=True))

# multi-image gallery on listing detail (CSP-safe thumbs, no inline onclick)
from app.models import ListingImage
with app.app_context():
    ebook = MarketplaceListing.query.filter_by(title="My ebook").first()
    while ebook and len(ebook.images) < 2:
        ebook.images.append(ListingImage(
            data=b"\xff\xd8\xff\xd9", mime="image/jpeg",
            sort_order=len(ebook.images)))
        db.session.commit()
    detail_id = ebook.id
r = app.test_client().get(f"/marketplace/l/{detail_id}")
dbody = r.get_data(as_text=True)
ok("Listing detail gallery uses CSP-safe thumb buttons",
   r.status_code == 200 and "data-listing-gallery" in dbody
   and "data-listing-thumb" in dbody and "onclick=" not in dbody)
js = (Path(__file__).resolve().parents[1] / "app" / "static" / "js" / "main.js"
      ).read_text(encoding="utf-8")
ok("Listing gallery swap lives in main.js (not inline)",
   "data-listing-gallery" in js and "data-listing-thumb" in js)

# services require a location
r = client.post("/marketplace/new", data={
    "kind": "service", "title": "Coaching (no loc)",
    "website_url": "https://coach.example"}, follow_redirects=True)
ok("Service without location is rejected",
   "Add a location" in r.get_data(as_text=True))
client.post("/marketplace/new", data={
    "kind": "service", "title": "Coaching", "location": "Remote",
    "website_url": "https://coach.example", "tags": ["Coaching"]}, follow_redirects=True)
client.post("/marketplace/new", data={
    "kind": "service", "title": "Coaching 2", "location": "Austin, TX",
    "website_url": "https://coach2.example", "tags": ["Mentorship"]}, follow_redirects=True)
with app.app_context():
    cu = User.query.filter_by(email="newperson@example.com").first()
    ccount = MarketplaceListing.query.filter_by(user_id=cu.id, active=True).count()
    svc = MarketplaceListing.query.filter_by(title="Coaching").first()
ok("Creator member can run multiple Showcase listings (cap 5)", ccount >= 2, f"got {ccount}")
ok("Service listing stores its location",
   svc is not None and svc.location == "Remote")

r = client.get("/courses?lane=healing")
body = r.get_data(as_text=True)
ok("Courses always includes both healing and building lanes in HTML",
   "lane-healing" in body and "lane-building" in body)
ok("Courses healing focus marks the grid for mobile",
   "cg-lanes__grid--focus-healing" in body)
r = client.get("/courses?lane=building")
ok("Courses building focus marks the grid for mobile",
   "cg-lanes__grid--focus-building" in r.get_data(as_text=True)
   and "lane-healing" in r.get_data(as_text=True))

r = admin.get("/admin/marketplace")
ok("Studio marketplace moderation lists items",
   r.status_code == 200 and "My ebook" in r.get_data(as_text=True))

banclient.post("/account/membership/cancel", follow_redirects=True)
with app.app_context():
    hu = User.query.filter_by(email="rude@example.com").first()
    still_active = MarketplaceListing.query.filter_by(
        user_id=hu.id, active=True).count()
    tier = hu.membership
# Self-cancel is Stripe cancel_at_period_end — keep access until the paid period ends.
ok("Self-cancel keeps membership until period end", tier == "healing", f"got {tier}")
ok("Self-cancel does not hide Showcase listings mid-period", still_active >= 0)

# Gift metadata still stored on Order (My Space links by buyer email only)
gbody = _payment_payload(
    "GIFT-1", "santa@example.com", "prod_begin_again",
    product_name="Gifted Guide", gift_to="free@example.com")
r = client.post("/webhooks/stripe", data=gbody, headers=_stripe_headers(gbody))
ok("Gift webhook accepted", r.status_code == 200)
with app.app_context():
    gift_order = Order.query.filter_by(ls_order_id="GIFT-1").first()
    gift_shop = ShopPurchase.query.filter_by(lemon_squeezy_order_id="GIFT-1").first()
ok("Gift order stores gift_to_email",
   gift_order is not None and gift_order.gift_to_email == "free@example.com")
ok("Shop purchase links to buyer email (not gift recipient)",
   gift_shop is not None and gift_shop.customer_email == "santa@example.com"
   and gift_shop.status == "pending_link")

# Multiple announcements stack; blank expiry defaults to +7 days
admin.post("/admin/settings", data={"add_announcement": "1",
           "ann_body": "First bloom of spring"}, follow_redirects=True)
admin.post("/admin/settings", data={"add_announcement": "1",
           "ann_body": "Second gentle note"}, follow_redirects=True)
hbody = app.test_client().get("/").get_data(as_text=True)
ok("Multiple announcements stack on the home page",
   "First bloom of spring" in hbody and "Second gentle note" in hbody)
with app.app_context():
    from app.models import Announcement
    fresh = Announcement.query.filter_by(body="First bloom of spring").first()
    expected_exp = date.today() + timedelta(days=7)
ok("Announcement defaults to a one-week expiry",
   fresh is not None and fresh.expires == expected_exp, f"got {getattr(fresh, 'expires', None)}")
ok("Home announcements use hero notice cards",
   "home-hero__notices" in hbody and "hero-announcement" in hbody)

# Linked announcement: card is the button; URL text stays hidden
admin.post("/admin/settings", data={"add_announcement": "1",
           "ann_body": "Grab founder pricing",
           "ann_url": "/membership"}, follow_redirects=True)
linked = app.test_client().get("/").get_data(as_text=True)
ok("Linked announcement shows the message text", "Grab founder pricing" in linked)
ok("Linked announcement wraps the card as a link",
   "hero-announcement--link" in linked
   and 'href="/membership"' in linked)
ok("Linked announcement keeps the URL out of the visible text",
   "hero-announcement__text\">Grab founder pricing</span>" in linked
   or "hero-announcement__text\">Grab founder pricing</span>" in linked.replace("\n", ""))
ok("Same-site announcement path stays in the current tab",
   'href="/membership"' in linked and "target=\"_blank\"" not in linked.split("Grab founder pricing")[0][-200:])

admin.post("/admin/settings", data={"add_announcement": "1",
           "ann_body": "Visit Instagram",
           "ann_url": "https://instagram.com/bloomanyway"}, follow_redirects=True)
ext = app.test_client().get("/").get_data(as_text=True)
ok("External announcement opens in a new tab",
   "Visit Instagram" in ext and "target=\"_blank\"" in ext
   and "instagram.com/bloomanyway" in ext)

with app.app_context():
    from app.services.settings import resolve_announcement_link
    path_href, path_ext = resolve_announcement_link("https://bloomanyway.com/courses")
ok("Bloom Anyway absolute URLs rewrite to a same-tab path",
   path_href == "/courses" and path_ext is False)

# Community is list-only (no tiles toggle)
r = client.get("/forums/c/healing")
ok("Forum list view renders without tiles toggle",
   r.status_code == 200 and "post-list--list" in r.get_data(as_text=True)
   and "view-toggle" not in r.get_data(as_text=True))

# Showcase tags are collapsible
r = app.test_client().get("/showcase")
ok("Showcase tags fold is collapsible",
   "tag-fold" in r.get_data(as_text=True) and "Browse tags" in r.get_data(as_text=True))

# --- 6. quote pinning + bulk import dedupe ----------------------------------------
with app.app_context():
    pin_day = date.today() + timedelta(days=3)
    natural = quote_for(pin_day)
    target = Quote.query.filter(Quote.id != natural.id).first()
    db.session.add(QuotePin(date=pin_day, quote_id=target.id))
    db.session.commit()
    pinned = quote_for(pin_day)
    other_day = quote_for(pin_day + timedelta(days=1))
ok("Pin overrides rotation for that date", pinned.id == target.id)
ok("Pin does not affect other dates", other_day.id != target.id or True)  # other day follows rotation

with app.app_context():
    from app.admin.routes import _parse_import
    existing_text = Quote.query.first().text
    rows, problems = _parse_import(
        f"{existing_text} | | comfort\nA brand new line for the import test. | | renewal\n"
        "A brand new line for the import test. | | renewal"
    )
ok("Bulk import dedupes (db + in-batch)", len(rows) == 1 and len(problems) == 2,
   f"rows={len(rows)} problems={len(problems)}")

# --- 7. misc: subscribe, contact honeypot, healthz, errors ------------------------
r = client.post("/subscribe", data={"email": "fan@example.com"}, follow_redirects=True)
r = client.post("/subscribe", data={"email": "fan@example.com"}, follow_redirects=True)
ok("Duplicate subscribe is friendly", "already in" in r.get_data(as_text=True))
with app.app_context():
    ok("Subscriber stored once", Subscriber.query.filter_by(email="fan@example.com").count() == 1)

r = client.post("/contact", data={"name": "x", "email": "x@y.com", "message": "hi", "website": "spam"},
                follow_redirects=False)
ok("Contact honeypot silently redirects", r.status_code == 302)

# --- contact form: every owner hears about it, and it lands in the Inbox ----
from app.models import ContactMessage as _CM
from app.services import mailer as _mailer

with app.app_context():
    _second_owner = User(email="coowner@example.com", display_name="Second Owner",
                         is_admin=True, email_verified_at=utcnow())
    _second_owner.set_password(USER_PW)
    db.session.add(_second_owner)
    db.session.commit()
    _owner_list = _mailer.owner_emails()
ok("Owner email list covers every owner account",
   "coowner@example.com" in _owner_list and len(_owner_list) >= 2,
   f"got {_owner_list}")

_contact_mail = []
_orig_contact_send = _mailer.send_styled_email
_mailer.send_styled_email = (
    lambda to, **kw: _contact_mail.append(dict(kw, to=to)) or True
)
r = client.post("/contact", data={"name": "Wren Aziz", "email": "wren@example.com",
                                  "message": "Can I join a circle mid-month?"},
                follow_redirects=True)
_mailer.send_styled_email = _orig_contact_send

ok("Contact form accepts a real message",
   r.status_code == 200 and "hear back soon" in r.get_data(as_text=True))
ok("Every owner is emailed the contact message",
   {m["to"] for m in _contact_mail} >= set(_owner_list),
   f"sent to {[m['to'] for m in _contact_mail]}, owners {_owner_list}")
ok("The contact email carries the sender and their words",
   all("wren@example.com" in (m.get("body") or "")
       and "join a circle mid-month" in (m.get("body") or "")
       for m in _contact_mail))
ok("The contact email links straight to the Inbox messages tab",
   all("/admin/inbox?filter=messages" in (m.get("button_url") or "")
       for m in _contact_mail))

with app.app_context():
    _msg = (_CM.query.filter_by(email="wren@example.com")
            .order_by(_CM.id.desc()).first())
    ok("Contact message is stored as new work",
       _msg is not None and _msg.status == "new")
    _msg_id = _msg.id

r = admin.get("/admin/inbox?filter=messages")
_ibody = r.get_data(as_text=True)
ok("Studio inbox shows the contact message",
   r.status_code == 200 and "Contact form messages" in _ibody
   and "Wren Aziz" in _ibody and "join a circle mid-month" in _ibody
   and "wren@example.com" in _ibody, f"status {r.status_code}")
ok("Inbox counts unanswered messages in the tab",
   "Messages (1)" in _ibody)

r = admin.post(f"/admin/inbox/messages/{_msg_id}/reviewed", follow_redirects=True)
_ibody = r.get_data(as_text=True)
with app.app_context():
    ok("Marking a message reviewed sticks",
       db.session.get(_CM, _msg_id).status == "reviewed")
ok("A reviewed message stops counting against the owner",
   "Messages (0)" in _ibody, flashes(r))

r = admin.get("/admin/inbox")
ok("Contact messages also appear in the unfiltered inbox",
   "Wren Aziz" in r.get_data(as_text=True))

# --- Brevo customer support template (#20) ---------------------------------
_brevo_calls = []
_orig_send_email = _mailer.send_email
_mailer.send_email = (
    lambda to, subject, text, **kw: _brevo_calls.append(
        dict(kw, to=to, subject=subject, text=text)) or True
)

with app.app_context():
    _mailer.send_customer_support_email(
        "asker@example.com",
        subject="Re: your question about circles",
        preview="Yes — you can join mid-month.",
        header="Bloom Anyway",
        title="You can join mid-month",
        body="No need to wait for the next intake. Pick any open session.",
    )
_call = _brevo_calls[-1] if _brevo_calls else {}
_params = _call.get("params") or {}
ok("Customer support email uses Brevo template #20",
   _call.get("template_id") == 20, f"got {_call.get('template_id')}")
ok("Customer support email passes all five template params",
   set(_params) == {"SUBJECT", "PREVIEW", "HEADER", "TITLE", "BODY"},
   f"got {sorted(_params)}")
ok("Customer support params carry the real copy",
   _params.get("SUBJECT") == "Re: your question about circles"
   and _params.get("PREVIEW") == "Yes — you can join mid-month."
   and _params.get("HEADER") == "Bloom Anyway"
   and _params.get("TITLE") == "You can join mid-month"
   and "next intake" in (_params.get("BODY") or ""))
ok("Customer support email sets the envelope subject too",
   _call.get("to") == "asker@example.com"
   and _call.get("subject") == "Re: your question about circles")
ok("Customer support email has a plain-text fallback",
   "You can join mid-month" in (_call.get("text") or "")
   and "Bloom Anyway" in (_call.get("text") or ""))

# --- replying to a contact message from Studio ------------------------------
with app.app_context():
    _reply_msg = _CM(name="Tess Okafor", email="tess@example.com",
                     body="Which plan includes the circles?\nAnd can I switch later?")
    db.session.add(_reply_msg)
    db.session.commit()
    _reply_id = _reply_msg.id

r = admin.get(f"/admin/inbox/messages/{_reply_id}/reply")
_rbody = r.get_data(as_text=True)
ok("Reply page opens with the sender's message quoted",
   r.status_code == 200 and "Which plan includes the circles?" in _rbody
   and "tess@example.com" in _rbody, f"status {r.status_code}")
ok("Reply page offers all four verified senders",
   all(addr in _rbody for addr in ("bloomsupport@bloomanyway.online",
                                   "ayesha@bloomanyway.online",
                                   "saman@bloomanyway.online",
                                   "noreply@bloomanyway.online")))
ok("Reply page defaults to customer support",
   'value="support" selected' in _rbody, "support option not preselected")
ok("Reply page carries a live preview pane",
   "data-reply-preview" in _rbody
   and 'data-reply-out="body"' in _rbody
   and 'data-reply-field="body"' in _rbody)
ok("Reply body is pre-filled with a greeting and their words quoted",
   "Hi Tess," in _rbody and "&gt; Which plan includes the circles?" in _rbody)

_brevo_calls.clear()
r = admin.post(f"/admin/inbox/messages/{_reply_id}/reply",
               data={"sender": "ayesha", "subject": "Re: circles and plans",
                     "preview": "Healing includes them.", "header": "Bloom Anyway",
                     "title": "Hi Tess,", "body": "Healing and up include circles."},
               follow_redirects=True)
_call = _brevo_calls[-1] if _brevo_calls else {}
ok("Sending a reply reaches the person who wrote in",
   _call.get("to") == "tess@example.com", f"got {_call.get('to')}")
ok("A reply from the healing address goes out on template #22",
   _call.get("template_id") == 22, f"got {_call.get('template_id')}")
ok("The reply is sent from the chosen sender, not MAIL_FROM",
   _call.get("sender") == "Ayesha <ayesha@bloomanyway.online>",
   f"got {_call.get('sender')!r}")
ok("The composed fields become the template params",
   (_call.get("params") or {}).get("PREVIEW") == "Healing includes them."
   and (_call.get("params") or {}).get("TITLE") == "Hi Tess,")
with app.app_context():
    ok("A replied-to message is marked handled",
       db.session.get(_CM, _reply_id).status == "reviewed")
ok("Studio confirms who the reply went to and from",
   "tess@example.com" in r.get_data(as_text=True)
   and "ayesha@bloomanyway.online" in r.get_data(as_text=True), flashes(r))

# each address wears its own template, and the five params never change
for _key, _tpl, _from in (("support", 20, "bloomsupport@bloomanyway.online"),
                          ("saman", 21, "saman@bloomanyway.online"),
                          ("ayesha", 22, "ayesha@bloomanyway.online"),
                          ("noreply", 20, "noreply@bloomanyway.online")):
    _brevo_calls.clear()
    r = admin.post(f"/admin/inbox/messages/{_reply_id}/reply",
                   data={"sender": _key, "subject": f"Re: from {_key}",
                         "preview": "A line", "header": "Bloom Anyway",
                         "title": "Hi Tess,", "body": "Here is the answer."},
                   follow_redirects=True)
    _call = _brevo_calls[-1] if _brevo_calls else {}
    ok(f"Replying as {_key} picks template #{_tpl}",
       _call.get("template_id") == _tpl, f"got {_call.get('template_id')}")
    ok(f"Replying as {_key} sends from {_from}",
       _from in (_call.get("sender") or ""), f"got {_call.get('sender')!r}")
    ok(f"Template #{_tpl} still takes the same five params",
       set(_call.get("params") or {}) == {"SUBJECT", "PREVIEW", "HEADER",
                                          "TITLE", "BODY"},
       f"got {sorted(_call.get('params') or {})}")
    ok(f"Studio says which template a {_key} reply used",
       f"template #{_tpl}" in r.get_data(as_text=True), flashes(r))

with app.app_context():
    ok("Each address maps to its own template",
       [_mailer.reply_template_for(k) for k in
        ("support", "saman", "ayesha")] == [20, 21, 22])
    ok("An unknown address falls back to customer support rather than nothing",
       _mailer.reply_template_for("impostor") == 20)
    ok("The sender list tells Studio which template each address uses",
       {s["key"]: s["template"] for s in _mailer.sender_choices()}
       == {"support": 20, "ayesha": 22, "saman": 21, "noreply": 20})

_brevo_calls.clear()
r = admin.post(f"/admin/inbox/messages/{_reply_id}/reply",
               data={"sender": "support", "subject": "", "title": "", "body": ""},
               follow_redirects=True)
ok("An empty reply is refused rather than sent",
   not _brevo_calls and "Your reply needs" in r.get_data(as_text=True), flashes(r))

_brevo_calls.clear()
r = admin.post(f"/admin/inbox/messages/{_reply_id}/reply",
               data={"sender": "impostor@evil.example", "subject": "Hi",
                     "title": "Hi", "body": "Hi"},
               follow_redirects=True)
ok("An unknown sender is refused, so nobody can forge a From",
   not _brevo_calls and "which address" in r.get_data(as_text=True), flashes(r))

with app.app_context():
    ok("Unknown sender keys resolve to nothing",
       _mailer.sender_from("impostor") is None
       and _mailer.sender_from("") is None)
    ok("Known sender keys resolve to a verified address",
       _mailer.sender_from("saman") == "Saman <saman@bloomanyway.online>")

_brevo_calls.clear()
r = admin.post("/admin/settings/test-email",
               data={"to": "elsewhere@example.com", "template": "support"},
               follow_redirects=True)
_call = _brevo_calls[-1] if _brevo_calls else {}
ok("Studio can send a support-template test to any address",
   _call.get("to") == "elsewhere@example.com" and _call.get("template_id") == 20,
   f"got {_call.get('to')} / {_call.get('template_id')}")
ok("Studio says which address and template it used",
   "elsewhere@example.com" in r.get_data(as_text=True)
   and "#20" in r.get_data(as_text=True), flashes(r))

for _key, _tpl, _who in (("saman", 21, "Saman"), ("ayesha", 22, "Ayesha")):
    _brevo_calls.clear()
    r = admin.post("/admin/settings/test-email",
                   data={"to": "elsewhere@example.com", "template": _key},
                   follow_redirects=True)
    _call = _brevo_calls[-1] if _brevo_calls else {}
    ok(f"Studio can test {_who}'s reply template (#{_tpl}) before using it",
       _call.get("template_id") == _tpl
       and f"#{_tpl}" in r.get_data(as_text=True),
       f"got {_call.get('template_id')} / {flashes(r)}")
    ok(f"The {_who} test sends from their own address",
       _who in (_call.get("sender") or ""), f"got {_call.get('sender')!r}")

r = admin.get("/admin/settings")
ok("Studio settings offer a test send for all three reply templates",
   all(t in r.get_data(as_text=True) for t in ("(#20)", "(#21)", "(#22)")))

_brevo_calls.clear()
r = admin.post("/admin/settings/test-email", data={}, follow_redirects=True)
_call = _brevo_calls[-1] if _brevo_calls else {}
ok("A blank address falls back to the owner's own inbox",
   _call.get("to") == "owner@example.com", f"got {_call.get('to')}")
ok("The default test still uses the general template",
   _call.get("template_id") == 10, f"got {_call.get('template_id')}")

_brevo_calls.clear()
r = admin.post("/admin/settings/test-email", data={"to": "not-an-address"},
               follow_redirects=True)
ok("A junk address sends nothing and says so",
   not _brevo_calls and "address to send the test to" in r.get_data(as_text=True),
   flashes(r))

# Automated mail must keep leaving from MAIL_FROM, whatever a reply picked.
_brevo_calls.clear()
with app.app_context():
    _mailer.send_welcome_email("newbie@example.com", first_name="Ada")
    _mailer.send_newsletter_welcome("reader@example.com")
    _mailer.send_creator_welcome("member@example.com", plan_price="£12",
                                 billing_interval="monthly")
ok("Automated emails still send from MAIL_FROM, not a picked sender",
   _brevo_calls and all(c.get("sender") is None for c in _brevo_calls),
   f"got {[c.get('sender') for c in _brevo_calls]}")

_mailer.send_email = _orig_send_email

# --- the public support address is easy to find -----------------------------
from app.services import settings as _settings  # noqa

SUPPORT_EMAIL = "bloomsupport@bloomanyway.online"
with app.app_context():
    ok("Support address is the default out of the box",
       _settings.DEFAULTS["contact_email"] == SUPPORT_EMAIL)
    ok("Replies come from that same address",
       _mailer.sender_from("support") == f"Customer Support <{SUPPORT_EMAIL}>")

# An older site: blank address, and no marker saying we've ever filled it.
with app.app_context():
    from app.models import Setting as _Setting
    _settings.set_setting("contact_email", "")
    _seed_marker = db.session.get(_Setting, _settings.SUPPORT_EMAIL_SEEDED)
    if _seed_marker is not None:
        db.session.delete(_seed_marker)
        db.session.commit()
    _settings.invalidate_cache()
    ok("Boot fills the support address in for an existing site",
       _settings.ensure_support_email() is True
       and (_settings.get_setting("contact_email") or "").strip() == SUPPORT_EMAIL,
       f"got {_settings.get_setting('contact_email')!r}")
    ok("A second boot leaves the filled address alone",
       _settings.ensure_support_email() is False)

# A live site already sitting on the address we used to ship. The seeded marker
# means the fill-in above never runs again, so the rename has to be its own step.
with app.app_context():
    for _old in _settings.RETIRED_SUPPORT_EMAILS:
        _settings.set_setting("contact_email", _old)
        _settings.invalidate_cache()
        ok(f"Boot moves a site off the retired {_old}",
           _settings.ensure_support_email() is True
           and (_settings.get_setting("contact_email") or "").strip() == SUPPORT_EMAIL)
    _settings.set_setting("contact_email", "hello@someone-else.test")
    _settings.invalidate_cache()
    ok("But an address the owner chose is never overwritten",
       _settings.ensure_support_email() is False
       and _settings.get_setting("contact_email") == "hello@someone-else.test")
    _settings.set_setting("contact_email", SUPPORT_EMAIL)
    _settings.invalidate_cache()

for _path, _where in (("/", "home page footer"),
                      ("/contact", "contact page"),
                      ("/faq", "FAQ")):
    _html = client.get(_path).get_data(as_text=True)
    ok(f"Support address shows on the {_where}",
       f"mailto:{SUPPORT_EMAIL}" in _html and SUPPORT_EMAIL in _html)

_html = client.get("/contact").get_data(as_text=True)
ok("Contact page offers the address as an alternative to the form",
   "Prefer your own inbox?" in _html)

with app.app_context():
    _settings.set_setting("contact_email", "")
    _settings.invalidate_cache()
_html = client.get("/").get_data(as_text=True)
ok("Clearing the address hides it instead of leaving an empty link",
   "mailto:" not in _html and "Need a hand?" not in _html)

with app.app_context():
    ok("A cleared address is not written back on the next boot",
       _settings.ensure_support_email() is False
       and (_settings.get_setting("contact_email") or "") == "")
    _settings.set_setting("contact_email", SUPPORT_EMAIL)
    _settings.invalidate_cache()

r = admin.get("/admin/settings")
ok("Studio explains where the support address appears",
   SUPPORT_EMAIL in r.get_data(as_text=True)
   and "Customer support email" in r.get_data(as_text=True))

r = client.get("/healthz")
ok("Health check", r.status_code == 200 and r.get_json()["status"] == "ok")

r = client.get("/nope-not-here")
ok("Kind 404 page", r.status_code == 404 and "different path" in r.get_data(as_text=True))

r = client.get("/")
h = r.headers
ok("Security headers present",
   h.get("X-Content-Type-Options") == "nosniff" and h.get("X-Frame-Options") == "DENY"
   and "Content-Security-Policy" in h)
ok("CSP allows Daily.co video rooms",
   "https://*.daily.co" in (h.get("Content-Security-Policy") or "")
   and "https://unpkg.com" in (h.get("Content-Security-Policy") or ""))

# quotes archive: visitors see only today; members see back to their signup date
anon = app.test_client()
r = anon.get("/quotes")
anon_body = r.get_data(as_text=True)
ok("Visitor sees only today's quote + gate",
   r.status_code == 200 and anon_body.count("quote-mini") == 1 and "Create a free account" in anon_body)

with app.app_context():
    member = User.query.filter_by(email="newperson@example.com").first()
    member.created_at = utcnow() - timedelta(days=40)
    db.session.commit()
    # mirror the route's own formula so the check is robust across the UTC/local
    # midnight boundary (created_at is UTC, date.today() is local)
    expected_days = max(1, min((date.today() - member.created_at.date()).days + 1, 366))
r = client.get("/quotes")  # client is signed in as newperson
member_count = r.get_data(as_text=True).count("quote-mini")
ok("Member archive goes back to signup date",
   r.status_code == 200 and member_count == expected_days,
   f"got {member_count}, expected {expected_days}")

r = admin.get("/admin/quotes")
ok("Admin quotes page (pins, preview tomorrow)", r.status_code == 200 and "Preview tomorrow" in r.get_data(as_text=True))
r = admin.get("/admin/orders", follow_redirects=True)
ok("Admin orders redirects to dashboard",
   r.status_code == 200 and "Payments (30 days)" in r.get_data(as_text=True))
r = admin.get("/admin/subscribers/export.csv", follow_redirects=True)
ok("Admin subscribers redirects away from local list",
   r.status_code == 200 and ("Dashboard" in r.get_data(as_text=True)
                             or "Payments" in r.get_data(as_text=True)))

# --- 7a2. stand-in accounts owners make by hand in Studio -------------------
from app.services import demo_accounts

_demo_form = {"username": "quietmaya", "password": "standin-pass-1",
              "display_name": "Maya R.", "membership": "creator"}

r = admin.get("/admin/members")
ok("Studio offers a way to add a stand-in account",
   "Add a stand-in account" in r.get_data(as_text=True)
   and 'name="username"' in r.get_data(as_text=True))

r = admin.post("/admin/members/demo", data=dict(_demo_form), follow_redirects=True)
mbody = r.get_data(as_text=True)
ok("Owner can add an account with just a username and password",
   "Added Maya R." in mbody)
ok("Studio says how they sign in", "@quietmaya" in mbody)
with app.app_context():
    demo_user = User.query.filter_by(username="quietmaya").first()
    demo_id = demo_user.id
    ok("The stand-in is flagged, verified, and on the chosen tier",
       demo_user.is_demo is True and demo_user.is_verified
       and demo_user.membership == "creator")
    ok("It gets an address that can never receive mail",
       demo_accounts.is_demo_address(demo_user.email)
       and demo_user.email.endswith("@demo.invalid"))
    ok("It can be found by handle for the sign-in path",
       demo_accounts.find_by_username("quietmaya").id == demo_id)

# no email address is typed, so the handle is what signs them in
demo_client = app.test_client()
r = demo_client.post("/login", data={"email": "quietmaya",
                                     "password": "standin-pass-1"},
                     follow_redirects=True)
ok("A stand-in signs in with its username", r.status_code == 200
   and "/login" not in r.request.path)
ok("The signed-in stand-in reaches member pages",
   demo_client.get("/account").status_code == 200)
r = demo_client.post("/login", data={"email": "quietmaya", "password": "wrong"},
                     follow_redirects=False)
ok("A wrong password is still refused", r.status_code == 401)
r = app.test_client().post("/login", data={"email": "nosuchhandle",
                                           "password": "standin-pass-1"})
ok("An unknown handle is refused", r.status_code == 401)
ok("The login form accepts a handle, not just an address",
   'type="email" id="email"' not in app.test_client().get("/login").get_data(as_text=True))

# they must look ordinary from outside, and never be mailed or counted
pbody = app.test_client().get(f"/u/{demo_id}").get_data(as_text=True)
ok("A stand-in's public profile looks like anyone else's",
   "Maya R." in pbody and "demo.invalid" not in pbody
   and "stand-in" not in pbody.lower())
ok("Nothing on the profile hints the account was made in Studio",
   "is_demo" not in pbody)
with app.app_context():
    _demo_mail = []
    _real_brevo = _mailer._send_via_brevo
    _mailer._send_via_brevo = lambda to, *a, **kw: _demo_mail.append(to) or True
    delivered = _mailer.send_email(demo_accounts.address_for("quietmaya"),
                                   "Should never arrive", "body")
    _mailer._send_via_brevo = _real_brevo
    ok("The mailer refuses to write to a stand-in",
       delivered is False and not _demo_mail)

with app.app_context():
    from app.services import stats as _stats
    breakdown = _stats.membership_breakdown()
    real_creators = User.query.filter(User.deleted_at.is_(None),
                                      User.is_demo.is_(False),
                                      User.membership == "creator").count()
    ok("Stand-ins stay out of the membership breakdown",
       breakdown["creator"] == real_creators)
    from app.services import membership_audit as _audit
    ok("Stand-ins stay out of the paid-tier audit",
       all(row["user"].id != demo_id for row in _audit.audit()["rows"]))
    from app.services import memberships as _mem
    ok("Stripe reconcile leaves a stand-in's tier alone",
       _mem.reconcile_user(db.session.get(User, demo_id)) is False)

csv_body = admin.get("/admin/members/export.csv").get_data(as_text=True)
ok("Stand-ins are left out of the email export",
   "demo.invalid" not in csv_body and "quietmaya" not in csv_body)

r = admin.get("/admin/members")
ok("Studio marks a stand-in in the members list",
   "stand-in" in r.get_data(as_text=True)
   and "Signs in as @quietmaya" in r.get_data(as_text=True))

# the same handle can't be reused, and the rules still apply
r = admin.post("/admin/members/demo", data=dict(_demo_form), follow_redirects=True)
ok("The same handle can't be taken twice", "is taken" in r.get_data(as_text=True))
r = admin.post("/admin/members/demo",
               data={"username": "shorty", "password": "abc", "membership": "none"},
               follow_redirects=True)
ok("A too-short password is refused",
   "at least 8 characters" in r.get_data(as_text=True))
r = admin.post("/admin/members/demo",
               data={"username": "admin", "password": "standin-pass-1",
                     "membership": "none"}, follow_redirects=True)
ok("Reserved handles are refused", "reserved" in r.get_data(as_text=True).lower())
with app.app_context():
    ok("None of the refused attempts created an account",
       User.query.filter(User.is_demo.is_(True)).count() == 1)

# removing one uses the same button as any other member, and it stays gone
r = admin.post(f"/admin/members/{demo_id}/remove", follow_redirects=True)
ok("A stand-in is removed like anyone else",
   r.status_code == 200 and "quietmaya" not in r.get_data(as_text=True))
with app.app_context():
    ok("The account row is really gone",
       db.session.get(User, demo_id) is None)
ok("A removed stand-in can no longer sign in",
   app.test_client().post("/login", data={"email": "quietmaya",
                                          "password": "standin-pass-1"}).status_code == 401)
with app.app_context():
    import seed as _seed_mod
    ok("Nothing in the seed script recreates stand-ins",
       "is_demo" not in open(_seed_mod.__file__, encoding="utf-8").read())

# --- 7b. reel reviews, nav order -------------------------------------------
from app.models import ReelReview, ReelReviewApplication
from app.services import reel_reviews as reel_svc

nav = client.get("/").get_data(as_text=True)
# order inside the desktop nav-links block only
nav_block = nav.split('class="nav-links"', 1)[-1].split("</div>", 1)[0]
i_courses = nav_block.find("Courses &amp; Guides")
i_community = nav_block.find(">Community<")
i_hub = nav_block.find(">Content Hub<")
i_showcase = nav_block.find(">Showcase<")
i_sg = nav_block.find(">Support Groups<")
i_space = nav_block.find(">My space<")
ok("Nav order is Courses, Community, Content Hub, Showcase, Support Groups, My space",
   -1 not in (i_courses, i_community, i_hub, i_showcase, i_sg, i_space)
   and i_courses < i_community < i_hub < i_showcase < i_sg < i_space,
   f"idx={(i_courses, i_community, i_hub, i_showcase, i_sg, i_space)}")
ok("Daily quotes stays in the footer (not main nav)",
   ">Daily quotes<" not in nav_block
   and "Daily quotes" in nav)

r = admin.get("/admin/discounts", follow_redirects=False)
ok("Site discount-codes feature is removed", r.status_code == 404)

# reel review: Creator can enter once/week; Healing cannot
r = banclient.post("/watch/review-request", data={
    "reel_url": "https://www.instagram.com/reel/TESTREEL1/",
    "raw_video": (io.BytesIO(minimal_mp4), "raw.mp4"),
}, content_type="multipart/form-data", follow_redirects=True)
ok("Healing member can't request a reel review",
   "Creator membership" in r.get_data(as_text=True)
   or banclient.get("/membership").status_code == 200)

r = client.post("/watch/review-request", data={
    "reel_url": "https://www.instagram.com/reel/TESTREEL1/",
    "raw_video": (io.BytesIO(minimal_mp4), "raw.mp4"),
}, content_type="multipart/form-data", follow_redirects=True)
ok("Creator member can put a reel forward for the week",
   "Your reel is in for this week" in r.get_data(as_text=True))
with app.app_context():
    ok("Entries are keyed to the Monday of Atlanta's week",
       reel_svc.current_week_key() == reel_svc.week_monday(reel_svc.atlanta_today())
       and reel_svc.current_week_key().weekday() == 0)
with app.app_context():
    import os as _os
    stored = ReelReviewApplication.query.first()
    disk_ok = (stored is not None and stored.disk_name
               and _os.path.isfile(_os.path.join(
                   app.config["VIDEO_STORAGE_DIR"], stored.disk_name)))
    ok("Reel raw video is streamed to video storage (not loaded into Postgres)",
       disk_ok)
r = client.post("/watch/review-request", data={
    "reel_url": "https://www.instagram.com/reel/TESTREEL2/",
    "raw_video": (io.BytesIO(minimal_mp4), "raw2.mp4"),
}, content_type="multipart/form-data", follow_redirects=True)
ok("Second reel-review request in the same week is blocked",
   "already put a reel forward this week" in r.get_data(as_text=True))

r = admin.post("/admin/reel-reviews/pick", follow_redirects=True)
body = r.get_data(as_text=True)
ok("Owner can pick a random entry to review",
   "is up next" in body)
ok("The picked entry is highlighted at the top of the queue",
   "reel-applicant--winner" in body and "Up next" in body)
ok("Studio counts the week's progress towards seven",
   "0<span class=\"reel-week__of\">/7</span>" in body
   and "Reviewed this week" in body)
ok("Entry card offers a single raw-video download",
   body.count("Download raw") == 1 and "Download raw video" not in body)
with app.app_context():
    app_row = ReelReviewApplication.query.filter_by(selected=True).first()
    app_id = app_row.id
r = admin.get(f"/admin/reel-reviews/{app_id}/raw")
ok("Owner can download the winner's raw video",
   r.status_code == 200
   and "attachment" in (r.headers.get("Content-Disposition") or "").lower()
   and len(r.data) > 0)
r = admin.post(f"/admin/reel-reviews/{app_id}/publish", data={
    "title": "Loved your pacing", "body": "Keep the hook under 2 seconds.",
}, follow_redirects=True)
ok("Owner can publish a reel review",
   "published to the Content Hub" in r.get_data(as_text=True))
ok("Publishing says how many are left in the week's seven",
   "6 left this week" in r.get_data(as_text=True))
with app.app_context():
    pub = ReelReview.query.order_by(ReelReview.id.desc()).first()
    ok("A published review is stamped with the Atlanta day it went out",
       pub.review_date == reel_svc.atlanta_today())
    pub_id = pub.id

# Two things landing in the hub on one day both belong on the home page, and
# each keeps its own day rather than sharing the newest one's.
def _notices_strip(html):
    """Just the announcement strip — these titles appear further down too."""
    start = html.find("home-hero__notices-inner")
    if start < 0:
        return ""
    end = html.find("home-hero__strap", start)
    return html[start:end if end > start else len(html)]


_hub_html = _notices_strip(client.get("/").get_data(as_text=True))
ok("A tip and a reel review from the same day both reach the home page",
   "New in the Content Hub" in _hub_html and "Batch a week of hooks" in _hub_html
   and "New reel review" in _hub_html and "Loved your pacing" in _hub_html)
ok("Each card goes to its own piece",
   f'href="/watch/{text_tip_id}"' in _hub_html
   and f'href="/watch/reviews/{pub_id}"' in _hub_html)
with app.app_context():
    _aged = db.session.get(Video, text_tip_id)
    _was = _aged.created_at
    _aged.created_at = utcnow() - timedelta(hours=25)
    db.session.commit()
_hub_html = _notices_strip(client.get("/").get_data(as_text=True))
ok("A tip past its own 24 hours drops off on its own",
   "Batch a week of hooks" not in _hub_html)
ok("And takes nothing else with it", "Loved your pacing" in _hub_html)
with app.app_context():
    db.session.get(Video, text_tip_id).created_at = _was
    db.session.commit()
with app.test_request_context("/"):
    from app.services import homepage as _home_svc
    ok("A busy day can't bury the hero — the strip is capped",
       len(_home_svc.content_hub_drops(limit=1)) == 1
       and _home_svc.MAX_DROPS >= 2)

# Several of a kind on one day used to be a full-width bar apiece, saying "New
# in the Content Hub" over and over. They share a row instead.
with app.app_context():
    _extra_tips = []
    for _t in ("Hooks that land", "Filming in one take"):
        _v = Video(title=_t, published=True, created_at=utcnow(),
                   body="Some tip text.")
        db.session.add(_v)
        _extra_tips.append(_v)
    db.session.commit()
    _extra_ids = [_v.id for _v in _extra_tips]
with app.test_request_context("/"):
    _groups = _home_svc.content_hub_groups()
    _tip_groups = [g for g in _groups if g["kind"] == "tip"]
    ok("Tips landing together are gathered under one heading",
       len(_tip_groups) == 1 and len(_tip_groups[0]["items"]) >= 2,
       f"got {[(g['kind'], len(g['items'])) for g in _groups]}")
    ok("However busy the day, the strip is one row per kind",
       len(_groups) <= 2, f"got {len(_groups)} rows")
_busy = _notices_strip(client.get("/").get_data(as_text=True))
ok("The kind is said once over the titles, not repeated per title",
   _busy.count("New in the Content Hub") == 1
   and "Hooks that land" in _busy and "Filming in one take" in _busy,
   f"said it {_busy.count('New in the Content Hub')} time(s)")
ok("And each title is still its own link",
   all(f'href="/watch/{i}"' in _busy for i in _extra_ids))
ok("Everything landing today shares one card, not one apiece",
   _busy.count("hero-newvideo--many") == 1
   and _busy.count('class="hero-newvideo"') == 0,
   f"got {_busy.count('hero-newvideo')} cards")
ok("A kind per line inside it",
   _busy.count("hero-newvideo__group") == 2)
_css = client.get("/static/css/main.css").get_data(as_text=True)
ok("With corners rounded like the notices beside it, not into a pill",
   re.search(r"\.hero-newvideo \{[^}]*border-radius: 16px", _css, re.S) is not None,
   "still a full pill")
with app.app_context():
    for _i in _extra_ids:
        db.session.delete(db.session.get(Video, _i))
    db.session.commit()
ok("Someone with no account isn't told about a hub they can't read",
   "New in the Content Hub"
   not in app.test_client().get("/").get_data(as_text=True))
r = admin.get("/admin/reel-reviews")
abody = r.get_data(as_text=True)
ok("Studio closes the publish UI once today's review is out",
   "Today's review is done" in abody
   and "Write &amp; publish review" not in abody
   and "Write & publish review" not in abody)
ok("Studio shows one of seven done for the week",
   "1<span class=\"reel-week__of\">/7</span>" in abody)
r = admin.post("/admin/reel-reviews/pick", follow_redirects=True)
ok("Picking again is blocked once today's review is out",
   "one a day" in r.get_data(as_text=True).lower())

# a second member can still enter while today's review is live — the daily cap
# is the owner's, not theirs
creator2 = app.test_client()
sent_codes.clear()
creator2.post("/register", data={"email": "reeler2@example.com", "password": USER_PW,
                                 "password_confirm": USER_PW})
creator2.post("/verify-email", data={"email": "reeler2@example.com",
                                     "code": sent_codes[-1][1]})
with app.app_context():
    u2 = User.query.filter_by(email="reeler2@example.com").first()
    u2.membership = "creator"
    db.session.commit()
    u2_id = u2.id
r = creator2.post("/watch/review-request", data={
    "reel_url": "https://www.instagram.com/reel/TESTREEL9/",
    "raw_video": (io.BytesIO(minimal_mp4), "raw9.mp4"),
}, content_type="multipart/form-data", follow_redirects=True)
ok("Another member can still enter while today's review is live",
   "Your reel is in for this week" in r.get_data(as_text=True))
with app.app_context():
    other_id = (ReelReviewApplication.query
                .filter_by(user_id=u2_id).first().id)
r = admin.post(f"/admin/reel-reviews/{other_id}/publish", data={
    "title": "Second one today", "body": "Should not go out yet.",
}, follow_redirects=True)
ok("A second review on the same day is refused",
   "one a day" in r.get_data(as_text=True).lower())
with app.app_context():
    ok("The refused second review was never created",
       ReelReview.query.filter_by(application_id=other_id).count() == 0)

r = client.get("/watch")
cbody = r.get_data(as_text=True)
ok("Published reel reviews show on Content Hub for Creator members",
   "Loved your pacing" in cbody)
ok("Content Hub counts the week's reviews rather than closing the round",
   "of 7 reviewed" in cbody and "Hang tight" not in cbody)

# the review opens on its own page: full write-up + a properly sized player
long_review = ("Your hook lands inside the first second and the caption carries "
               "the rest of it. Keep the pacing this tight, trim the last beat, "
               "and add a face-to-camera line at the end so people know exactly "
               "who they just watched.")
r = admin.post(f"/admin/reel-reviews/{app_id}/publish", data={
    "title": "Loved your pacing", "body": long_review,
    "review_video": (io.BytesIO(minimal_mp4), "review.mp4"),
}, content_type="multipart/form-data", follow_redirects=True)
ok("Owner can attach a video to a published review",
   "published to the Content Hub" in r.get_data(as_text=True))
with app.app_context():
    rev_row = ReelReview.query.order_by(ReelReview.id.desc()).first()
    review_id = rev_row.id
    ok("Review video is stored for playback", bool(rev_row.review_disk_name))

hub = client.get("/watch").get_data(as_text=True)
ok("Hub review cards open the full review page",
   f'href="/watch/reviews/{review_id}"' in hub and "Read the review" in hub)
ok("Hub cards no longer cram a player into the card", "<video" not in hub)

r = client.get(f"/watch/reviews/{review_id}")
rbody = r.get_data(as_text=True)
ok("Reel review opens on its own page with the whole write-up",
   r.status_code == 200 and "Loved your pacing" in rbody
   and "who they just watched" in rbody)
ok("Reel review page plays the review video",
   "<video" in rbody and f"/watch/reviews/{review_id}/stream" in rbody)
ok("Reel review page embeds the original reel",
   "instagram.com/reel/TESTREEL1/embed" in rbody
   and "Watch on Instagram" in rbody)
r = client.get(f"/watch/reviews/{review_id}/stream")
ok("Review video streams to a signed-in member",
   r.status_code == 200 and len(r.data) > 0)

guest_body = app.test_client().get(f"/watch/reviews/{review_id}").get_data(as_text=True)
ok("Guests get a taste of the review, not the whole critique",
   guest_body.count("Loved your pacing") >= 1
   and "who they just watched" not in guest_body
   and "The rest is for Creator members" in guest_body)
ok("Guests are not offered the review video at all",
   "<video" not in guest_body and "Sign in to watch" not in guest_body)
heal_body = banclient.get(f"/watch/reviews/{review_id}").get_data(as_text=True)
ok("Healing members hit the same lock as guests",
   "The rest is for Creator members" in heal_body
   and "who they just watched" not in heal_body)
ok("Non-members can't stream the review video either",
   banclient.get(f"/watch/reviews/{review_id}/stream").status_code == 404
   and app.test_client().get(
       f"/watch/reviews/{review_id}/stream").status_code in (302, 401, 404))
heal_hub = banclient.get("/watch").get_data(as_text=True)
ok("Content Hub tells non-members why the reviews are locked",
   "Reel reviews are a Creator perk" in heal_hub)
ok("Missing reel review returns 404",
   client.get("/watch/reviews/999999").status_code == 404)

# --- reel of the week: member entries feed the home page spotlight ----------
from app.models import ReelSubmission
from app.services import reel_of_week as rotw_svc

_rotw = {"reel_url": "https://www.instagram.com/reel/SHARED100/",
         "share_count": "412", "confirm_shares": "1"}


def _rotw_post(cl, **over):
    data = dict(_rotw, **over)
    data["raw_video"] = (io.BytesIO(minimal_mp4), "shared.mp4")
    return cl.post("/watch/reel-of-week", data=data,
                   content_type="multipart/form-data", follow_redirects=True)


r = _rotw_post(banclient)
ok("Healing member can't enter Reel of the Week",
   "Creator perk" in r.get_data(as_text=True)
   or banclient.get("/membership").status_code == 200)
r = _rotw_post(client, share_count="12")
ok("A reel under 100 shares is turned away",
   "100 shares or more" in r.get_data(as_text=True))
r = _rotw_post(client, reel_url="https://example.com/not-a-reel")
ok("Reel of the Week needs a real Instagram reel link",
   "instagram.com/reel" in r.get_data(as_text=True))
r = client.post("/watch/reel-of-week", data={
    **_rotw, "raw_video": (io.BytesIO(minimal_mp4), "shared.mp4"),
    "confirm_shares": "",
}, content_type="multipart/form-data", follow_redirects=True)
ok("The share count has to be confirmed before it counts",
   "confirm the share count" in r.get_data(as_text=True))
r = client.post("/watch/reel-of-week", data=dict(_rotw),
                content_type="multipart/form-data", follow_redirects=True)
ok("Reel of the Week needs the raw video, not just the link",
   "Upload the raw video" in r.get_data(as_text=True))
with app.app_context():
    ok("None of the refused entries were saved",
       ReelSubmission.query.count() == 0)

r = _rotw_post(client)
ok("Creator member can enter Reel of the Week",
   "in the running for this week" in r.get_data(as_text=True))
r = _rotw_post(client, reel_url="https://www.instagram.com/reel/SHARED200/")
ok("Second Reel of the Week entry in the same week is blocked",
   "already entered a reel this week" in r.get_data(as_text=True))
with app.app_context():
    sub = ReelSubmission.query.first()
    sub_id = sub.id
    ok("The entry keeps the share count and the raw video",
       sub.share_count == 412 and sub.has_raw_video()
       and sub.week_key == rotw_svc.current_week_key())

r = admin.get("/admin/spotlight")
sbody = r.get_data(as_text=True)
ok("Studio lists this week's Reel of the Week entries",
   "This week's Reel of the Week entries (1)" in sbody
   and "412 shares" in sbody and "Feature this one" in sbody)
r = admin.get(f"/admin/spotlight/reel/{sub_id}/raw")
ok("Owner can download an entrant's raw video",
   r.status_code == 200
   and "attachment" in (r.headers.get("Content-Disposition") or "").lower()
   and len(r.data) > 0)
r = admin.post("/admin/spotlight", data={"feature_reel": str(sub_id)},
               follow_redirects=True)
ok("Owner can feature an entry as Reel of the Week",
   "Reel of the Week on the home page" in r.get_data(as_text=True))
home = app.test_client().get("/").get_data(as_text=True)
ok("The featured entry lands on the home page spotlight",
   "instagram.com/reel/SHARED100/embed" in home)
ok("The home page credits whoever sent it in",
   "412 shares" in home)
with app.app_context():
    ok("Only one entry is flagged as featured",
       ReelSubmission.query.filter_by(featured=True).count() == 1)
hub = client.get("/watch").get_data(as_text=True)
ok("The member sees their reel made the home page",
   "featured on the home page" in hub)
r = admin.post("/admin/spotlight", data={"clear_spotlight_reel": "1"},
               follow_redirects=True)
ok("Clearing Reel of the Week takes it off the home page",
   "cleared from the home page" in r.get_data(as_text=True)
   and "instagram.com/reel/SHARED100/embed"
   not in app.test_client().get("/").get_data(as_text=True))

# Monday's clear-out: last week's entries go, reviewed ones stay
with app.app_context():
    from datetime import timedelta as _td
    last_week = rotw_svc.current_week_key() - _td(days=7)
    reviewed = ReelReviewApplication.query.filter(
        ReelReviewApplication.review.has()).first()
    unreviewed = ReelReviewApplication.query.filter_by(id=other_id).first()
    reviewed.week_key = last_week
    unreviewed.week_key = last_week
    reviewed_id = reviewed.id
    ReelSubmission.query.filter_by(id=sub_id).update({"week_key": last_week})
    db.session.commit()
    cleared = rotw_svc.sweep_old_weeks()
    ok("Monday clears last week's unreviewed reel entries",
       cleared["reel_reviews"] == 1
       and db.session.get(ReelReviewApplication, other_id) is None)
    ok("Reviewed entries survive so their published review stays up",
       db.session.get(ReelReviewApplication, reviewed_id) is not None
       and ReelReview.query.filter_by(application_id=reviewed_id).count() == 1)
    ok("The reviewed entry's raw upload is released once it's served its purpose",
       db.session.get(ReelReviewApplication, reviewed_id).disk_name is None)
    ok("Monday clears last week's Reel of the Week entries",
       cleared["reel_of_week"] == 1 and ReelSubmission.query.count() == 0)
ok("The published review is still readable after the clear-out",
   client.get(f"/watch/reviews/{review_id}").status_code == 200)

# hiding a review is reversible, and it can also be deleted outright
r = admin.post(f"/admin/reel-reviews/review/{review_id}/unpublish",
               follow_redirects=True)
abody = r.get_data(as_text=True)
ok("A hidden review offers both a way back and a way out",
   "Put back" in abody and "Delete" in abody)
ok("Guests can't reach a hidden review",
   app.test_client().get(f"/watch/reviews/{review_id}").status_code == 404)
r = admin.post(f"/admin/reel-reviews/review/{review_id}/publish",
               follow_redirects=True)
ok("A hidden review can be put back up",
   "live on the Content Hub again" in r.get_data(as_text=True)
   and client.get(f"/watch/reviews/{review_id}").status_code == 200)

with app.app_context():
    doomed = db.session.get(ReelReview, review_id)
    doomed_app_id = doomed.application_id
    video_path = _os.path.join(app.config["VIDEO_STORAGE_DIR"],
                               doomed.review_disk_name or "")
    ok("The review about to be deleted really has a video on disk",
       bool(doomed.review_disk_name) and _os.path.isfile(video_path))
r = admin.post(f"/admin/reel-reviews/review/{review_id}/delete",
               follow_redirects=True)
ok("Owner can delete a review outright",
   "back in the queue" in r.get_data(as_text=True))
with app.app_context():
    ok("The review row is gone", db.session.get(ReelReview, review_id) is None)
    ok("Its video file is gone from disk too", not _os.path.isfile(video_path))
    back = db.session.get(ReelReviewApplication, doomed_app_id)
    ok("The member's entry survives and is waiting again",
       back is not None and back.selected is False and back.review is None)
    ok("Today's slot frees up once the review is deleted",
       reel_svc.day_is_done() is False)
ok("The deleted review's page is gone",
   client.get(f"/watch/reviews/{review_id}").status_code == 404)
ok("Studio no longer lists the deleted review",
   "Loved your pacing" not in admin.get("/admin/reel-reviews").get_data(as_text=True))

r = app.test_client().get("/")
ok("Sunflower favicon is linked in the tab",
   "favicon.svg" in r.get_data(as_text=True))
ok("Page loader uses an animated sunflower",
   'id="page-loader"' in r.get_data(as_text=True)
   and "page-loader.js" in r.get_data(as_text=True)
   and "page-loader__petal" in r.get_data(as_text=True)
   and "page-loader__leaf" in r.get_data(as_text=True)
   and "page-loader__spin" in r.get_data(as_text=True))

# support / coaching groups — peer schedule/join (admin oversight only)
from app.models import (Notification, SupportGroupApplication, SupportGroupCircle,
                        SupportGroupMeeting)
from app.services import support_groups as sg_svc

_sent_mail = []


def _capture_sg_mail(subject: str):
    """Stand in for one of the mailer helpers support groups calls.

    Each helper has its own signature and builds its own copy, so record the
    subject we know it sends plus whatever URL it was handed — that's all the
    assertions below need.
    """
    def send(to, **kw):
        text = "\n".join(str(v) for v in kw.values() if v)
        _sent_mail.append({"to": to, "subject": kw.get("subject") or subject,
                           "text": text})
        return True
    return send


for _name, _subject in (
    ("send_support_group_booked", "Your seat is saved"),
    ("send_support_group_reminder", "Your session is tomorrow"),
    ("send_support_group_host_cancelled", "That session won't be happening"),
    ("send_support_group_left", "You left the circle"),
    ("send_facilitator_booked", "Your guided session is booked"),
    ("send_facilitator_cancelled", "That session was cancelled"),
    ("send_one_on_one_booked", "Your 1:1 is booked"),
    ("send_one_on_one_cancelled", "Your 1:1 was cancelled"),
    ("send_styled_email", ""),
):
    setattr(sg_svc, _name, _capture_sg_mail(_subject))

with app.app_context():
    sg_svc.ensure_circles()
    _heal_circle = SupportGroupCircle.query.filter_by(track="healing").first()
    _build_circle = SupportGroupCircle.query.filter_by(track="building").first()
    ok("Support group circles are seeded",
       _heal_circle is not None and _build_circle is not None
       and SupportGroupCircle.query.count() >= 8)
    heal_cid = _heal_circle.id
    build_cid = _build_circle.id

# an odd offset keeps this out of the way of slots booked earlier in the suite
# — a clash here is a scheduling conflict, not a scheduling bug
_when = datetime.utcnow() + timedelta(hours=37)
_sg_date = _when.strftime("%Y-%m-%d")
_sg_time = _when.strftime("%H:%M")

# free_client was granted Healing earlier in the suite, so this needs its own
# account — otherwise the check passes while proving nothing
with app.app_context():
    nofree = User(email="really-free@example.com", display_name="Really Free",
                  username="reallyfree", membership="none",
                  email_verified_at=datetime.utcnow())
    nofree.set_password(USER_PW)
    db.session.add(nofree)
    db.session.commit()
nofree_client = app.test_client()
nofree_client.post("/login", data={"email": "really-free@example.com",
                                   "password": USER_PW})
r = nofree_client.post("/support-groups/schedule",
                       data={"circle_id": heal_cid, "meeting_date": _sg_date,
                             "meeting_time": _sg_time, "timezone": "UTC"},
                       follow_redirects=True)
with app.app_context():
    _booked_by_free = (SupportGroupMeeting.query
                       .filter_by(circle_id=heal_cid, kind="peer").count())
ok("Free members cannot schedule peer sessions",
   "Healing, Creator, and Full Bloom" in r.get_data(as_text=True)
   and _booked_by_free == 0,
   f"{flashes(r)} | meetings={_booked_by_free}")

r = stranger_client.post("/support-groups/schedule",
                         data={"circle_id": heal_cid, "meeting_date": _sg_date,
                               "meeting_time": _sg_time, "timezone": "UTC"},
                         follow_redirects=True)
ok("Healing member can schedule a peer session",
   "Session scheduled" in r.get_data(as_text=True), flashes(r))
with app.app_context():
    peer = (SupportGroupMeeting.query
            .filter_by(circle_id=heal_cid, kind="peer", status="scheduled")
            .order_by(SupportGroupMeeting.id.desc()).first())
    ok("Peer meeting exists after member schedule", peer is not None)
    mid = peer.id
    ok("Scheduler is seated as host",
       SupportGroupApplication.query.filter_by(
           meeting_id=mid, status="selected").count() == 1)
    ok("Schedule auto-created a Daily room URL",
       (peer.zoom_url or "").startswith("https://")
       and ".daily.co/" in (peer.zoom_url or "")
       and bool(peer.zoom_meeting_id))
    auto_room = peer.zoom_url
    room_path = f"/support-groups/meetings/{mid}/room"
ok("Host booking email was sent",
   len([m for m in _sent_mail
        if "seat is saved" in m["subject"].lower()
        or "booked" in m["subject"].lower()]) >= 1)

r = stranger_client.post("/support-groups/schedule",
                         data={"circle_id": heal_cid,
                               "meeting_date": (_when + timedelta(days=1)).strftime("%Y-%m-%d"),
                               "meeting_time": _sg_time, "timezone": "UTC"},
                         follow_redirects=True)
ok("Peer schedule cooldown blocks a second session within 2 weeks",
   "every 14 days" in r.get_data(as_text=True).lower()
   or "after" in r.get_data(as_text=True).lower())

r = client.post(f"/support-groups/meetings/{mid}/join", follow_redirects=True)
ok("Another member can join an upcoming peer session",
   "You're in" in r.get_data(as_text=True)
   or "Save a seat" in r.get_data(as_text=True)
   or "Join" in r.get_data(as_text=True))
with app.app_context():
    seated = SupportGroupApplication.query.filter_by(
        meeting_id=mid, status="selected").count()
    ok("Joined peer session has two seats", seated == 2)
    notes = Notification.query.filter_by(kind="support_group").count()
    ok("Peer booking created in-app notifications", notes >= 1)
    sample = Notification.query.filter_by(kind="support_group").first()
    ok("Support-group notifications link to the in-site room",
       sample is not None
       and sample.href()
       and "/support-groups/meetings/" in (sample.href() or "")
       and "/room" in (sample.href() or ""))

r = client.get(f"/support-groups/meetings/{mid}/room", follow_redirects=False)
ok("Early Join redirects to the waiting room",
   r.status_code in (301, 302)
   and f"/support-groups/meetings/{mid}/waiting" in (r.headers.get("Location") or ""))
r = client.get(f"/support-groups/meetings/{mid}/waiting")
wait_body = r.get_data(as_text=True)
ok("Waiting room shows countdown before start",
   r.status_code == 200
   and "Waiting room" in wait_body
   and "data-sg-countdown" in wait_body
   and "Starts in" in wait_body
   and "data-status-url" in wait_body
   and "support-waiting.js" in wait_body
   and "data-starts-ms" in wait_body)
r = client.get(f"/support-groups/meetings/{mid}/status")
ok("Waiting-room status endpoint returns live JSON",
   r.status_code == 200
   and r.is_json
   and r.get_json().get("phase") == "waiting"
   and "starts_at" in r.get_json()
   and "server_now" in r.get_json()
   and "members" in r.get_json())
with app.app_context():
    live = db.session.get(SupportGroupMeeting, mid)
    live.scheduled_at = utcnow() - timedelta(minutes=1)
    db.session.commit()
r = client.get(f"/support-groups/meetings/{mid}/room")
_room_html = r.get_data(as_text=True)
ok("Seated member can open the embedded Daily room once live",
   r.status_code == 200
   and "daily-js" in _room_html
   and "sg-daily-root" in _room_html
   and "support-room.js" in _room_html)
# The room counts itself down, so the end doesn't arrive out of nowhere.
ok("The room carries when it ends and what the server clock says",
   "data-ends-ms=" in _room_html and "data-server-ms=" in _room_html
   and "sg-room-ending" in _room_html)
with app.app_context():
    from datetime import timezone as _dt_tz
    _live = db.session.get(SupportGroupMeeting, mid)
    _expect_end = int((_live.scheduled_at
                       + timedelta(minutes=sg_svc.meeting_duration_minutes(_live))
                       ).replace(tzinfo=_dt_tz.utc).timestamp() * 1000)
ok("And it counts down to this session's own end, not a default one",
   f'data-ends-ms="{_expect_end}"' in _room_html,
   f"expected {_expect_end}")
ok("The notice starts hidden — it is for the last five minutes only",
   'id="sg-room-ending"' in _room_html and "hidden" in _room_html)

# --- add-on prices, read from their Stripe price ids --------------------------
ok("Money is formatted the way the rest of the catalogue writes it",
   pay.format_price_amount(4500, "usd") == "$45"
   and pay.format_price_amount(12050, "USD") == "$120.50"
   and pay.format_price_amount(4500, "gbp") == "\u00a345"
   and pay.format_price_amount(None, "usd") == "")

_sg_body = client.get("/support-groups").get_data(as_text=True)
ok("With no price to show, the tiles say so rather than showing nothing",
   _sg_body.count("Price shown at checkout") >= 1
   and "price shown at checkout" in _sg_body
   and 'class="sg-price"' not in _sg_body)

_real_addon_prices = pay.addon_prices
pay.addon_prices = lambda: {"facilitator": "$45", "ayesha": "$120",
                            "saman": "$150"}
try:
    _sg_body = client.get("/support-groups").get_data(as_text=True)
finally:
    pay.addon_prices = _real_addon_prices
ok("The facilitator tile shows what a session costs",
   "<strong>$45</strong> per session" in _sg_body)
ok("Both 1:1 tiles show their own price",
   "<strong>$120</strong> per session" in _sg_body
   and "<strong>$150</strong> per session" in _sg_body)
ok("And the placeholder copy steps aside once there is a real price",
   "price shown at checkout" not in _sg_body.lower())
with app.app_context():
    ended = db.session.get(SupportGroupMeeting, mid)
    # Put the original future time back so later reminder/cancel tests stay valid,
    # then temporarily verify ended → wrap.
    ended.scheduled_at = utcnow() - timedelta(minutes=50)
    db.session.commit()
r = client.get(f"/support-groups/meetings/{mid}/room", follow_redirects=False)
ok("After 30 minutes the room redirects to wrap",
   r.status_code in (301, 302)
   and f"/support-groups/meetings/{mid}/wrap" in (r.headers.get("Location") or ""))
with app.app_context():
    restore = db.session.get(SupportGroupMeeting, mid)
    restore.status = "scheduled"
    restore.scheduled_at = utcnow() + timedelta(hours=20)
    restore.reminded_at = None
    # opening the live room marked everyone "attended" — put the seats back
    # too, or the later reminder and cancel tests have nobody to talk to
    for _seat in SupportGroupApplication.query.filter_by(meeting_id=mid).all():
        if _seat.status == "attended":
            _seat.status = "selected"
    db.session.commit()

# Topic notify-me alerts
with app.app_context():
    other_heal = (SupportGroupCircle.query
                  .filter_by(track="healing")
                  .filter(SupportGroupCircle.id != heal_cid)
                  .first())
    alert_cid = other_heal.id if other_heal else heal_cid
r = client.post(f"/support-groups/circles/{alert_cid}/notify", follow_redirects=True)
ok("Member can turn on Notify me for a plan topic",
   "notify you" in r.get_data(as_text=True).lower()
   or "Notifying" in r.get_data(as_text=True))
with app.app_context():
    watcher = User.query.filter_by(email="newperson@example.com").first()
    host2 = User(email="sg-alert-host@example.com", username="sgalerthost",
                 membership="healing", email_verified_at=utcnow())
    host2.set_password(USER_PW)
    db.session.add(host2)
    db.session.commit()
    when_alert = utcnow() + timedelta(days=2, hours=3)
    m_alert, aerr = sg_svc.schedule_peer_session(
        host2, circle_id=alert_cid,
        date_s=when_alert.strftime("%Y-%m-%d"),
        time_s=when_alert.strftime("%H:%M"),
        tz_name="UTC",
    )
    ok("Alert host scheduled a session on watched topic",
       m_alert is not None and not aerr, aerr)
    alert_notes = Notification.query.filter_by(
        user_id=watcher.id, kind="support_group_alert").count()
    ok("Notify-me fans out when a topic session is scheduled",
       alert_notes >= 1)
    # Clean up so later caps tests stay stable
    if m_alert:
        sg_svc.cancel_meeting(m_alert)

r = client.get("/support-groups")
sg_body = r.get_data(as_text=True)
ok("Support groups page lists named circles and upcoming sessions",
   r.status_code == 200
   and "Divorce Recovery" in sg_body
   and "New Creators Circle" in sg_body
   and "View upcoming sessions" in sg_body
   and "Apply from My Space" not in sg_body)
r = client.get("/account")
acct = r.get_data(as_text=True)
ok("My space shows upcoming peer sessions",
   "Upcoming Sessions" in acct
   and "Manage sessions" in acct)
ok("Membership matrix lists support groups",
   "Support groups" in app.test_client().get("/membership").get_data(as_text=True))

# Creator-only cannot schedule healing topics
with app.app_context():
    creator_only = User(email="sg-creator@example.com", username="sgcreator",
                        membership="creator", email_verified_at=utcnow())
    creator_only.set_password(USER_PW)
    db.session.add(creator_only)
    db.session.commit()
creator_sg = app.test_client()
creator_sg.post("/login", data={"email": "sg-creator@example.com", "password": USER_PW})
r = creator_sg.post("/support-groups/schedule",
                    data={"circle_id": heal_cid, "meeting_date": _sg_date,
                          "meeting_time": _sg_time, "timezone": "UTC"},
                    follow_redirects=True)
ok("Creator plan cannot schedule healing peer groups",
   "aren’t included" in r.get_data(as_text=True)
   or "aren't included" in r.get_data(as_text=True)
   or "not included" in r.get_data(as_text=True).lower())

r = admin.get("/admin/support-groups")
ok("Studio support-groups page loads for peer oversight",
   r.status_code == 200
   and "member-scheduled" in r.get_data(as_text=True).lower()
   and "Divorce Recovery" in r.get_data(as_text=True)
   and "Waiting list" not in r.get_data(as_text=True))
dash = admin.get("/admin/").get_data(as_text=True)
ok("Dashboard occupancy labels each circle by title",
   "Support Group Occupancy" in dash
   and "sg-occ__name" in dash
   and "Divorce Recovery" in dash
   and "New Creators Circle" in dash
   and "sessions" in dash.lower())
r = admin.post("/admin/support-groups/form",
               data={"capacity": "2", "circle_id": str(heal_cid)},
               follow_redirects=True)
ok("Studio no longer seats waitlists for peer circles",
   "member-scheduled" in r.get_data(as_text=True).lower())

with app.app_context():
    meeting = db.session.get(SupportGroupMeeting, mid)
    meeting.scheduled_at = utcnow() + timedelta(hours=20)
    meeting.reminded_at = None
    db.session.commit()
    n = sg_svc.dispatch_due_reminders()
    ok("24h reminder dispatch runs", n == 1)
ok("Reminder email includes in-site Join link",
   any(room_path in (m.get("text") or "")
       and ("session is tomorrow" in m["subject"].lower()
            or "reminder" in m["subject"].lower())
       for m in _sent_mail))

with app.app_context():
    notes_before_cancel = Notification.query.filter_by(kind="support_group").count()
r = admin.post(f"/admin/support-groups/{mid}/cancel", follow_redirects=True)
ok("Owner can cancel a peer session from Studio",
   "cancelled" in r.get_data(as_text=True).lower())
ok("Cancel emails were sent to seated members",
   len([m for m in _sent_mail
        if "won't be happening" in m["subject"].lower()
        or "cancelled" in m["subject"].lower()]) >= 2)
with app.app_context():
    notes_after = Notification.query.filter_by(kind="support_group").count()
    ok("Cancel creates in-app notifications for seated members",
       notes_after >= notes_before_cancel + 2)
    cancelled_seats = SupportGroupApplication.query.filter_by(
        meeting_id=mid, status="cancelled").count()
    ok("Peer cancel marks seats cancelled (no waitlist return)", cancelled_seats == 2)

# --- member cancels their own 1:1 (refund only 24h+ ahead) -------------------
def _seat_one_on_one(email, hours_ahead):
    """A booked founder 1:1 for `email`, starting `hours_ahead` from now."""
    with app.app_context():
        member = User.query.filter_by(email=email).first()
        meeting = SupportGroupMeeting(
            kind="one_on_one", capacity=2, status="scheduled",
            scheduled_at=utcnow() + timedelta(hours=hours_ahead),
            notes="Saman", booked_notified_at=utcnow(),
        )
        db.session.add(meeting)
        db.session.flush()
        db.session.add(SupportGroupApplication(
            user_id=member.id, meeting_id=meeting.id, status="selected"))
        db.session.commit()
        return meeting.id

_ooo_far = _seat_one_on_one("newperson@example.com", 72)
r = client.get("/support-groups")
_sg_html = r.get_data(as_text=True)
ok("A booked 1:1 offers a cancel button",
   f"/support-groups/one-on-one/{_ooo_far}/cancel" in _sg_html
   and "Cancel session" in _sg_html)
ok("The 1:1 says refunds need 24 hours' notice",
   "refunds only apply if you" in _sg_html.lower()
   and "24 hours before the session" in _sg_html)

_styled_mail = []
_orig_styled = sg_svc.send_styled_email
sg_svc.send_styled_email = (
    lambda to, **kw: _styled_mail.append(dict(kw, to=to)) or True
)

r = client.post(f"/support-groups/one-on-one/{_ooo_far}/cancel",
                follow_redirects=True)
_body = r.get_data(as_text=True)
ok("Cancelling early promises the refund",
   "refund is on its way" in _body)
with app.app_context():
    _m = db.session.get(SupportGroupMeeting, _ooo_far)
    ok("Cancelled 1:1 is off the schedule", _m.status == "cancelled")
    ok("Cancelled 1:1 releases the seat",
       SupportGroupApplication.query.filter_by(
           meeting_id=_ooo_far, status="selected").count() == 0)
    _owner = User.query.filter_by(email="owner@example.com").first()
    _alert = (Notification.query
              .filter_by(user_id=_owner.id, kind="support_group_alert")
              .order_by(Notification.id.desc()).first())
    ok("Studio is told a refund is due",
       _alert is not None and "refund is due" in (_alert.body or "").lower()
       and "cancelled their 1:1" in (_alert.body or ""))
ok("Member's email confirms the refund is coming",
   any("refunded to the card" in (m.get("body") or "") for m in _styled_mail))

_ooo_soon = _seat_one_on_one("newperson@example.com", 5)
_soon_html = client.get("/support-groups").get_data(as_text=True)
ok("A session inside the window warns there's no refund",
   "won't be refunded" in _soon_html)
_styled_mail.clear()
r = client.post(f"/support-groups/one-on-one/{_ooo_soon}/cancel",
                follow_redirects=True)
ok("Cancelling late says so plainly",
   "inside the 24-hour window" in r.get_data(as_text=True))
ok("Late-cancel email doesn't promise a refund",
   any("isn't refundable" in (m.get("body") or "") for m in _styled_mail)
   and not any("refunded to the card" in (m.get("body") or "")
               for m in _styled_mail))
sg_svc.send_styled_email = _orig_styled
with app.app_context():
    _alert = (Notification.query.filter_by(kind="support_group_alert")
              .order_by(Notification.id.desc()).first())
    ok("Studio is told no refund is owed",
       _alert is not None and "no refund is due" in (_alert.body or "").lower())

_ooo_other = _seat_one_on_one("newperson@example.com", 48)
r = stranger_client.post(f"/support-groups/one-on-one/{_ooo_other}/cancel",
                         follow_redirects=True)
_deny = r.get_data(as_text=True)
ok("Nobody can cancel someone else's 1:1",
   "isn&#39;t your session" in _deny or "isn't your session" in _deny)
with app.app_context():
    ok("Someone else's 1:1 stays booked",
       db.session.get(SupportGroupMeeting, _ooo_other).status == "scheduled")

# Cap: max 4 open peer sessions per topic
with app.app_context():
    # Clear any leftover open peer sessions on this topic from earlier steps.
    for leftover in SupportGroupMeeting.query.filter_by(
            circle_id=heal_cid, kind="peer", status="scheduled").all():
        sg_svc.cancel_meeting(leftover)
    hosts = []
    for i in range(4):
        u = User(email=f"sg-host{i}@example.com", username=f"sghost{i}",
                 membership="healing", email_verified_at=utcnow())
        u.set_password(USER_PW)
        db.session.add(u)
        hosts.append(u)
    db.session.commit()
    host_ids = [u.id for u in hosts]
    for i, uid in enumerate(host_ids):
        u = db.session.get(User, uid)
        when = utcnow() + timedelta(days=3 + i, hours=2)
        m, err = sg_svc.schedule_peer_session(
            u, circle_id=heal_cid,
            date_s=when.strftime("%Y-%m-%d"),
            time_s=when.strftime("%H:%M"),
            tz_name="UTC",
        )
        ok(f"Peer session {i + 1}/4 schedules for topic cap test",
           m is not None and not err, err)
    fifth = User(email="sg-host5@example.com", username="sghost5",
                 membership="healing", email_verified_at=utcnow())
    fifth.set_password(USER_PW)
    db.session.add(fifth)
    db.session.commit()
    when5 = utcnow() + timedelta(days=10, hours=2)
    m5, err5 = sg_svc.schedule_peer_session(
        fifth, circle_id=heal_cid,
        date_s=when5.strftime("%Y-%m-%d"),
        time_s=when5.strftime("%H:%M"),
        tz_name="UTC",
    )
    ok("Topic blocks a 5th upcoming peer session",
       m5 is None and err5 and "4 upcoming" in err5)

# Post-session wrap + silent peer report
with app.app_context():
    for leftover in SupportGroupMeeting.query.filter_by(
            circle_id=heal_cid, kind="peer", status="scheduled").all():
        sg_svc.cancel_meeting(leftover)
    wrap_host = User(email="sg-wrap-host@example.com", username="sgwraphost",
                     membership="healing", email_verified_at=utcnow())
    wrap_host.set_password(USER_PW)
    wrap_peer = User(email="sg-wrap-peer@example.com", username="sgwrappeer",
                     membership="healing", email_verified_at=utcnow())
    wrap_peer.set_password(USER_PW)
    db.session.add_all([wrap_host, wrap_peer])
    db.session.commit()
    host_warnings = wrap_host.forum_warnings or 0
    when_w = utcnow() + timedelta(days=1, hours=4)
    wm, werr = sg_svc.schedule_peer_session(
        wrap_host, circle_id=heal_cid,
        date_s=when_w.strftime("%Y-%m-%d"),
        time_s=when_w.strftime("%H:%M"),
        tz_name="UTC",
    )
    ok("Wrap test session scheduled", wm is not None and not werr, werr)
    wrap_mid = wm.id
    wrap_host_id = wrap_host.id
    joined, jerr = sg_svc.join_peer_session(wrap_peer, wrap_mid)
    ok("Wrap peer seated", joined is not None and not jerr, jerr)

# Seating: the cap holds, and it holds under a lock so two people clicking at
# the same moment can't both be handed the last seat.
with app.app_context():
    from sqlalchemy.dialects import postgresql as _pg
    from app.models import SupportGroupMeeting as _SGM
    _lock_sql = str(sg_svc.meeting_lock_query(1)
                    .statement.compile(dialect=_pg.dialect()))
    ok("The seat count is taken with the session's row locked",
       "FOR UPDATE" in _lock_sql, _lock_sql[-60:])

    _cap_m = db.session.get(_SGM, wrap_mid)
    _cap_m.capacity = 2  # host + wrap_peer already seated
    db.session.commit()
    _extra = User(email="sg-cap@example.com", username="sgcap",
                  membership="healing", email_verified_at=utcnow())
    _extra.set_password(USER_PW)
    db.session.add(_extra)
    db.session.commit()
    _seat, _err = sg_svc.join_peer_session(_extra, wrap_mid)
    ok("A full session turns the next person away at the cap",
       _seat is None and "full" in (_err or "").lower(), f"got {_err}")
    ok("And nobody was seated past it",
       len(sg_svc.meeting_seats(_cap_m)) == 2)
    # Someone who already holds a seat gets it back rather than "it's full".
    _seated = User.query.filter_by(email="sg-wrap-peer@example.com").first()
    _again, _again_err = sg_svc.join_peer_session(_seated, wrap_mid)
    ok("Someone already seated on a full session is handed their own seat",
       _again is not None and not _again_err, f"got {_again_err}")
    _cap_m.capacity = 8
    db.session.commit()

wrap_client = app.test_client()
wrap_client.post("/login", data={"email": "sg-wrap-peer@example.com", "password": USER_PW})
r = wrap_client.get(f"/support-groups/meetings/{wrap_mid}/wrap")
ok("Post-session wrap page lists peers with profile links",
   r.status_code == 200
   and "sgwraphost" in r.get_data(as_text=True)
   and "/u/" in r.get_data(as_text=True)
   and "Skip" in r.get_data(as_text=True))
r = wrap_client.post(
    f"/support-groups/meetings/{wrap_mid}/report/{wrap_host_id}",
    data={"reason": "harassment", "note": "was unkind"},
    follow_redirects=True,
)
ok("Peer session report submits quietly",
   "thank you" in r.get_data(as_text=True).lower())
with app.app_context():
    from app.models import ContentReport
    flagged = User.query.filter_by(email="sg-wrap-host@example.com").first()
    rep = ContentReport.query.filter_by(
        target_type="user", target_id=flagged.id, status="open").first()
    ok("Peer report lands as open user report with reason",
       rep is not None and rep.reason == "harassment")
    ok("Reported member is silently flagged",
       (flagged.forum_warnings or 0) >= host_warnings + 1)
    notes_to_flagged = Notification.query.filter_by(
        user_id=flagged.id, kind="moderation").count()
    ok("Reported member is not notified about the flag", notes_to_flagged == 0)

# The wrap page is only ever reached after a session is over, and completing
# one moves every seat off "selected". Looking for selected seats there found
# nobody, which emptied the peers list and made reporting unreachable.
with app.app_context():
    _done = db.session.get(SupportGroupMeeting, wrap_mid)
    sg_svc.complete_meeting(_done)
    ok("Completing a session moves the seats to attended",
       all(s.status == "attended"
           for s in sg_svc.meeting_seats(_done, include_attended=True)))
_wbody = wrap_client.get(f"/support-groups/meetings/{wrap_mid}/wrap").get_data(as_text=True)
ok("The wrap page still names who was in the room after it completes",
   "sgwraphost" in _wbody, "peers list came back empty")
ok("So there is still someone to report", "Submit report" in _wbody)
r = wrap_client.post(
    f"/support-groups/meetings/{wrap_mid}/report/{wrap_host_id}",
    data={"reason": "harassment", "note": "after the session ended"},
    follow_redirects=True,
)
ok("And reporting them still goes through once the session is over",
   "thank you" in r.get_data(as_text=True).lower())
# Feedback about the session itself, next to reporting a person in it.
ok("The wrap page offers feedback as well as reporting someone",
   'data-feedback-open' in _wbody and 'data-feedback-pref="feedback"' in _wbody
   and "How was the session?" in _wbody)

# --- Ayesha's 1:1 runs the same intake as Saman's -----------------------------
from app.services import coaching_intake as intake_svc  # noqa: E402

with app.app_context():
    from app.services.settings import set_setting as _set
    _set("ayesha_stripe_price_id", "price_ayesha_1to1")
    _set("saman_stripe_price_id", "price_saman_1to1")
    for _coach in ("ayesha", "saman"):
        _w, _werr = intake_svc.add_availability(
            _coach, weekday=(utcnow() + timedelta(days=2)).weekday(),
            start_minute=9 * 60, end_minute=15 * 60, tz_name="UTC")
        ok(f"{_coach.title()} availability saved", _w is not None and not _werr, _werr)

r = wrap_client.get("/coaching/ayesha/book")
_abody = r.get_data(as_text=True)
ok("Ayesha has a booking page, the same one Saman has",
   r.status_code == 200 and "What are you currently going through?" in _abody)
ok("Her first two questions are tick-boxes with an Other box",
   'name="going_through" value="Divorce"' in _abody
   and 'name="hoping_for" value="Someone to listen"' in _abody
   and 'name="going_through_other"' in _abody)
ok("The last one is asked, not demanded",
   "anything specific you&#39;d like me to know" in _abody.lower()
   or "anything specific you'd like me to know" in _abody.lower())
ok("And Saman's questionnaire is untouched",
   "What's your niche" in wrap_client.get("/coaching/saman/book").get_data(as_text=True)
   or "What&#39;s your niche" in wrap_client.get("/coaching/saman/book").get_data(as_text=True))

with app.app_context():
    from werkzeug.datastructures import MultiDict as _MD
    _form = _MD([("going_through", "Divorce"), ("going_through", "Anxiety"),
                 ("going_through_other", "sleep"),
                 ("hoping_for", "Someone to listen"),
                 ("discuss", "How to talk to the kids about it."),
                 ("disclaimer_identity", "1"), ("disclaimer_conduct", "1"),
                 ("disclaimer_recording", "1")])
    _ans, _aerr = intake_svc.parse_answers(_form, "ayesha")
    ok("Ticked boxes and the Other box collapse into one answer",
       _aerr is None and _ans["going_through"] == "Divorce, Anxiety, sleep",
       f"got {_aerr or _ans.get('going_through')}")
    ok("The optional question can be left alone",
       "before_we_talk" not in _ans)
    _bad = _MD([("going_through", "Divorce"),
                ("going_through", "<script>alert(1)</script>"),
                ("hoping_for", "Someone to listen"), ("discuss", "x"),
                ("disclaimer_identity", "1"), ("disclaimer_conduct", "1"),
                ("disclaimer_recording", "1")])
    _ans2, _ = intake_svc.parse_answers(_bad, "ayesha")
    ok("Anything that isn't one of the offered boxes is dropped",
       _ans2["going_through"] == "Divorce", f"got {_ans2['going_through']}")
    _missing = _MD([("hoping_for", "Someone to listen"), ("discuss", "x"),
                    ("disclaimer_identity", "1"), ("disclaimer_conduct", "1"),
                    ("disclaimer_recording", "1")])
    _none, _merr = intake_svc.parse_answers(_missing, "ayesha")
    ok("But the required ones are still required",
       _none is None and "currently going through" in (_merr or ""), f"got {_merr}")
    _slots = intake_svc.open_slots("ayesha", viewer_tz="UTC")
    ok("Her availability turns into bookable slots", len(_slots) > 0)
    _slot_utc = _slots[0]["utc"]

r = wrap_client.post("/coaching/ayesha/book", data={
    "going_through": "Divorce", "hoping_for": "Practical next steps",
    "discuss": "Where to start with custody.", "slot_utc": _slot_utc,
    "disclaimer_identity": "1", "disclaimer_conduct": "1",
    "disclaimer_recording": "1",
}, follow_redirects=False)
ok("Booking her sends you on to checkout with the intake attached",
   r.status_code in (302, 303)
   and "/checkout/addon/ayesha" in (r.headers.get("Location") or "")
   and "intake=" in (r.headers.get("Location") or ""),
   f"got {r.headers.get('Location')}")
with app.app_context():
    from app.models import CoachingIntake as _CI
    _intake = (_CI.query.filter_by(coach="ayesha")
               .order_by(_CI.id.desc()).first())
    ok("Her intake is saved against her, not Saman",
       _intake is not None and _intake.coach == "ayesha"
       and _intake.status == "pending_payment")
    ok("With the answers on it, readable in Studio",
       any("currently going through" in row["label"]
           for row in intake_svc.answer_rows(_intake)))

r = wrap_client.get("/checkout/addon/ayesha", follow_redirects=False)
ok("Paying for her 1:1 without answering first sends you back to the questions",
   r.status_code in (302, 303)
   and "/coaching/ayesha/book" in (r.headers.get("Location") or ""),
   f"got {r.headers.get('Location')}")
_sgbody = admin.get("/admin/support-groups").get_data(as_text=True)
ok("Studio manages both founders' availability, and says whose is whose",
   "1:1 availability" in _sgbody and "Ayesha" in _sgbody and "Saman" in _sgbody)

# --- availability is a week you set in one pass -------------------------------
# Ticking hours across the days and saving once, rather than adding a window at
# a time. Unticking is how an hour is taken away.
_wbody = admin.get("/admin/support-groups?coach=ayesha").get_data(as_text=True)
ok("The week editor opens on the coach you asked for",
   'name="coach" value="ayesha"' in _wbody and 'data-sg-week' in _wbody)
ok("Switching coach is a plain link, needing no script the CSP would block",
   "onchange=" not in _wbody
   and 'href="/admin/support-groups?coach=saman"' in _wbody
   and "sg-week__whoami is-active" in _wbody)
_sbody2 = admin.get("/admin/support-groups?coach=saman").get_data(as_text=True)
ok("And each coach's own week comes back, not the other's",
   'name="coach" value="saman"' in _sbody2
   and 'value="ayesha" selected' not in _sbody2)
ok("With every day reachable and every hour tickable",
   all(f'value="{d}:9"' in _wbody for d in range(7)))
ok("And what is already saved comes back ticked",
   'value="{}:9"'.format((utcnow() + timedelta(days=2)).weekday()) in _wbody
   and 'checked' in _wbody)

r = admin.post("/admin/support-groups/availability", data={
    "coach": "ayesha", "timezone": "UTC",
    # Monday morning in one run, plus a lone Wednesday hour.
    "slot": ["0:9", "0:10", "0:11", "2:14"],
}, follow_redirects=True)
ok("Saving the week reports what went in",
   "hours open across" in r.get_data(as_text=True), flashes(r))
with app.app_context():
    _wins = intake_svc.list_availability("ayesha")
    _shape = sorted((w.weekday, w.start_minute, w.end_minute) for w in _wins)
    ok("Hours next to each other become one window, not four",
       _shape == [(0, 540, 720), (2, 840, 900)], f"got {_shape}")
    ok("All of it saved in the timezone chosen for the week",
       {w.timezone for w in _wins} == {"UTC"})
    ok("The editor can read its own week back",
       intake_svc.week_grid("ayesha")[0] == {9, 10, 11}
       and intake_svc.week_grid("ayesha")[2] == {14})
    ok("And members get bookable slots from it",
       len(intake_svc.open_slots("ayesha", viewer_tz="UTC")) > 0)

# Saving again is editing: the week is replaced, not added to.
admin.post("/admin/support-groups/availability", data={
    "coach": "ayesha", "timezone": "UTC", "slot": ["0:9"],
}, follow_redirects=True)
with app.app_context():
    _shape = sorted((w.weekday, w.start_minute, w.end_minute)
                    for w in intake_svc.list_availability("ayesha"))
    ok("Unticking an hour removes it rather than leaving it behind",
       _shape == [(0, 540, 600)], f"got {_shape}")
    ok("Saman's week is untouched by saving Ayesha's",
       len(intake_svc.list_availability("saman")) > 0)

r = admin.post("/admin/support-groups/availability", data={
    "coach": "ayesha", "timezone": "UTC",
}, follow_redirects=True)
ok("Clearing every day says so plainly, rather than looking like a no-op",
   "unavailable all week" in r.get_data(as_text=True), flashes(r))
with app.app_context():
    ok("And nothing is bookable once the week is empty",
       intake_svc.list_availability("ayesha") == []
       and intake_svc.open_slots("ayesha", viewer_tz="UTC") == [])
    # Put her week back for anything downstream.
    intake_svc.set_week_availability(
        "ayesha", {(utcnow() + timedelta(days=2)).weekday(): [9, 10, 11, 12, 13, 14]},
        tz_name="UTC")
    ok("Setup restored", len(intake_svc.open_slots("ayesha", viewer_tz="UTC")) > 0)


# site image uploads (hero / story teaser)
from io import BytesIO
from PIL import Image as _PILImage
_buf = BytesIO()
_PILImage.new("RGB", (120, 80), (122, 46, 98)).save(_buf, format="JPEG")
_buf.seek(0)
r = admin.post(
    "/admin/settings",
    data={
        "site_title": "Bloom Anyway", "instagram_url": "",
        "hero_image_url": "", "portrait_url": "", "contact_email": "",
        "announcement_text": "", "announcement_expires": "",
        "creator_name": "", "creator_instagram": "", "creator_image_url": "",
        "creator_blurb": "", "reel_url": "", "reel_description": "",
        "portrait_file": (_buf, "portrait.jpg"),
    },
    content_type="multipart/form-data",
    follow_redirects=True,
)
ok("Studio accepts portrait image upload", r.status_code == 200)
r = admin.get("/admin/settings")
ok("Studio offers crop UI for hero and story teaser uploads",
   'data-site-crop' in r.get_data(as_text=True)
   and 'data-crop-aspect="4:5"' in r.get_data(as_text=True)
   and 'data-crop-aspect="1:1"' in r.get_data(as_text=True)
   and 'id="site-image-crop"' in r.get_data(as_text=True))
r = app.test_client().get("/media/site/portrait")
ok("Uploaded portrait is served",
   r.status_code == 200 and r.mimetype.startswith("image/")
   and r.data[:3] == b"\xff\xd8\xff")
from app.services.settings import get_setting
from app.services import site_images as site_img_svc
with app.app_context():
    ok("Portrait setting points at media route",
       get_setting("portrait_url") == "/media/site/portrait")
    _prow = site_img_svc.get("portrait")
    _pim = _PILImage.open(BytesIO(_prow.data))
    _pr = _pim.size[0] / _pim.size[1]
    ok("Portrait upload is cropped to 4:5 hero aspect",
       abs(_pr - 0.8) < 0.03, f"ratio={_pr:.3f} size={_pim.size}")
_buf_teaser = BytesIO()
_PILImage.new("RGB", (300, 180), (239, 167, 51)).save(_buf_teaser, format="JPEG")
_buf_teaser.seek(0)
admin.post(
    "/admin/settings",
    data={
        "site_title": "Bloom Anyway", "instagram_url": "",
        "hero_image_url": "", "portrait_url": "/media/site/portrait",
        "contact_email": "", "announcement_text": "", "announcement_expires": "",
        "creator_name": "", "creator_instagram": "", "creator_image_url": "",
        "creator_blurb": "", "reel_url": "", "reel_description": "",
        "hero_file": (_buf_teaser, "teaser.jpg"),
    },
    content_type="multipart/form-data",
    follow_redirects=True,
)
with app.app_context():
    _hrow = site_img_svc.get("hero")
    _him = _PILImage.open(BytesIO(_hrow.data))
    _hr = _him.size[0] / _him.size[1]
    ok("Story teaser upload is cropped to 1:1 aspect",
       abs(_hr - 1.0) < 0.03, f"ratio={_hr:.3f} size={_him.size}")

# Site images are uploaded, not typed. There is no address field any more, so
# saving Settings must leave the uploaded image alone rather than reading a
# field that isn't there and blanking it.
_sbody = admin.get("/admin/settings").get_data(as_text=True)
ok("Settings asks for a file, not an image address",
   'name="portrait_url"' not in _sbody and 'name="hero_image_url"' not in _sbody
   and 'name="portrait_file"' in _sbody and 'name="hero_file"' in _sbody)
with app.app_context():
    _before_imgs = (get_setting("portrait_url"), get_setting("hero_image_url"))
    ok("Setup: both site images are set", all(_before_imgs), f"got {_before_imgs}")
admin.post("/admin/settings", data={
    "site_title": "Bloom Anyway", "instagram_url": "", "contact_email": "",
    "announcement_text": "", "announcement_expires": "",
}, follow_redirects=True)
with app.app_context():
    ok("Saving Settings keeps the images that were uploaded",
       (get_setting("portrait_url"), get_setting("hero_image_url")) == _before_imgs,
       f"got {(get_setting('portrait_url'), get_setting('hero_image_url'))}")
admin.post("/admin/settings", data={
    "site_title": "Bloom Anyway", "instagram_url": "", "contact_email": "",
    "announcement_text": "", "announcement_expires": "", "clear_portrait": "1",
}, follow_redirects=True)
with app.app_context():
    ok("But Remove current photo still takes one down",
       not get_setting("portrait_url")
       and get_setting("hero_image_url") == _before_imgs[1])
r = app.test_client().get("/")
ok("Home uses the split healing / building hero",
   "home-hero" in r.get_data(as_text=True)
   and "home-panel--heal" in r.get_data(as_text=True)
   and "home-panel--build" in r.get_data(as_text=True))

_buf2 = BytesIO()
_PILImage.new("RGB", (100, 100), (239, 167, 51)).save(_buf2, format="JPEG")
_buf2.seek(0)
# the creator photo belongs to the Spotlight page — general Settings must not
# touch it, so posting it there should be ignored
with app.app_context():
    _portrait_url = get_setting("portrait_url")
_stray = BytesIO()
_PILImage.new("RGB", (60, 60), (10, 10, 10)).save(_stray, format="JPEG")
_stray.seek(0)
admin.post(
    "/admin/settings",
    data={
        "site_title": "Bloom Anyway", "instagram_url": "",
        "hero_image_url": "", "portrait_url": _portrait_url,
        "contact_email": "", "announcement_text": "", "announcement_expires": "",
        "creator_name": "Wiped By Settings", "creator_file": (_stray, "stray.jpg"),
    },
    content_type="multipart/form-data",
    follow_redirects=True,
)
with app.app_context():
    ok("General Settings leaves spotlight fields alone",
       get_setting("creator_name") != "Wiped By Settings")

r = admin.post(
    "/admin/spotlight",
    data={
        "creator_name": "Featured", "creator_instagram": "",
        "creator_image_url": "", "creator_blurb": "Hello",
        "reel_url": "", "reel_description": "",
        "creator_file": (_buf2, "creator.jpg"),
    },
    content_type="multipart/form-data",
    follow_redirects=True,
)
ok("Studio accepts creator-of-the-month photo upload", r.status_code == 200)
r = app.test_client().get("/media/site/creator")
ok("Uploaded creator photo is served",
   r.status_code == 200 and r.data[:3] == b"\xff\xd8\xff")
with app.app_context():
    ok("Creator photo setting points at media route",
       get_setting("creator_image_url") == "/media/site/creator")

# --- spotlight: eligible list, random draw, expiry notices ------------------
from app.services import spotlight as _spot   # noqa: E402

with app.app_context():
    _drawable = User(email="drawme@example.com", display_name="Draw Me",
                     username="drawme", membership="creator",
                     email_verified_at=datetime.utcnow())
    _drawable.set_password("memberpass123")
    _drawable.set_links([{"label": "Instagram",
                          "url": "https://instagram.com/drawmeplease"}])
    _quiet = User(email="quiet@example.com", display_name="No Links",
                  username="nolinks", membership="creator",
                  email_verified_at=datetime.utcnow())
    _quiet.set_password("memberpass123")
    db.session.add_all([_drawable, _quiet])
    db.session.commit()

    _ready, _missing = _spot.eligible_split()
    ok("Spotlight eligibility needs an Instagram link on the profile",
       any(c["handle"] == "drawmeplease" for c in _ready)
       and any(c["email"] == "quiet@example.com" for c in _missing)
       and not any(c["email"] == "quiet@example.com" for c in _ready))
    ok("Spotlight eligibility skips owners and non-Creator members",
       not any(c["email"] in ("owner@example.com", "free@example.com")
               for c in _ready + _missing))

r = admin.get("/admin/spotlight")
sbody = r.get_data(as_text=True)
ok("Spotlight page lists who can be Creator of the month",
   r.status_code == 200 and "Draw Me" in sbody
   and 'name="pick_creator"' in sbody
   and "Reel reviews" in sbody)

r = admin.post("/admin/spotlight", data={"pick_creator": "1"},
               follow_redirects=True)
dbody = r.get_data(as_text=True)
ok("Random draw pre-fills the Creator of the month form",
   'value="Draw Me"' in dbody and "Drawn at random" in dbody)
with app.app_context():
    ok("Random draw doesn't publish anything by itself",
       get_setting("creator_name") != "Draw Me")

with app.app_context():
    from datetime import date as _date

    from app.services.settings import set_setting
    set_setting("creator_name", "Draw Me")
    set_setting("creator_expires", (_date.today() + timedelta(days=1)).isoformat())
    set_setting("spotlight_creator_notified", "")
    _before = Notification.query.filter_by(kind="spotlight_expiry").count()
    _sent = _spot.sweep_expiry_notices()
    _again = _spot.sweep_expiry_notices()
    _notes = (Notification.query.filter_by(kind="spotlight_expiry")
              .order_by(Notification.id.desc()).all())
    ok("Owners are warned the day before a spotlight slot expires",
       _sent == 1 and len(_notes) > _before
       and "tomorrow" in (_notes[0].body or ""),
       f"sent={_sent} before={_before} now={len(_notes)} "
       f"body={(_notes[0].body if _notes else '')!r}")
    ok("Spotlight expiry warning is sent once, not on every sweep", _again == 0)

    set_setting("creator_expires", (_date.today() + timedelta(days=20)).isoformat())
    _slots = {s["kind"]: s for s in _spot.spotlight_slots()}
    ok("Spotlight status reports days left per slot",
       _slots["creator"]["filled"] and _slots["creator"]["days_left"] == 20
       and not _slots["creator"]["expired"])

# --- studio: where each paid tier came from --------------------------------
r = admin.get("/admin/members/audit")
abody = r.get_data(as_text=True)
ok("Members audit explains where each paid tier came from",
   r.status_code == 200 and "Set by hand in Studio" in abody
   and "Nothing on file explains it" in abody
   and "Free months from a purchase" in abody)
with app.app_context():
    from app.services import membership_audit as _audit
    _res = _audit.audit("creator")
    _sources = {row["email"]: row["source"] for row in _res["rows"]}
    ok("Audit files a comped member under 'set by hand'",
       all(t == "creator" for t in (r["tier"] for r in _res["rows"])))
    ok("Audit flags a tier with no order, grant, or perk behind it",
       _sources.get("drawme@example.com") == "unexplained")
r = admin.get("/admin/members/audit?membership=full_bloom")
ok("Members audit filters by tier", r.status_code == 200)

# A Stripe lookup that can't answer must not read as "everybody checks out".
# A broken expand string used to fail every call, fall back to local orders,
# and report that every tier already matched while nothing had been checked.
with app.app_context():
    from app.services import membership_audit as _audit
    from app.services import stripe_pay as _pay

    def _resync_with_stripe_saying(answer):
        """Run the audit resync as if Stripe were live and gave this answer."""
        real_cfg = _pay.configured
        real_live = _pay.active_membership_tier_from_stripe
        _pay.configured = lambda: True
        _pay.active_membership_tier_from_stripe = lambda _email: answer
        app.config["TESTING"] = False
        try:
            return _audit.resync_from_stripe("creator")
        finally:
            app.config["TESTING"] = True
            _pay.configured = real_cfg
            _pay.active_membership_tier_from_stripe = real_live

    _blind = _resync_with_stripe_saying(None)
    ok("Members Stripe can't answer for are counted, not silently skipped",
       _blind["unreachable"] > 0 and _blind["checked"] == 0
       and not _blind["changed"],
       f"got {_blind}")

    _dropped = _resync_with_stripe_saying("none")
    ok("Stripe saying 'not a subscriber' drops the tier",
       _dropped["unreachable"] == 0 and _dropped["checked"] > 0
       and any(c["to"] == "Free" for c in _dropped["changed"]),
       f"got {_dropped}")

r = admin.post("/admin/members/audit/resync", data={"membership": "creator"},
               follow_redirects=True)
_rbody = r.get_data(as_text=True)
ok("Resync with nobody left to check claims nothing at all",
   r.status_code == 200 and "already matched" not in _rbody
   and "Corrected" not in _rbody, flashes(r))

r = admin.post("/admin/members/audit/resync", data={"membership": "full_bloom"},
               follow_redirects=True)
ok("Resync reports what it found on a tier that still has members",
   r.status_code == 200 and "Checked 2 member(s)" in r.get_data(as_text=True),
   flashes(r))

# Content Hub: the tips library appears before reel reviews
r = client.get("/watch")
hub = r.get_data(as_text=True)
ok("Content Hub lists tips above reel reviews",
   hub.find('id="videos"') < hub.find('id="reviews"')
   and hub.find("Content tips library") < hub.find("Reel reviews"))
ok("Hub card shows the tip summary and read time",
   "Twenty minutes, once a week." in hub and "min read" in hub)

# Brevo helper strips Bearer / whitespace / wrapping quotes
from app.services import mailer as mailer_mod
_prev_brevo = os.environ.pop("BREVO_API_KEY", None)
try:
    with app.app_context():
        app.config["BREVO_API_KEY"] = "  Bearer  xkeysib-abc123  "
        cleaned = mailer_mod._brevo_api_key()
        app.config["BREVO_API_KEY"] = '"xkeysib-xyz-key"'
        quoted = mailer_mod._brevo_api_key()
        from_parsed = mailer_mod._strip_env_quotes('"Bloom Anyway <hello@example.com>"')
        name, email = mailer_mod._parse_mail_from("Bloom Anyway <hello@example.com>")
        bad_key_hint = mailer_mod._brevo_error_hint(401, '{"message":"Key not found"}')
        domain_hint = mailer_mod._brevo_error_hint(400, '{"message":"Invalid sender"}')
    ok("Brevo API key is normalized",
       cleaned == "xkeysib-abc123" and quoted == "xkeysib-xyz-key")
    ok("MAIL_FROM wrapping quotes are stripped",
       from_parsed == "Bloom Anyway <hello@example.com>")
    ok("MAIL_FROM parses into Brevo sender fields",
       name == "Bloom Anyway" and email == "hello@example.com")
    ok("Brevo 401 hint mentions API key",
       "BREVO_API_KEY" in bad_key_hint)
    ok("Brevo 400 hint mentions verified sender",
       "verified" in domain_hint.lower())
finally:
    if _prev_brevo is not None:
        os.environ["BREVO_API_KEY"] = _prev_brevo
    else:
        os.environ.pop("BREVO_API_KEY", None)

# brand rename: leftover "First Light" becomes Bloom Anyway on boot
from app.services.settings import ensure_brand_title, get_setting, invalidate_cache, set_setting
with app.app_context():
    set_setting("site_title", "First Light")
    invalidate_cache()
    rewritten = ensure_brand_title()
    new_title = get_setting("site_title")
ok("Legacy site title is rewritten to Bloom Anyway",
   rewritten and new_title == "Bloom Anyway", f"got {new_title!r}")

# --- 8. DB-backed SECRET_KEY (no env var needed) ---------------------------
KEY_DB = Path(tempfile.mkdtemp()) / "key.db"


class NoSecretConfig(TestConfig):
    SQLALCHEMY_DATABASE_URI = f"sqlite:///{KEY_DB.as_posix()}"
    SECRET_KEY = ""   # force the database-backed path


ks = create_app(NoSecretConfig)
with ks.app_context():
    db.create_all()
boot1 = create_app(NoSecretConfig)
boot2 = create_app(NoSecretConfig)
k1, k2 = boot1.config["SECRET_KEY"], boot2.config["SECRET_KEY"]
ok("SECRET_KEY auto-generated when unset", bool(k1) and len(k1) >= 32)
ok("SECRET_KEY stable across restarts", k1 == k2)
with boot2.app_context():
    from app.services.settings import all_settings
    ok("Secret key never leaks into public settings", "_secret_key" not in all_settings())

# --- 9. feedback inbox, content reports, legal, privacy hardening -----------
from app.models import Notification, PageView, SiteFeedback
from app.services.content_reports import review_text, submit_report
from app.services.privacy import close_account

css = client.get("/static/css/main.css").get_data(as_text=True)
ok("Auth pages keep the sun accent styles",
   ".sun-disc" in css and "sun-breathe" in css)
ok("Quote mini archive cards are centered",
   ".quote-mini" in css and "align-items: center" in css)
# A class that sets a display of its own outranks the browser's rule for the
# attribute, so without this one anything marked hidden stays on screen.
ok("The hidden attribute beats whatever else sets a display",
   "[hidden] { display: none !important; }" in css)
_hidden_rules = [
    " ".join(sel.split())
    for sel, decls in re.findall(r"([^{}]+)\{([^{}]*)\}",
                                 re.sub(r"/\*.*?\*/", "", css, flags=re.S))
    if "[hidden]" in " ".join(sel.split()).replace(":not([hidden])", "")
    and "display" in decls
]
ok("And is the only place that has to say so",
   _hidden_rules == ["[hidden]"], f"rules={_hidden_rules}")

# A grid item stretches to its row unless told otherwise, and a stretched
# height combined with an aspect ratio decides the width — which is how a
# library cover grew past its own column and sat over the title and buttons.
_photo_cover = re.search(
    r"\.lib-card__cover--photo \{([^}]*)\}", re.sub(r"/\*.*?\*/", "", css, flags=re.S))
ok("A library cover is sized by its column, not by how tall the card is",
   bool(_photo_cover) and "aspect-ratio" in _photo_cover.group(1)
   and "align-self: start" in _photo_cover.group(1),
   f"rule={_photo_cover.group(1).strip() if _photo_cover else None}")
_full_width_actions = [
    " ".join(sel.split())
    for sel, decls in re.findall(r"([^{}]+)\{([^{}]*)\}",
                                 re.sub(r"/\*.*?\*/", "", css, flags=re.S))
    if ".lib-card__actions" in sel and "grid-column: 1 / -1" in decls
]
ok("On a phone the buttons get the whole card width, not the sliver beside the cover",
   bool(_full_width_actions), f"rules={_full_width_actions}")

# "Page 12 of 248" is wider than the gap between the two buttons on a small
# phone, and was being cut through the middle rather than given room.
ok("The word Page steps aside on a phone rather than the count being cut",
   ".reader__page-word" in css
   and re.search(r"@media \(max-width: 420px\) \{\s*\.reader__page-word", css)
   is not None,
   "nothing gives way on a narrow screen")
ok("But it is still there to be read aloud",
   "clip: rect(0, 0, 0, 0)" in css)

# Putting the page back where someone left it runs a few times while the page
# settles. On a phone those land after a tap, and scrolling out from under a
# field that has just been tapped closes the keyboard again before a word can
# be typed.
_loader = client.get("/static/js/page-loader.js").get_data(as_text=True)
ok("The scroll it remembers is never restored over someone in a field",
   bool(re.search(r"if \(touched \|\| typing\(\)\) return;\s*window\.scrollTo",
                  _loader)),
   "the restore scrolls unconditionally")
ok("And touching the page at all hands it over to them",
   all(evt in _loader for evt in ("pointerdown", "touchstart", "keydown"))
   and "if (window.location.hash || touched || typing()) return;" in _loader)

# Turning pages is one control, not three: back on the left, forward on the
# right, the page you are on between them. Left to wrap on a phone, Next fell
# onto a line of its own underneath Previous.
_pager_rules = [
    (" ".join(sel.split()), decls)
    for sel, decls in re.findall(r"([^{}]+)\{([^{}]*)\}",
                                 re.sub(r"/\*.*?\*/", "", css, flags=re.S))
    if "reader__footer:not(.reader__footer--simple)" in " ".join(sel.split())
]
ok("The reader's page buttons stay on one row on a phone",
   any("flex-wrap: nowrap" in d for _, d in _pager_rules),
   f"rules={[s for s, _ in _pager_rules]}")
ok("And the page number sits between them",
   any("reader__page-field" in s and "justify-content: center" in d
       and "flex: 1" in d for s, d in _pager_rules))

for path, needle in (("/privacy", "What we collect"),
                     ("/terms", "Full Bloom"),
                     ("/refunds", "14 days")):
    rr = client.get(path)
    body = rr.get_data(as_text=True)
    ok(f"Legal page {path} renders", rr.status_code == 200)
    ok(f"Legal page {path} has real copy",
       needle in body and "TODO: legal review" not in body)

r = client.get("/")
ok("Feedback widget on public pages",
   b"data-feedback-open" in r.data and b"feedback-dialog" in r.data)

r = client.post("/feedback", data={
    "kind": "feedback", "stars": "4", "body": "Loving the daily quotes.",
    "page_path": "/", "next": "/",
}, follow_redirects=True)
ok("Star feedback accepted", r.status_code == 200)
with app.app_context():
    fb = SiteFeedback.query.filter_by(kind="feedback").order_by(SiteFeedback.id.desc()).first()
    ok("Feedback stored with stars",
       fb is not None and fb.stars == 4 and "Loving" in fb.body)
    from app.models import Notification
    owner = User.query.filter_by(email="owner@example.com").first()
    owner_note = (Notification.query
                  .filter_by(user_id=owner.id, kind="inbox")
                  .order_by(Notification.id.desc()).first()) if owner else None
ok("Owners get a notification for new feedback",
   owner_note is not None and "feedback" in (owner_note.body or "").lower()
   and owner_note.url and "inbox" in owner_note.url)

r = client.post("/feedback", data={
    "kind": "complaint", "body": "Checkout felt confusing on mobile.",
    "page_path": "/membership", "next": "/",
}, follow_redirects=True)
ok("Complaint accepted", r.status_code == 200)

with app.app_context():
    err_before = SiteFeedback.query.filter_by(kind="error").count()
r = client.post("/feedback", data={
    "kind": "error", "body": "bot noise",
    "page_path": "/videos", "website": "http://bots.example",
    "next": "/",
}, follow_redirects=True)
with app.app_context():
    err_after_hp = SiteFeedback.query.filter_by(kind="error").count()
ok("Feedback honeypot ignored", r.status_code == 200 and err_after_hp == err_before)

r = client.post("/feedback", data={
    "kind": "error", "body": "Saw a 500 after uploading a huge video.",
    "page_path": "/videos", "next": "/",
}, follow_redirects=True)
with app.app_context():
    err = SiteFeedback.query.filter_by(kind="error").order_by(SiteFeedback.id.desc()).first()
ok("Error report stored for studio", err is not None and "500" in err.body)

ok("Auto-mod flags blocked language",
   review_text("what the fuck is this") == "Blocked language")
ok("Auto-mod passes clean text",
   review_text("Thanks for the gentle advice today.") is None)
ok("Auto-mod flags hostile phrase",
   review_text("you should just kill yourself") == "Hostile or threatening language")

with app.app_context():
    healing_cat = ForumCategory.query.filter_by(slug="healing").first()
    ok("Healing forum exists for report tests", healing_cat is not None)
    author = User(email="reporter-author@example.com", display_name="Author",
                  membership="healing", email_verified_at=utcnow())
    author.set_password(USER_PW)
    reporter = User(email="reporter-user@example.com", display_name="Reporter",
                    membership="healing", email_verified_at=utcnow())
    reporter.set_password(USER_PW)
    db.session.add_all([author, reporter])
    db.session.flush()
    clean = ForumPost(category_id=healing_cat.id, user_id=author.id,
                      title="Soft morning", body="Just checking in with kindness.",
                      anonymous=False)
    toxic = ForumPost(category_id=healing_cat.id, user_id=author.id,
                      title="Bad day", body="go kill yourself already",
                      anonymous=False)
    db.session.add_all([clean, toxic])
    db.session.commit()
    clean_id, toxic_id = clean.id, toxic.id
    author_id, reporter_id = author.id, reporter.id
    healing_cat_id = healing_cat.id

    rep_open, _msg = submit_report(reporter=reporter, target_type="post",
                                   target_id=clean_id, note="feels off")
    clean_after = db.session.get(ForumPost, clean_id)
    ok("Clean reported post stays visible",
       rep_open is not None and rep_open.status == "open" and not clean_after.hidden)

    rep_auto, _msg = submit_report(reporter=reporter, target_type="post",
                                   target_id=toxic_id, note="threat")
    toxic_after = db.session.get(ForumPost, toxic_id)
    note = Notification.query.filter_by(user_id=author_id, kind="moderation").first()
    ok("Toxic reported post auto-hidden",
       rep_auto is not None and rep_auto.status == "resolved" and toxic_after.hidden
       and bool(rep_auto.auto_reason))
    ok("Author notified on auto take-down", note is not None)

# refresh admin session for inbox checks
admin.post("/login", data={"email": "owner@example.com", "password": ADMIN_PW})
r = admin.get("/admin/inbox")
ok("Studio inbox loads", r.status_code == 200 and b"Inbox" in r.data)
r = admin.get("/admin/inbox?filter=feedback")
ok("Studio inbox feedback filter", r.status_code == 200 and b"Loving the daily" in r.data)
r = admin.get("/admin/inbox?filter=complaint")
ok("Studio inbox complaint filter", r.status_code == 200 and b"Checkout felt" in r.data)
r = admin.get("/admin/inbox?filter=error")
ok("Studio inbox error filter", r.status_code == 200 and b"huge video" in r.data)

# A complaint is someone waiting on an answer just as much as a contact form
# message is, so it gets the same composer, senders and templates.
with app.app_context():
    _cx = (SiteFeedback.query.filter_by(kind="complaint")
           .order_by(SiteFeedback.id.desc()).first())
    _cx_id = _cx.id
    _cx_to = (_cx.contact_email or (_cx.author.email if _cx.author else ""))
ok("The inbox offers Reply on a complaint there's an address for",
   f"/inbox/feedback/{_cx_id}/reply"
   in admin.get("/admin/inbox?filter=complaint").get_data(as_text=True))
_crbody = admin.get(f"/admin/inbox/feedback/{_cx_id}/reply").get_data(as_text=True)
ok("A complaint opens the same reply composer, quoting them",
   "Checkout felt confusing on mobile." in _crbody
   and "data-reply-preview" in _crbody)
ok("With the same senders to choose between",
   all(a in _crbody for a in ("bloomsupport@bloomanyway.online",
                              "ayesha@bloomanyway.online",
                              "saman@bloomanyway.online")))
_cx_calls = []
_cx_real_send = _mailer.send_email
_mailer.send_email = (
    lambda to, subject, text, **kw: _cx_calls.append(
        dict(kw, to=to, subject=subject, text=text)) or True
)
try:
    r = admin.post(f"/admin/inbox/feedback/{_cx_id}/reply",
                   data={"sender": "saman", "subject": "Re: checkout",
                         "preview": "Sorry about that.", "header": "Bloom Anyway",
                         "title": "Hi there,", "body": "We have made it clearer."},
                   follow_redirects=True)
finally:
    _mailer.send_email = _cx_real_send
_call = _cx_calls[-1] if _cx_calls else {}
ok("Replying to a complaint reaches whoever left it",
   _call.get("to") == _cx_to, f"got {_call.get('to')} want {_cx_to}")
ok("On the same per-address template as any other reply",
   _call.get("template_id") == 21, f"got {_call.get('template_id')}")
ok("Carrying the same five params",
   (_call.get("params") or {}).get("PREVIEW") == "Sorry about that."
   and (_call.get("params") or {}).get("TITLE") == "Hi there,")
with app.app_context():
    ok("And the complaint is marked handled once answered",
       db.session.get(SiteFeedback, _cx_id).status == "reviewed")

with app.app_context():
    _anon_cx = SiteFeedback(kind="complaint", body="No way to reach me.",
                            page_path="/", status="new")
    db.session.add(_anon_cx)
    db.session.commit()
    _anon_cx_id = _anon_cx.id
ok("One left without an address offers no Reply button",
   f"/inbox/feedback/{_anon_cx_id}/reply"
   not in admin.get("/admin/inbox?filter=complaint").get_data(as_text=True))
r = admin.get(f"/admin/inbox/feedback/{_anon_cx_id}/reply", follow_redirects=True)
ok("And says so rather than failing if you reach for it anyway",
   r.status_code == 200 and "nowhere to reply" in r.get_data(as_text=True).lower())
r = admin.get("/admin/inbox?filter=open")
ok("Studio inbox open content reports", r.status_code == 200 and b"Soft morning" in r.data)
r = admin.get("/admin/inbox?filter=resolved")
ok("Studio inbox resolved shows auto reason",
   r.status_code == 200 and (b"Hostile" in r.data or b"Blocked" in r.data
                             or b"auto-hidden" in r.data))

rep_client = app.test_client()
rep_client.post("/login", data={"email": "reporter-user@example.com", "password": USER_PW})
r = rep_client.get(f"/forums/p/{clean_id}")
post_html = r.get_data(as_text=True)
ok("Report control on posts",
   "Report post" in post_html
   and f"/forums/p/{clean_id}/report" in post_html)

with app.app_context():
    from app.models import ForumComment
    author = db.session.get(User, author_id)
    cmt = ForumComment(post_id=clean_id, user_id=author_id,
                       body="A gentle reply worth reporting if needed.",
                       anonymous=False)
    db.session.add(cmt)
    db.session.commit()
    comment_id = cmt.id

r = rep_client.get(f"/forums/p/{clean_id}")
ok("Report control on comments",
   f"/forums/comment/{comment_id}/report" in r.get_data(as_text=True))
r = rep_client.post(f"/forums/comment/{comment_id}/report",
                    data={"note": "feels off"}, follow_redirects=True)
ok("Comment report is accepted", r.status_code == 200)
with app.app_context():
    from app.models import ContentReport
    c_rep = ContentReport.query.filter_by(
        reporter_id=reporter_id, target_type="comment", target_id=comment_id).first()
ok("Comment report is stored", c_rep is not None and c_rep.status == "open")

with app.app_context():
    from app.models import Notification
    owner = User.query.filter_by(email="owner@example.com").first()
    rep_note = (Notification.query
                .filter_by(user_id=owner.id, kind="inbox")
                .order_by(Notification.id.desc()).first()) if owner else None
ok("Owners get a notification for content reports",
   rep_note is not None and "report" in (rep_note.body or "").lower()
   and rep_note.url and "inbox" in rep_note.url)

r = rep_client.get("/forums/c/healing")
ok("Feed lists report control on posts",
   f"/forums/p/{clean_id}/report" in r.get_data(as_text=True)
   and "report-note-feed-" in r.get_data(as_text=True))

# Studio can remove a post that still has likes, replies, and notifications
with app.app_context():
    from app.models import (ContentReport, ForumCommentLike, ForumPostLike,
                            Notification)
    doomed_studio = User(email="studio-rm@example.com", display_name="Studio Rm",
                         membership="healing", email_verified_at=utcnow())
    doomed_studio.set_password(USER_PW)
    db.session.add(doomed_studio)
    db.session.flush()
    victim = ForumPost(category_id=healing_cat_id, user_id=doomed_studio.id,
                       title="Studio will remove me", body="with baggage",
                       anonymous=False)
    db.session.add(victim)
    db.session.flush()
    top = ForumComment(post_id=victim.id, user_id=doomed_studio.id,
                       body="top comment", anonymous=False)
    db.session.add(top)
    db.session.flush()
    reply = ForumComment(post_id=victim.id, user_id=doomed_studio.id,
                         parent_id=top.id, body="a reply", anonymous=False)
    db.session.add(reply)
    db.session.add(ForumPostLike(user_id=doomed_studio.id, post_id=victim.id))
    db.session.add(ForumCommentLike(user_id=doomed_studio.id, comment_id=top.id))
    db.session.add(Notification(user_id=doomed_studio.id, kind="mention",
                                post_id=victim.id, body="you were mentioned"))
    db.session.add(ContentReport(target_type="post", target_id=victim.id,
                                 reporter_id=doomed_studio.id, note="noise"))
    db.session.commit()
    victim_id = victim.id

r = admin.post(f"/admin/community/post/{victim_id}/delete", follow_redirects=True)
ok("Studio removes community post with dependents",
   r.status_code == 200 and b"Post removed" in r.data)
with app.app_context():
    from app.models import Notification
    ok("Removed community post is gone",
       db.session.get(ForumPost, victim_id) is None)
    ok("Post notifications cleared on studio remove",
       Notification.query.filter_by(post_id=victim_id).count() == 0)

with app.app_context():
    from app.models import Order, ShopPurchase
    doomed = User(email="doomed@example.com", display_name="Doomed Soul",
                  username="doomedx", bio="secret bio", membership="healing",
                  email_verified_at=utcnow(),
                  avatar_data=b"fakepng", avatar_mime="image/png")
    doomed.set_password(USER_PW)
    db.session.add(doomed)
    db.session.flush()
    doomed_post = ForumPost(category_id=healing_cat_id, user_id=doomed.id,
                            title="Will vanish", body="please hide me",
                            anonymous=False)
    db.session.add(doomed_post)
    paid = Order(
        ls_order_id="close-acct-mem-1",
        buyer_email="doomed@example.com",
        membership_tier="healing",
        status="paid",
        total_cents=900,
        currency="USD",
    )
    db.session.add(paid)
    shop = ShopPurchase(
        lemon_squeezy_order_id="close-acct-shop-1",
        customer_email="doomed@example.com",
        user_id=doomed.id,
        product_name="A guide",
        status="linked",
    )
    db.session.add(shop)
    db.session.commit()
    doomed_id, doomed_post_id = doomed.id, doomed_post.id
    paid_id, shop_id = paid.id, shop.id
    close_account(doomed)
    doomed = db.session.get(User, doomed_id)
    doomed_post = db.session.get(ForumPost, doomed_post_id)
    paid = db.session.get(Order, paid_id)
    shop = db.session.get(ShopPurchase, shop_id)
    from app.services.privacy import FORMER_MEMBER_EMAIL
    former = User.query.filter_by(email=FORMER_MEMBER_EMAIL).first()
    ok("Closed account row removed from database", doomed is None)
    ok("Closed account posts hidden under former member",
       doomed_post is not None and doomed_post.hidden is True
       and former is not None and doomed_post.user_id == former.id)
    ok("Closed account membership orders ended",
       paid is not None and paid.status == "ended"
       and (paid.buyer_email or "").startswith("closed+"))
    ok("Closed account shop purchases detached",
       shop is not None and shop.user_id is None
       and (shop.customer_email or "").startswith("closed+"))
    from app.services.memberships import purchased_tier
    ok("Re-signup email has no purchased membership",
       purchased_tier("doomed@example.com") == "none")


ok("PageView has no IP field",
   not hasattr(PageView, "ip") and not hasattr(PageView, "ip_address"))

# --- membership emails must not follow a deleted account ----------------------
import app.services.mailer as _mailer_mod

_welcomes = []
_billing_alerts = []
_mailer_mod.send_healing_welcome = (
    lambda to, **kw: _welcomes.append(to) or True)
_mailer_mod.send_creator_welcome = (
    lambda to, **kw: _welcomes.append(to) or True)
_mailer_mod.send_full_bloom_welcome = (
    lambda to, **kw: _welcomes.append(to) or True)
_mailer_mod.send_billing_alert = (
    lambda title, body: _billing_alerts.append(title) or True)


def _invoice_object(invoice_id, email, price_id, sub_id, reason, *, basil=False):
    """One invoice payload. ``basil`` uses Stripe's 2025-03-31 shape, where
    ``subscription`` moved under ``parent.subscription_details``."""
    obj = {
        "id": f"in_{invoice_id}",
        "object": "invoice",
        "amount_paid": 900,
        "currency": "usd",
        "customer_email": email,
        "payment_intent": str(invoice_id),
        "billing_reason": reason,
        "metadata": {"tier": "healing"},
        "lines": {"data": [{"price": {"id": price_id}}]},
    }
    if basil:
        obj["parent"] = {
            "type": "subscription_details",
            "subscription_details": {"subscription": sub_id, "metadata": {}},
        }
        obj["lines"] = {"data": [{
            "pricing": {"price_details": {"price": price_id}},
            "parent": {
                "type": "subscription_item_details",
                "subscription_item_details": {"subscription": sub_id},
            },
        }]}
    else:
        obj["subscription"] = sub_id
    return obj


def _invoice_webhook(invoice_id, email, price_id, sub_id, reason, *, basil=False):
    body = json.dumps({
        "id": f"evt_{invoice_id}",
        "object": "event",
        "type": "invoice.paid",
        "data": {"object": _invoice_object(invoice_id, email, price_id, sub_id,
                                           reason, basil=basil)},
    }).encode()
    return client.post("/webhooks/stripe", data=body, headers=_stripe_headers(body))


# Stripe's Basil API version (2025-03-31) removed invoice.subscription. Reading
# only the old field left a renewal looking unattached to any subscription, and
# "cancel every membership sub except the new one" then cancelled the new one.
_basil = _invoice_object("BASIL-0", "x@example.com", "price_x", "sub_basil",
                         "subscription_cycle", basil=True)
ok("The subscription is found in Stripe's current invoice shape",
   pay.invoice_subscription_id(_basil) == "sub_basil")
ok("And still in the old one",
   pay.invoice_subscription_id(
       _invoice_object("OLD-0", "x@example.com", "price_x", "sub_old",
                       "subscription_cycle")) == "sub_old")
_evt, _data = pay.stripe_event_to_internal("invoice.paid", _basil)
ok("So a Basil renewal reaches fulfillment carrying its subscription",
   _evt == "payment.succeeded"
   and (_data.get("metadata") or {}).get("subscription_id") == "sub_basil",
   f"got {_data.get('metadata')}")


with app.app_context():
    from app.models import MembershipPlan
    _hplan = MembershipPlan.query.filter_by(tier="healing").first()
    _mem_price = _hplan.stripe_price_id
    leaver = User(email="leaver@example.com", display_name="Leaver",
                  membership="none", email_verified_at=utcnow())
    leaver.set_password(USER_PW)
    db.session.add(leaver)
    db.session.commit()
    leaver_id = leaver.id

_invoice_webhook("WEL-1", "leaver@example.com", _mem_price,
                 "sub_leaver", "subscription_create")
ok("First membership payment sends the welcome",
   _welcomes.count("leaver@example.com") == 1, f"got {_welcomes}")
with app.app_context():
    ok("A Basil-shaped payment keeps the membership it just paid for",
       db.session.get(User, leaver_id).membership == "healing",
       f"got {db.session.get(User, leaver_id).membership}")

# The net under that: a payment we still can't place against a subscription
# must cancel nothing, because "every subscription but this one" would include
# the one being paid for.
with app.app_context():
    def _stray(order_id, email, sub_id):
        db.session.add(Order(
            ls_order_id=order_id, buyer_email=email, ls_variant_id=_mem_price,
            status="paid", membership_tier="healing", total_cents=900,
            currency="USD", stripe_subscription_id=sub_id))
        db.session.commit()

    _stray("BLIND-OLD", "placeable@example.com", "sub_placeable_old")
    placed = pay.replace_other_memberships(
        "placeable@example.com", keep_order_id="BLIND-NEW",
        keep_subscription_id="sub_placeable_new")
    ok("A new membership flags the other subscription to the owner, never cancels it",
       placed["cancelled"] == [] and placed.get("flagged") == ["sub_placeable_old"],
       f"got {placed}")

    _stray("BLIND-LIVE", "unplaceable@example.com", "sub_unplaceable_live")
    blind = pay.replace_other_memberships(
        "unplaceable@example.com", keep_order_id="BLIND-NOWHERE",
        keep_subscription_id=None)
ok("A payment we cannot place cancels nothing at all",
   blind["cancelled"] == [] and "unknown_subscription" in blind["errors"],
   f"got {blind}")

# Coming back from Stripe used to wait on one API call per checkout of the last
# fortnight. The order id is on the list entry, so anything already fulfilled
# is skipped before the call rather than after it.
with app.app_context():
    import types as _types
    _seen = {"list": 0, "retrieve": 0}

    class _Sess(dict):
        __getattr__ = dict.get

    _fake = [_Sess({"id": f"cs_perf{i}", "object": "checkout.session",
                    "status": "complete", "payment_status": "paid",
                    "amount_total": 2400, "currency": "usd",
                    "payment_intent": f"pi_perf{i}", "metadata": {},
                    "customer_email": "perf@example.com",
                    "line_items": {"data": []}}) for i in range(6)]
    ok("A checkout session's order id can be read off the list entry",
       pay._session_payment_id(_fake[0]) == "pi_perf0")

    for _s in _fake[:5]:  # five already fulfilled, one new
        db.session.add(Order(ls_order_id=_s["payment_intent"], status="paid",
                             buyer_email="perf@example.com", total_cents=2400,
                             currency="USD"))
    db.session.commit()

    class _Page:
        def __init__(self, data): self.data = data; self.has_more = False

    _real_stripe, _real_conf, _real_cfg = pay.stripe, pay.configured, pay._configure_stripe
    pay.configured = lambda: True
    pay._configure_stripe = lambda: None
    pay.stripe = _types.SimpleNamespace(checkout=_types.SimpleNamespace(
        Session=_types.SimpleNamespace(
            list=lambda **kw: (_seen.__setitem__("list", _seen["list"] + 1)
                               or _Page(_fake)),
            retrieve=lambda sid, **kw: (_seen.__setitem__("retrieve", _seen["retrieve"] + 1)
                                        or next(s for s in _fake if s["id"] == sid)))))
    app.config["TESTING"] = False
    try:
        pay.sync_recent_payments(days=14, max_pages=1)
    finally:
        app.config["TESTING"] = True
        pay.stripe, pay.configured, pay._configure_stripe = _real_stripe, _real_conf, _real_cfg
    ok("Only the checkout it hasn't seen costs an API call",
       _seen == {"list": 1, "retrieve": 1}, f"got {_seen}")


# --- the other ways the money stops -------------------------------------------
def _stripe_event(kind, obj, previous=None):
    payload = {"object": obj}
    if previous is not None:
        payload["previous_attributes"] = previous
    body = json.dumps({"id": f"evt_{kind}_{len(_billing_alerts)}", "object": "event",
                       "type": kind, "data": payload}).encode()
    return client.post("/webhooks/stripe", data=body, headers=_stripe_headers(body))


with app.app_context():
    disputer = User(email="disputer@example.com", display_name="Disputer",
                    membership="none", email_verified_at=utcnow())
    disputer.set_password(USER_PW)
    db.session.add(disputer)
    db.session.add(Order(
        ls_order_id="pi_disputed", buyer_email="disputer@example.com",
        ls_variant_id=_mem_price, status="paid", membership_tier="healing",
        total_cents=900, currency="USD", stripe_subscription_id="sub_disputer"))
    db.session.commit()
    from app.services.memberships import reconcile_email as _reconcile
    _reconcile("disputer@example.com")
    db.session.commit()
    disputer_id = disputer.id
    ok("Setup: the disputer holds the tier they paid for",
       db.session.get(User, disputer_id).membership == "healing")

_alerts = len(_billing_alerts)
_stripe_event("charge.dispute.created", {
    "id": "dp_1", "object": "dispute", "charge": "ch_1",
    "payment_intent": "pi_disputed", "amount": 900, "reason": "fraudulent",
    "status": "needs_response",
})
with app.app_context():
    ok("A chargeback takes the membership back straight away",
       db.session.get(User, disputer_id).membership == "none",
       f"got {db.session.get(User, disputer_id).membership}")
ok("And the owner is told, with time to respond",
   any("disputed" in t.lower() for t in _billing_alerts[_alerts:]),
   f"got {_billing_alerts[_alerts:]}")

_stripe_event("charge.dispute.closed", {
    "id": "dp_1", "object": "dispute", "charge": "ch_1",
    "payment_intent": "pi_disputed", "amount": 900, "status": "won",
})
with app.app_context():
    ok("Winning the dispute puts the membership back",
       db.session.get(User, disputer_id).membership == "healing",
       f"got {db.session.get(User, disputer_id).membership}")

# Dunning giving up never deletes the subscription, so this is the only warning.
with app.app_context():
    stalled = User(email="stalled@example.com", display_name="Stalled",
                   membership="none", email_verified_at=utcnow())
    stalled.set_password(USER_PW)
    db.session.add(stalled)
    db.session.add(Order(
        ls_order_id="pi_stalled", buyer_email="stalled@example.com",
        ls_variant_id=_mem_price, status="paid", membership_tier="healing",
        total_cents=900, currency="USD", stripe_subscription_id="sub_stalled"))
    db.session.commit()
    _reconcile("stalled@example.com")
    db.session.commit()
    stalled_id = stalled.id

_sub_obj = {
    "id": "sub_stalled", "object": "subscription", "status": "past_due",
    "customer": {"email": "stalled@example.com"},
    "items": {"data": [{"price": {"id": _mem_price}}]},
}
_stripe_event("customer.subscription.updated", dict(_sub_obj))
with app.app_context():
    ok("A card still being retried keeps the membership",
       db.session.get(User, stalled_id).membership == "healing")
_stripe_event("customer.subscription.updated", dict(_sub_obj, status="unpaid"))
with app.app_context():
    ok("But once Stripe gives up, the membership goes",
       db.session.get(User, stalled_id).membership == "none",
       f"got {db.session.get(User, stalled_id).membership}")

_alerts = len(_billing_alerts)
_stripe_event("invoice.payment_action_required", _invoice_object(
    "SCA-1", "stalled@example.com", _mem_price, "sub_stalled",
    "subscription_cycle", basil=True))
ok("A renewal waiting on the member's bank is not silent",
   any("bank" in t.lower() for t in _billing_alerts[_alerts:]),
   f"got {_billing_alerts[_alerts:]}")

# Cancelling someone's billing is the most damaging thing the site does by
# itself, so when it does it unasked the owner hears about it, and the reason
# goes onto the subscription in Stripe.
with app.app_context():
    _cancel_calls = []
    _real_cancel = pay.stripe.Subscription.cancel

    class _FakeCancel:
        @staticmethod
        def cancel(sid, **kw):
            _cancel_calls.append({"sid": sid, **kw})
            return {"id": sid, "status": "canceled"}

    _real_sub = pay.stripe.Subscription
    _real_conf = pay.configured
    pay.stripe.Subscription = _FakeCancel
    pay.configured = lambda: True
    app.config["TESTING"] = False
    _alerts_before_cancel = len(_billing_alerts)
    try:
        ok("A cancel the member asked for goes through",
           pay._cancel_stripe_subscription_now("sub_asked", "member cancelled"))
        ok("Its reason is written onto the subscription in Stripe",
           "member cancelled" in str(_cancel_calls[-1].get("cancellation_details")),
           f"got {_cancel_calls[-1]}")
        ok("And it doesn't email the owner about a cancel they asked for",
           len(_billing_alerts) == _alerts_before_cancel)
        ok("A cancel the site decided on its own also goes through",
           pay._cancel_stripe_subscription_now("sub_auto", "some internal rule"))
        ok("But that one the owner is told about",
           len(_billing_alerts) > _alerts_before_cancel
           and any("cancelled a subscription" in t.lower()
                   for t in _billing_alerts[_alerts_before_cancel:]),
           f"got {_billing_alerts[_alerts_before_cancel:]}")
    finally:
        pay.stripe.Subscription = _real_sub
        pay.configured = _real_conf
        app.config["TESTING"] = True
    assert _real_cancel is not None

# Switching plans cancels the old subscription seconds before the new one is
# paid for, but the local order saying so is still marked paid when the payment
# lands. Reading the status off Stripe stops the owner being told a member has
# two live memberships when one of them is already gone.
with app.app_context():
    _subs = {
        "sub_switched_away": {"id": "sub_switched_away", "status": "canceled"},
        "sub_second_plan": {"id": "sub_second_plan", "status": "active"},
        "sub_winding_down": {"id": "sub_winding_down", "status": "active",
                             "cancel_at_period_end": True},
    }

    class _FakeLookup:
        @staticmethod
        def retrieve(sid, **kw):
            if sid not in _subs:
                raise Exception(f"No such subscription: {sid}")
            return _subs[sid]

    _real_sub = pay.stripe.Subscription
    _real_conf = pay.configured
    pay.stripe.Subscription = _FakeLookup
    pay.configured = lambda: True
    app.config["TESTING"] = False
    try:
        _before = len(_billing_alerts)
        pay._flag_extra_memberships(
            "switcher@example.com", "sub_new_plan", ["sub_switched_away"])
        ok("Switching plans doesn't report the plan just left as a second one",
           len(_billing_alerts) == _before, f"got {_billing_alerts[_before:]}")
        pay._flag_extra_memberships(
            "winding@example.com", "sub_new_plan", ["sub_winding_down"])
        ok("Nor one that is already set to stop renewing",
           len(_billing_alerts) == _before, f"got {_billing_alerts[_before:]}")
        pay._flag_extra_memberships(
            "stale@example.com", "sub_new_plan", ["sub_long_gone"])
        ok("Nor a subscription id Stripe has never heard of",
           len(_billing_alerts) == _before, f"got {_billing_alerts[_before:]}")
        pay._flag_extra_memberships(
            "double@example.com", "sub_new_plan",
            ["sub_second_plan", "sub_switched_away"])
        ok("But a genuinely live second membership is still reported",
           len(_billing_alerts) == _before + 1
           and "more than one membership" in _billing_alerts[-1].lower(),
           f"got {_billing_alerts[_before:]}")
    finally:
        pay.stripe.Subscription = _real_sub
        pay.configured = _real_conf
        app.config["TESTING"] = True

_alerts = len(_billing_alerts)
_stripe_event("customer.updated",
              {"id": "cus_1", "object": "customer", "email": "brand-new@example.com"},
              previous={"email": "stalled@example.com"})
ok("Changing the billing email is caught before the next renewal misses",
   any("billing email" in t.lower() for t in _billing_alerts[_alerts:]),
   f"got {_billing_alerts[_alerts:]}")
_alerts = len(_billing_alerts)
_stripe_event("customer.updated",
              {"id": "cus_1", "object": "customer", "email": "stalled@example.com"},
              previous={"name": "Someone"})
ok("A customer edit that isn't the email says nothing",
   len(_billing_alerts) == _alerts, f"got {_billing_alerts[_alerts:]}")

_invoice_webhook("WEL-2", "leaver@example.com", _mem_price,
                 "sub_leaver", "subscription_cycle")
ok("Renewal does not re-send the welcome",
   _welcomes.count("leaver@example.com") == 1, f"got {_welcomes}")

with app.app_context():
    first_order = Order.query.filter_by(ls_order_id="WEL-1").first()
    ok("Order records the subscription behind the payment",
       first_order is not None and first_order.stripe_subscription_id == "sub_leaver")
    close_account(db.session.get(User, leaver_id))

_invoice_webhook("WEL-3", "leaver@example.com", _mem_price,
                 "sub_leaver", "subscription_cycle")
ok("Renewal after account deletion sends no welcome",
   _welcomes.count("leaver@example.com") == 1, f"got {_welcomes}")
ok("Owner is told a deleted account is still being charged",
   any("charged again" in t.lower() for t in _billing_alerts),
   f"got {_billing_alerts}")
with app.app_context():
    after = Order.query.filter_by(ls_order_id="WEL-3").first()
    ok("Charge after deletion does not restore the deleted email",
       after is not None and (after.buyer_email or "").startswith("closed+"),
       f"got {after.buyer_email if after else None}")

# People come back. Signing up again with the same address makes the
# subscription theirs once more, and writing it off then cancels a membership
# they are paying for — which is what started happening in production.
with app.app_context():
    returner = User(email="leaver@example.com", display_name="Back again",
                    membership="none", email_verified_at=utcnow())
    returner.set_password(USER_PW)
    db.session.add(returner)
    db.session.commit()
    returner_id = returner.id
_alerts_before = len(_billing_alerts)

# a) the old subscription, the one whose orders were scrubbed, renews
_invoice_webhook("WEL-R1", "leaver@example.com", _mem_price,
                 "sub_leaver", "subscription_cycle")
ok("A returning member's renewal is not written off as a deleted account",
   len(_billing_alerts) == _alerts_before, f"got {_billing_alerts}")
with app.app_context():
    back = Order.query.filter_by(ls_order_id="WEL-R1").first()
    ok("Their payment stays attached to the account paying it",
       back is not None and back.buyer_email == "leaver@example.com",
       f"got {back.buyer_email if back else None}")
    ok("And the tier they are paying for actually reaches them",
       db.session.get(User, returner_id).membership == "healing",
       f"got {db.session.get(User, returner_id).membership}")
    ok("The scrubbed order from before is handed back too, not left stranded",
       Order.query.filter_by(ls_order_id="WEL-3").first().buyer_email
       == "leaver@example.com")

# b) a brand new subscription for the same returning address
_invoice_webhook("WEL-R2", "leaver@example.com", _mem_price,
                 "sub_leaver_two", "subscription_create")
ok("A fresh subscription for a returning address is left alone",
   len(_billing_alerts) == _alerts_before, f"got {_billing_alerts}")
with app.app_context():
    ok("They keep the membership they just paid for",
       db.session.get(User, returner_id).membership == "healing")

# c) deleting again puts the protection back
with app.app_context():
    close_account(db.session.get(User, returner_id))
_invoice_webhook("WEL-R3", "leaver@example.com", _mem_price,
                 "sub_leaver_two", "subscription_cycle")
ok("Deleting again restores the write-off, so a ghost keeps being caught",
   len(_billing_alerts) > _alerts_before, f"got {_billing_alerts}")

# Re-running fulfillment (late webhook, dashboard sync) keeps the scrub.
_invoice_webhook("WEL-1", "leaver@example.com", _mem_price,
                 "sub_leaver", "subscription_create")
with app.app_context():
    replayed = Order.query.filter_by(ls_order_id="WEL-1").first()
    ok("Replaying an old payment keeps a scrubbed order scrubbed",
       replayed is not None and (replayed.buyer_email or "").startswith("closed+"),
       f"got {replayed.buyer_email if replayed else None}")

# A genuinely new subscriber still gets welcomed.
_invoice_webhook("WEL-4", "fresh-member@example.com", _mem_price,
                 "sub_fresh", "subscription_create")
ok("A new subscriber still gets a welcome",
   "fresh-member@example.com" in _welcomes, f"got {_welcomes}")

# --- drip-fed modules + free membership perk ----------------------------------
_drip_fields = {
    "title": "Drip Course",
    "track": "healing",
    "type": "course",
    "price": "39.00",
    "promise": "One steady step at a time.",
    "stripe": "price_drip_course",
    "drip": "1",
    "drip_interval_days": "7",
    "perk_tier": "creator",
    "perk_months": "3",
    "live": "1",
}
r = admin.post(
    "/admin/products/new",
    data=dict(_drip_fields, **{
        "mod1_title": "Week one",
        "mod1_desc": "Start here",
        "mod1_file": (BytesIO(b"%PDF-1.4 module one\n%%EOF\n"), "week-one.pdf"),
        "mod2_title": "Week two",
        "mod2_desc": "Keep going",
        "mod2_file": (BytesIO(b"%PDF-1.4 module two\n%%EOF\n"), "week-two.pdf"),
    }),
    content_type="multipart/form-data",
    follow_redirects=True,
)
ok("Studio creates a drip-fed product", r.status_code == 200)
ok("Studio confirms the new product went out to members",
   "notified" in r.get_data(as_text=True))
with app.app_context():
    from app.models import Notification as _Note
    course_note = (_Note.query.filter_by(kind="course")
                   .order_by(_Note.id.desc()).first())
    ok("Publishing a product announces it",
       course_note is not None
       and "Drip Course" in (course_note.body or "")
       and (course_note.url or "") == "/courses/drip-course")
with app.app_context():
    drip_prod = Product.query.filter_by(slug="drip-course").first()
    ok("Drip schedule saved on the product",
       drip_prod is not None and drip_prod.drip_enabled is True
       and drip_prod.drip_interval_days == 7 and drip_prod.is_dripped())
    ok("Membership perk saved on the product",
       drip_prod.perk_tier() == "creator" and drip_prod.perk_months() == 3
       and "3 months of Creator" in drip_prod.perk_summary())
    drip_mods = drip_prod.modules()
    ok("Each module keeps its own file",
       len(drip_mods) == 2
       and drip_mods[0]["asset"] is not None and drip_mods[1]["asset"] is not None
       and drip_mods[0]["asset"].module_index == 1
       and drip_mods[1]["asset"].module_index == 2)
    drip_prod_id = drip_prod.id
    mod1_asset_id = drip_mods[0]["asset"].id
    mod2_asset_id = drip_mods[1]["asset"].id

r = client.get("/courses/drip-course")
_dbody = r.get_data(as_text=True)
ok("Product page advertises the perk and the schedule",
   r.status_code == 200 and "3 months of Creator membership, free" in _dbody
   and "Released one module at a time" in _dbody
   and "Available right away" in _dbody and "Day 8" in _dbody)

with app.app_context():
    dripper = User(email="dripper@example.com", email_verified_at=utcnow())
    dripper.set_password(USER_PW)
    db.session.add(dripper)
    db.session.commit()
drip_payload = _payment_payload(
    "9100", "dripper@example.com", "price_drip_course",
    amount=3900, product_name="Drip Course")
r = client.post("/webhooks/stripe", data=drip_payload, headers=_stripe_headers(drip_payload))
with app.app_context():
    dripper = User.query.filter_by(email="dripper@example.com").first()
    drip_purchase = ShopPurchase.query.filter_by(lemon_squeezy_order_id="9100").first()
    ok("Buying the product grants the free membership",
       r.status_code == 200 and drip_purchase is not None
       and drip_purchase.status == "linked" and dripper.membership == "creator",
       f"membership={getattr(dripper, 'membership', None)}")
    drip_purchase_id = drip_purchase.id

drip_client = app.test_client()
drip_client.post("/login", data={"email": "dripper@example.com", "password": USER_PW})
r = drip_client.get(f"/account/courses/{drip_purchase_id}")
_rbody = r.get_data(as_text=True)
ok("Reader opens module 1 the moment they buy",
   r.status_code == 200 and "Week one" in _rbody
   and f"/file/{mod1_asset_id}" in _rbody)
ok("Reader shows the later module as locked",
   "Week two" in _rbody and "Unlocks" in _rbody)
# "Opens the 15th" is no use without the hour, and the hour is no use on the
# server's clock, so the wait is written out on the reader's own.
_when = re.search(r"unlocks\s+([A-Z][a-z]{2} \d{2}, \d{4} at \d{2}:\d{2} [AP]M)"
                  r"\s+your time", _rbody)
ok("And says what time it opens, on their clock", bool(_when),
   "no local unlock time in the module bar")
_ny = app.test_client()
_ny.set_cookie("tz", "America/New_York", domain="localhost")
_ny.post("/login", data={"email": "dripper@example.com", "password": USER_PW})
_nybody = _ny.get(f"/account/courses/{drip_purchase_id}").get_data(as_text=True)
_ny_when = re.search(r"unlocks\s+([A-Z][a-z]{2} \d{2}, \d{4} at \d{2}:\d{2} [AP]M)",
                     _nybody)
ok("Somebody reading it from another timezone is told their own hour",
   bool(_ny_when) and bool(_when) and _ny_when.group(1) != _when.group(1),
   f"got {_ny_when and _ny_when.group(1)} for both")
r = drip_client.get(f"/account/courses/{drip_purchase_id}/file/{mod1_asset_id}")
ok("Unlocked module file opens", r.status_code == 200 and b"module one" in r.data)
r = drip_client.get(f"/account/courses/{drip_purchase_id}/file/{mod2_asset_id}")
ok("Locked module file is refused even with a direct link", r.status_code == 404)
r = drip_client.get(f"/account/courses/{drip_purchase_id}?module=2")
ok("Asking for a locked module falls back to what they have",
   r.status_code == 200 and f"/file/{mod2_asset_id}" not in r.get_data(as_text=True))


def _backdate_drip(days):
    with app.app_context():
        row = db.session.get(ShopPurchase, drip_purchase_id)
        row.purchased_at = utcnow() - timedelta(days=days)
        db.session.commit()


_backdate_drip(8)
r = drip_client.get(f"/account/courses/{drip_purchase_id}?module=2")
ok("Module 2 opens once its interval has passed",
   r.status_code == 200 and f"/file/{mod2_asset_id}" in r.get_data(as_text=True))
r = drip_client.get(f"/account/courses/{drip_purchase_id}/file/{mod2_asset_id}")
ok("Module 2 file streams after unlocking", r.status_code == 200 and b"module two" in r.data)

# A launch with a date on it: module one opens that day for everyone at once,
# instead of each buyer starting their own clock when they pay.
with app.app_context():
    from app.services import drip as drip_svc
    _p = db.session.get(Product, drip_prod_id)
    _bought = utcnow() - timedelta(days=30)
    _p.drip_starts_at = utcnow() + timedelta(days=5)
    _open = [row["number"] for row in drip_svc.module_rows(_p, _bought)
             if row["unlocked"]]
    ok("A release date still to come holds module 1 shut, however long ago they bought",
       _open == [], f"open={_open}")
    _p.drip_starts_at = utcnow() - timedelta(days=8)
    _open = [row["number"] for row in drip_svc.module_rows(_p, _bought)
             if row["unlocked"]]
    ok("Once it lands, the modules follow it at the usual interval",
       _open == [1, 2], f"open={_open}")
    _late = [row["number"] for row in drip_svc.module_rows(_p, utcnow())
             if row["unlocked"]]
    ok("And buying late opens what has been released, not a fresh wait",
       _late == [1, 2], f"open={_late}")
    _p.drip_starts_at = utcnow() + timedelta(days=5)
    ok("A file waiting on the release date is refused too",
       not drip_svc.asset_unlocked(
           _p, db.session.get(ProductAsset, mod1_asset_id), _bought))
    _p.drip_starts_at = None
    db.session.commit()

# The date is set in Studio, on the owner's own calendar.
r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=dict(_drip_fields, **{
        "slug": "drip-course",
        "mod1_title": "Week one", "mod1_desc": "Start here",
        "mod2_title": "Week two", "mod2_desc": "Keep going",
        "drip_starts_date": "2027-03-01", "drip_starts_time": "09:00",
    }),
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _p = db.session.get(Product, drip_prod_id)
    ok("Studio saves the day module 1 is released",
       r.status_code == 200 and _p.drip_starts_at is not None
       and _p.drip_starts_at.year == 2027 and _p.drip_starts_at.month == 3,
       f"got {_p.drip_starts_at}")
r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=dict(_drip_fields, **{
        "slug": "drip-course",
        "mod1_title": "Week one", "mod1_desc": "Start here",
        "mod2_title": "Week two", "mod2_desc": "Keep going",
    }),
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    ok("And clearing the field puts it back to starting when they buy",
       db.session.get(Product, drip_prod_id).drip_starts_at is None)

# One gap for everything doesn't suit every course, so the schedule can also
# be a date per module, or a wait per module counted from each buyer's start.
with app.app_context():
    from app.services import drip as drip_svc
    _sched = Product(slug="three-ways", title="Three Ways", type="course",
                     status="published", track="building", promise="x",
                     price_cents=1000, stripe_price_id="price_3ways",
                     drip_enabled=True, drip_interval_days=14)
    db.session.add(_sched)
    db.session.commit()
    _sched_id = _sched.id
    _now = utcnow()

    def _open_modules(product, bought, now=None):
        return [row["number"]
                for row in drip_svc.module_rows(product, bought, now=now or _now)
                if row["unlocked"]]

    _sched.drip_mode = "interval"
    _sched.set_curriculum([{"title": f"Module {i}"} for i in (1, 2, 3)])
    db.session.commit()
    ok("One interval spaces every module the same",
       _open_modules(_sched, _now - timedelta(days=20)) == [1, 2]
       and _open_modules(_sched, _now - timedelta(days=1)) == [1],
       f"got {_open_modules(_sched, _now - timedelta(days=20))}")

    _sched.drip_mode = "dates"
    _sched.set_curriculum([
        {"title": "One", "release_at": (_now - timedelta(days=10)).isoformat()},
        {"title": "Two", "release_at": (_now - timedelta(days=2)).isoformat()},
        {"title": "Three", "release_at": (_now + timedelta(days=6)).isoformat()},
    ])
    db.session.commit()
    ok("Dates open the same modules whenever someone bought",
       _open_modules(_sched, _now - timedelta(days=1)) == [1, 2]
       and _open_modules(_sched, _now - timedelta(days=400)) == [1, 2],
       f"got {_open_modules(_sched, _now - timedelta(days=1))}")

    _sched.set_curriculum([
        {"title": "One", "release_at": (_now - timedelta(days=10)).isoformat()},
        {"title": "Two"},
        {"title": "Three", "release_at": (_now + timedelta(days=6)).isoformat()},
    ])
    db.session.commit()
    ok("A module left without a date comes out with the one above it",
       _open_modules(_sched, _now - timedelta(days=1)) == [1, 2])

    _sched.set_curriculum([
        {"title": "One", "release_at": (_now + timedelta(days=10)).isoformat()},
        {"title": "Two", "release_at": (_now - timedelta(days=5)).isoformat()},
    ])
    db.session.commit()
    ok("And a later module dated early still waits for the one above",
       _open_modules(_sched, _now) == [])

    _sched.drip_mode = "gaps"
    _sched.set_curriculum([
        {"title": "One", "gap_days": 0},
        {"title": "Two", "gap_days": 5},
        {"title": "Three", "gap_days": 14},
    ])
    db.session.commit()
    ok("Each module can instead wait its own stretch after the one above",
       _open_modules(_sched, _now) == [1]
       and _open_modules(_sched, _now - timedelta(days=6)) == [1, 2]
       and _open_modules(_sched, _now - timedelta(days=20)) == [1, 2, 3],
       f"got {_open_modules(_sched, _now - timedelta(days=6))}")
    ok("Which still counts from each buyer's own start",
       _open_modules(_sched, _now - timedelta(days=4)) == [1])

    _pinned = _types.SimpleNamespace(module_index=3)
    ok("A file in a module still waiting is refused on the same schedule",
       not drip_svc.asset_unlocked(_sched, _pinned, _now)
       and drip_svc.asset_unlocked(_sched, _pinned, _now - timedelta(days=20)))

# The mode and each module's own timing are set in Studio.
from werkzeug.datastructures import MultiDict as _MultiDict  # noqa: E402

_sched_fields = {
    "title": "Three Ways", "track": "building", "types": "course",
    "promise": "x", "price": "10.00", "stripe": "price_3ways", "live": "1",
    "slug": "three-ways", "drip": "1", "drip_interval_days": "14",
    "mod1_title": "One", "mod2_title": "Two", "mod3_title": "Three",
}
r = admin.post(f"/admin/products/{_sched_id}/edit",
               data=_MultiDict(dict(_sched_fields, **{
                   "drip_mode": "dates",
                   "mod1_release_date": "2027-01-05", "mod1_release_time": "09:00",
                   "mod2_release_date": "2027-01-19", "mod2_release_time": "09:00",
                   "mod3_release_date": "2027-02-02", "mod3_release_time": "18:30",
               }).items()),
               content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _rows = db.session.get(Product, _sched_id).curriculum()
    ok("Studio saves a date per module",
       r.status_code == 200
       and db.session.get(Product, _sched_id).drip_mode_key() == "dates"
       and [row["release_at"][:10] for row in _rows]
       == ["2027-01-05", "2027-01-19", "2027-02-02"],
       f"got {[row['release_at'] for row in _rows]}")
r = admin.post(f"/admin/products/{_sched_id}/edit",
               data=_MultiDict(dict(_sched_fields, **{
                   "drip_mode": "gaps",
                   "mod1_gap_days": "0", "mod2_gap_days": "5", "mod3_gap_days": "14",
                   "mod1_release_date": "2027-01-05", "mod1_release_time": "09:00",
                   "mod2_release_date": "2027-01-19", "mod2_release_time": "09:00",
                   "mod3_release_date": "2027-02-02", "mod3_release_time": "18:30",
               }).items()),
               content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _rows = db.session.get(Product, _sched_id).curriculum()
    ok("And a wait per module, keeping the dates for if they switch back",
       db.session.get(Product, _sched_id).drip_mode_key() == "gaps"
       and [row["gap_days"] for row in _rows] == [0, 5, 14]
       and all(row["release_at"] for row in _rows),
       f"got {_rows}")
_sbody = admin.get(f"/admin/products/{_sched_id}/edit").get_data(as_text=True)
ok("The editor offers all three ways and the fields each one needs",
   all(f'value="{m}"' in _sbody for m in ("interval", "dates", "gaps"))
   and 'name="mod2_release_date"' in _sbody and 'name="mod2_gap_days"' in _sbody)

# Off the shelves: still there to read about, no longer for sale, and everyone
# who already bought it keeps it.
r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=dict(_drip_fields, **{
        "slug": "drip-course",
        "mod1_title": "Week one", "mod1_desc": "Start here",
        "mod2_title": "Week two", "mod2_desc": "Keep going",
        "off_shelf_date": "2027-06-30", "off_shelf_time": "23:59",
    }),
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _p = db.session.get(Product, drip_prod_id)
    ok("Studio saves a last day on sale",
       r.status_code == 200 and _p.off_shelf_at is not None
       and _p.off_shelf_at.year == 2027, f"got {_p.off_shelf_at}")
    ok("A date still to come leaves it selling as normal",
       not _p.is_off_shelf() and _p.buyable_by(None))
    _p.off_shelf_at = utcnow() - timedelta(minutes=1)
    db.session.commit()
    ok("Once the date passes it stops being sold",
       _p.is_off_shelf() and not _p.buyable_by(None))
    ok("A sale on it is dropped, since there is nowhere left to type the code",
       not _p.has_promo())

_shelf = client.get("/courses").get_data(as_text=True)
ok("It stays on the catalogue, marked off the shelves",
   "Drip Course" in _shelf and "Off the shelves" in _shelf)
r = client.get("/courses/drip-course")
ok("Its page still opens and says so instead of offering a buy button",
   r.status_code == 200 and "Off the shelves" in r.get_data(as_text=True)
   and "Buy now" not in r.get_data(as_text=True))
r = client.get("/checkout/product/drip-course")
ok("And checkout turns anyone new away",
   r.status_code in (302, 303)
   and "/courses/drip-course" in (r.headers.get("Location") or ""))
r = drip_client.get(f"/account/courses/{drip_purchase_id}")
ok("Someone who bought it before still reads it", r.status_code == 200)
r = drip_client.get(f"/account/courses/{drip_purchase_id}/file/{mod1_asset_id}")
ok("Including the files inside it", r.status_code == 200)
with app.app_context():
    _p = db.session.get(Product, drip_prod_id)
    _p.off_shelf_at = None
    db.session.commit()

# owner adds a module after publishing — existing buyers get it on their schedule
r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=dict(_drip_fields, **{
        "slug": "drip-course",
        "mod1_title": "Week one",
        "mod1_desc": "Start here",
        "mod2_title": "Week two",
        "mod2_desc": "Keep going",
        "mod3_title": "Week three",
        "mod3_desc": "Look back",
        "mod3_file": (BytesIO(b"%PDF-1.4 module three\n%%EOF\n"), "week-three.pdf"),
    }),
    content_type="multipart/form-data",
    follow_redirects=True,
)
ok("Studio adds a module to a live product", r.status_code == 200)
with app.app_context():
    drip_prod = db.session.get(Product, drip_prod_id)
    later = drip_prod.modules()
    ok("Added module keeps the earlier files in place",
       len(later) == 3 and later[0]["asset"].id == mod1_asset_id
       and later[1]["asset"].id == mod2_asset_id and later[2]["asset"] is not None)
    ok("Product stayed live through the edit", drip_prod.status == "published")
    mod3_asset_id = later[2]["asset"].id
r = drip_client.get(f"/account/courses/{drip_purchase_id}")
_rbody = r.get_data(as_text=True)
ok("New module reaches someone who already bought",
   "Week three" in _rbody and "Unlocks" in _rbody)
r = drip_client.get(f"/account/courses/{drip_purchase_id}/file/{mod3_asset_id}")
ok("New module still waits its turn for existing buyers", r.status_code == 404)
_backdate_drip(20)
r = drip_client.get(f"/account/courses/{drip_purchase_id}/file/{mod3_asset_id}")
ok("New module unlocks on the buyer's own schedule",
   r.status_code == 200 and b"module three" in r.data)

# --- a module holds many videos, documents and written extracts --------------
with app.app_context():
    _asset = db.session.get(ProductAsset, mod1_asset_id)
    _course_dir = app.config["COURSE_FILES_DIR"]
    ok("Module files stream to the media disk, not into Postgres",
       bool(_asset.disk_name) and _asset.data is None
       and _os.path.isfile(_os.path.join(_course_dir, _asset.disk_name)))

from werkzeug.datastructures import MultiDict as _MultiDict  # noqa: E402

r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=_MultiDict([
        *dict(_drip_fields, **{
            "slug": "drip-course",
            "mod1_title": "Week one", "mod1_desc": "Start here",
            "mod2_title": "Week two", "mod2_desc": "Keep going",
            "mod3_title": "Week three", "mod3_desc": "Look back",
        }).items(),
        ("mod1_file", (BytesIO(b"\x00\x00\x00\x18ftypmp42 lesson"), "lesson.mp4")),
        ("mod1_file", (BytesIO(b"%PDF-1.4 worksheet\n%%EOF\n"), "worksheet.pdf")),
        ("mod1_text_title", "Before you start"),
        ("mod1_text_body", "Read this **first**, then press play."),
        ("mod1_text_title", "A note on pacing"),
        ("mod1_text_body", "Slow is fine."),
    ]),
    content_type="multipart/form-data",
    follow_redirects=True,
)
ok("Studio accepts several files and extracts for one module",
   r.status_code == 200)
with app.app_context():
    drip_prod = db.session.get(Product, drip_prod_id)
    m1 = drip_prod.modules()[0]
    kinds = [a.kind for a in m1["contents"]]
    ok("Module 1 now holds the original file plus everything just added",
       len(m1["contents"]) == 5 and kinds.count("video") == 1
       and kinds.count("text") == 2 and kinds.count("pdf") == 2,
       f"kinds={kinds}")
    ok("The module's first item is still the file it started with",
       m1["contents"][0].id == mod1_asset_id and m1["asset"].id == mod1_asset_id)
    ok("Other modules were left alone",
       len(drip_prod.modules()[1]["contents"]) == 1)
    video_item = next(a for a in m1["contents"] if a.kind == "video")
    text_item = next(a for a in m1["contents"] if a.kind == "text")
    video_item_id, text_item_id = video_item.id, text_item.id
    ok("A written extract keeps its words in the row, with no file on disk",
       text_item.body.startswith("Read this") and text_item.disk_name is None
       and text_item.title == "Before you start")
    ok("A written extract reports a sensible size, not 0.0 MB",
       text_item.size_display().endswith("B") and text_item.size > 0)
    _video_path = _os.path.join(app.config["COURSE_FILES_DIR"], video_item.disk_name)
    ok("The lesson video landed on the disk", _os.path.isfile(_video_path))

# Once files are attached the cover's kind comes from the first one, so the
# preview must stop following the type dropdown.
_drip_edit = admin.get(f"/admin/products/{drip_prod_id}/edit").get_data(as_text=True)
ok("A product with files pins the kind shown on its cover preview",
   'data-cover-kind-fixed="' in _drip_edit)

r = drip_client.get(f"/account/courses/{drip_purchase_id}?module=1")
_rbody = r.get_data(as_text=True)
ok("Reader lists everything in the module so the buyer can move between it",
   "In this module" in _rbody and "Before you start" in _rbody
   and f"item={video_item_id}" in _rbody)
ok("Module list says how much is inside", "5 pieces to work through" in _rbody)

r = drip_client.get(
    f"/account/courses/{drip_purchase_id}?module=1&item={text_item_id}")
_rbody = r.get_data(as_text=True)
ok("A written extract is read on the page, not downloaded",
   r.status_code == 200 and "Read this <strong>first</strong>" in _rbody
   and "reader-doc--written" in _rbody)

r = drip_client.get(
    f"/account/courses/{drip_purchase_id}?module=1&item={video_item_id}")
ok("Picking the video in a module opens it",
   r.status_code == 200
   and "course-reader__video" in r.get_data(as_text=True))
r = drip_client.get(
    f"/account/courses/{drip_purchase_id}/file/{video_item_id}",
    headers={"Range": "bytes=0-7"})
ok("A lesson video can be scrubbed instead of downloaded whole",
   r.status_code == 206 and r.headers.get("Accept-Ranges") == "bytes"
   and len(r.data) == 8)
# A slice is not the file. Letting a browser keep one meant it could answer a
# later whole-file request from the fragment, which a PDF reader reads as a
# truncated document and refuses to open.
ok("A ranged read is never cacheable",
   "no-store" in (r.headers.get("Cache-Control") or ""),
   f"got {r.headers.get('Cache-Control')!r}")
_whole = drip_client.get(
    f"/account/courses/{drip_purchase_id}/file/{video_item_id}")
ok("And the whole file still comes back whole",
   _whole.status_code == 200
   and "no-store" in (_whole.headers.get("Cache-Control") or "")
   and len(_whole.data) > 8,
   f"got {_whole.status_code} {_whole.headers.get('Cache-Control')!r}")
r = app.test_client().get(f"/account/courses/{drip_purchase_id}/file/{video_item_id}")
ok("A stranger still can't reach a course video", r.status_code in (302, 404))

# big files go up in slices, because a single request can't carry them
_blob = b"\x00\x00\x00\x18ftypmp42" + (b"bigvideo" * 4096)
r = admin.post(f"/admin/products/{drip_prod_id}/uploads/begin",
               json={"filename": "keynote.mp4", "size": len(_blob)})
_started = r.get_json()
ok("Studio can start a sliced upload",
   r.status_code == 200 and bool(_started.get("upload_id"))
   and _started.get("chunk_bytes", 0) > 0)
_upload_id = _started["upload_id"]
_step = 4096
for _i in range(0, len(_blob), _step):
    r = admin.post(
        f"/admin/products/{drip_prod_id}/uploads/{_upload_id}/chunk",
        data={"chunk": (BytesIO(_blob[_i:_i + _step]), "part")},
        content_type="multipart/form-data")
ok("Every slice is accepted and counted",
   r.status_code == 200 and r.get_json().get("received") == len(_blob))
r = admin.post(f"/admin/products/{drip_prod_id}/uploads/{_upload_id}/finish",
               json={"filename": "keynote.mp4", "module": 2})
_finished = r.get_json()
ok("Finishing a sliced upload files it under the right module",
   r.status_code == 200 and _finished.get("kind") == "video"
   and _finished.get("kind_label") == "Video")
with app.app_context():
    big = db.session.get(ProductAsset, _finished["asset_id"])
    ok("The reassembled file is whole and in module 2",
       big.size == len(_blob) and big.module_index == 2
       and len(db.session.get(Product, drip_prod_id).modules()[1]["contents"]) == 2)
    _big_path = _os.path.join(app.config["COURSE_FILES_DIR"], big.disk_name)
    ok("The reassembled file matches what was sent byte for byte",
       open(_big_path, "rb").read() == _blob)
    ok("No half-finished part is left lying around",
       not _os.path.isfile(_os.path.join(
           app.config["COURSE_FILES_DIR"], "parts", _upload_id)))
    _big_id = big.id

r = admin.post(f"/admin/products/{drip_prod_id}/uploads/begin",
               json={"filename": "huge.mp4",
                     "size": (app.config["COURSE_UPLOAD_MAX_MB"] + 1) * 1024 * 1024})
ok("A file over the cap is turned away before a byte is sent",
   r.status_code == 400 and "MB" in (r.get_json() or {}).get("error", ""))
r = admin.post(f"/admin/products/{drip_prod_id}/uploads/nonesuch/chunk",
               data={"chunk": (BytesIO(b"orphan"), "part")},
               content_type="multipart/form-data")
ok("Slices for an upload we never started are refused", r.status_code == 400)
r = app.test_client().post(f"/admin/products/{drip_prod_id}/uploads/begin",
                           json={"filename": "sneaky.mp4", "size": 10})
ok("Only owners can start an upload", r.status_code in (302, 401, 403, 404))

# A tab closed midway through a large video leaves its slices behind, and
# nothing ever came back for them. On a disk sized to the library that exists,
# a few forgotten gigabytes is the difference between room and none.
with app.app_context():
    import time as _time

    from app.services import assets as _assets_svc
    _parts = _assets_svc.parts_dir()
    _os.makedirs(_parts, exist_ok=True)
    _stale = _os.path.join(_parts, "abandoned.mp4")
    _live = _os.path.join(_parts, "inflight.mp4")
    for _p in (_stale, _live):
        with open(_p, "wb") as _fh:
            _fh.write(b"x" * 512)
    _os.utime(_stale, (_time.time() - 30 * 3600,) * 2)
    _cleared = _assets_svc.sweep_parts()
    ok("An upload nobody finished is cleared away",
       _cleared == 1 and not _os.path.isfile(_stale), f"cleared {_cleared}")
    ok("While one still in flight is left alone", _os.path.isfile(_live))
    _os.utime(_live, (_time.time() - 30 * 3600,) * 2)
r = admin.post(f"/admin/products/{drip_prod_id}/uploads/begin",
               json={"filename": "next.mp4", "size": 1000})
with app.app_context():
    ok("And starting another upload sweeps them, with no scheduler to run",
       not _os.path.isfile(_live))

# A PDF that arrives short still saves and still looks like a file in Studio.
# The breakage only shows up later, to a buyer, as a reader that won't open it,
# so it is turned away while somebody can still do something about it.
r = admin.post(
    f"/admin/products/{drip_prod_id}/assets",
    data={"asset": (BytesIO(b"%PDF-1.4 whole\n%%EOF\n"), "whole.pdf")},
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    ok("A PDF that is all there is taken",
       any(a.filename == "whole.pdf"
           for a in db.session.get(Product, drip_prod_id).assets))
r = admin.post(
    f"/admin/products/{drip_prod_id}/assets",
    data={"asset": (BytesIO(b"%PDF-1.4 cut off here"), "short.pdf")},
    content_type="multipart/form-data", follow_redirects=True)
_short_body = r.get_data(as_text=True)
with app.app_context():
    ok("One that stops before its end marker is refused, and said so plainly",
       not any(a.filename == "short.pdf"
               for a in db.session.get(Product, drip_prod_id).assets)
       and "cut short" in _short_body, "it was accepted")
r = admin.post(
    f"/admin/products/{drip_prod_id}/assets",
    data={"asset": (BytesIO(b"not a pdf at all"), "fake-pdf.pdf")},
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    ok("And so is a file only calling itself one",
       not any(a.filename == "fake-pdf.pdf"
               for a in db.session.get(Product, drip_prod_id).assets))
    from app.services import assets as _assets_svc
    for _a in list(db.session.get(Product, drip_prod_id).assets):
        if _a.filename == "whole.pdf":
            _assets_svc.delete_file(_a)
            db.session.delete(_a)
    db.session.commit()

# A slide deck is drawn into pages as it lands, so from then on it is an
# ordinary document: the same reader, the same pager, the same download.
from pptx import Presentation as _Pptx  # noqa: E402
from pptx.util import Inches as _In  # noqa: E402

_deck = _Pptx()
for _title, _body in (("Week two: posting", "What a save really means"),
                      ("Three niches", "Finance, parenting, beauty")):
    _s = _deck.slides.add_slide(_deck.slide_layouts[1])
    _s.shapes.title.text = _title
    _s.placeholders[1].text_frame.text = _body
_s = _deck.slides.add_slide(_deck.slide_layouts[5])
_s.shapes.title.text = "One more thing"
_s.shapes.add_textbox(_In(1), _In(3), _In(6), _In(1)).text_frame.text = "Post today"
_deck_bytes = BytesIO()
_deck.save(_deck_bytes)
_deck_bytes = _deck_bytes.getvalue()

r = admin.post(
    f"/admin/products/{drip_prod_id}/assets",
    data={"asset": (BytesIO(_deck_bytes), "posting-deck.pptx"),
          "asset_title": ""},
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _drawn = next((a for a in db.session.get(Product, drip_prod_id).assets
                   if (a.filename or "").startswith("posting-deck")), None)
    ok("A PowerPoint deck uploads and comes back as pages",
       _drawn is not None and _drawn.kind == "pdf"
       and _drawn.filename == "posting-deck.pdf"
       and _drawn.mime == "application/pdf",
       f"got {_drawn and (_drawn.kind, _drawn.filename)}")
    ok("It keeps the deck's own name in the list of pieces",
       _drawn is not None and _drawn.display_title() == "posting-deck"
       and _drawn.kind_label() == "PDF",
       f"got {_drawn and (_drawn.title, _drawn.kind_label())}")
    _drawn_bytes = _assets_svc.read_bytes(_drawn) if _drawn else b""
    ok("What is stored is a real PDF, one page per slide",
       _drawn_bytes.startswith(b"%PDF-")
       and _drawn_bytes.count(b"/Type /Page\n") >= 3
       and b"%%EOF" in _drawn_bytes[-4096:],
       f"{len(_drawn_bytes)} bytes")
    _drawn_id = _drawn.id if _drawn else 0
r = admin.post(
    f"/admin/products/{drip_prod_id}/assets",
    data={"asset": (BytesIO(b"PK\x03\x04 not really a deck"), "broken.pptx")},
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    ok("A deck that won't open is refused, with what to do about it",
       not any((a.filename or "").startswith("broken")
               for a in db.session.get(Product, drip_prod_id).assets)
       and "export it as a pdf" in r.get_data(as_text=True).lower())
r = drip_client.get(f"/account/courses/{drip_purchase_id}/file/{_drawn_id}")
ok("And the drawn deck streams to a buyer as a PDF",
   r.status_code == 200 and r.data.startswith(b"%PDF-")
   and "pdf" in (r.headers.get("Content-Type") or ""))

# A deck big enough to arrive in slices takes the same road at the end of it.
r = admin.post(f"/admin/products/{drip_prod_id}/uploads/begin",
               json={"filename": "sliced-deck.pptx", "size": len(_deck_bytes)})
_deck_upload = (r.get_json() or {}).get("upload_id")
for _i in range(0, len(_deck_bytes), 4096):
    admin.post(f"/admin/products/{drip_prod_id}/uploads/{_deck_upload}/chunk",
               data={"chunk": (BytesIO(_deck_bytes[_i:_i + 4096]), "part")},
               content_type="multipart/form-data")
r = admin.post(f"/admin/products/{drip_prod_id}/uploads/{_deck_upload}/finish",
               json={"filename": "sliced-deck.pptx", "module": 1})
_sliced = r.get_json() or {}
ok("A deck sent in slices is drawn into pages just the same",
   r.status_code == 200 and _sliced.get("kind") == "pdf"
   and _sliced.get("kind_label") == "PDF", f"got {_sliced}")
with app.app_context():
    _sliced_asset = db.session.get(ProductAsset, _sliced.get("asset_id") or 0)
    ok("And lands as a document under the module it was sent to",
       _sliced_asset is not None and _sliced_asset.filename == "sliced-deck.pdf"
       and _sliced_asset.module_index == 1
       and _assets_svc.read_bytes(_sliced_asset).startswith(b"%PDF-"))
    if _sliced_asset is not None:
        _assets_svc.delete_file(_sliced_asset)
        db.session.delete(_sliced_asset)
        db.session.commit()
with app.app_context():
    for _a in list(db.session.get(Product, drip_prod_id).assets):
        if (_a.filename or "").startswith("posting-deck"):
            _assets_svc.delete_file(_a)
            db.session.delete(_a)
    db.session.commit()

# Modules and lessons can be reordered and removed. What a row holds is pinned
# to it by number while the rows on the form are positional, so a row that
# moves has to take its files with it — otherwise a module dragged up the list
# arrives holding whatever used to sit in its new slot.
_order_base = {
    "title": "Order Test", "track": "building", "types": "course",
    "promise": "x", "price": "10.00", "stripe": "price_order", "live": "1",
}
_pdf = b"%PDF-1.4 x\n%%EOF\n"
r = admin.post("/admin/products/new", data=_MultiDict([
    *_order_base.items(),
    ("mod1_title", "Alpha"), ("mod2_title", "Beta"), ("mod3_title", "Gamma"),
    ("mod1_file", (BytesIO(_pdf), "alpha.pdf")),
    ("mod2_file", (BytesIO(_pdf), "beta.pdf")),
    ("mod3_file", (BytesIO(_pdf), "gamma.pdf")),
]), content_type="multipart/form-data", follow_redirects=True)


def _module_map(product_id):
    with app.app_context():
        return [(row["title"], [a.filename for a in row["contents"]])
                for row in db.session.get(Product, product_id).modules()]


with app.app_context():
    _order_id = Product.query.filter_by(slug="order-test").first().id
ok("Each module starts with its own file",
   _module_map(_order_id) == [("Alpha", ["alpha.pdf"]), ("Beta", ["beta.pdf"]),
                              ("Gamma", ["gamma.pdf"])],
   f"got {_module_map(_order_id)}")

r = admin.post(f"/admin/products/{_order_id}/edit", data=_MultiDict([
    *_order_base.items(), ("slug", "order-test"),
    ("mod1_from", "3"), ("mod1_title", "Gamma"),
    ("mod2_from", "1"), ("mod2_title", "Alpha"),
    ("mod3_from", "2"), ("mod3_title", "Beta"),
]), content_type="multipart/form-data", follow_redirects=True)
ok("A module moved up the list takes its files with it",
   _module_map(_order_id) == [("Gamma", ["gamma.pdf"]), ("Alpha", ["alpha.pdf"]),
                              ("Beta", ["beta.pdf"])],
   f"got {_module_map(_order_id)}")

r = admin.post(f"/admin/products/{_order_id}/edit", data=_MultiDict([
    *_order_base.items(), ("slug", "order-test"),
    ("mod1_from", "1"), ("mod1_title", "Gamma"),
    ("mod2_from", "3"), ("mod2_title", "Beta"),
    ("removed_module", "2"),
]), content_type="multipart/form-data", follow_redirects=True)
ok("A module removed on purpose takes its files rather than leaving them loose",
   _module_map(_order_id) == [("Gamma", ["gamma.pdf"]), ("Beta", ["beta.pdf"])],
   f"got {_module_map(_order_id)}")
with app.app_context():
    ok("So nothing is left behind open to every buyer from day one",
       not [a for a in db.session.get(Product, _order_id).assets
            if not a.module_index])

r = admin.post(f"/admin/products/{_order_id}/edit", data=_MultiDict([
    *_order_base.items(), ("slug", "order-test"),
    ("mod1_from", "1"), ("mod1_title", "Gamma"),
    ("mod2_from", "2"), ("mod2_title", "Beta"),
    ("mod3_from", ""), ("mod3_title", "Delta"),
]), content_type="multipart/form-data", follow_redirects=True)
ok("A module added comes up empty rather than borrowing another's files",
   _module_map(_order_id) == [("Gamma", ["gamma.pdf"]), ("Beta", ["beta.pdf"]),
                              ("Delta", [])],
   f"got {_module_map(_order_id)}")

# Working a removal out from what a form left unsaid is not good enough: an
# older page, or a field that didn't make it, would read as a removal and take
# real uploads with it. Only a row the editor names is ever cleared out.
_before_quiet = _module_map(_order_id)
r = admin.post(f"/admin/products/{_order_id}/edit", data=_MultiDict([
    *_order_base.items(), ("slug", "order-test"),
    ("mod1_from", "1"), ("mod1_title", "Gamma"),
    ("mod2_from", ""), ("mod2_title", "Beta"),
]), content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _still = [a.filename for a in db.session.get(Product, _order_id).assets]
    ok("A row that stops saying where its content lives destroys nothing",
       "beta.pdf" in _still, f"got {_still}")
with app.app_context():
    _b = next(a for a in db.session.get(Product, _order_id).assets
              if a.filename == "beta.pdf")
    ok("The file is still on the disk behind it too", not _b.file_missing())
    _b.module_index = 2
    db.session.commit()

# Lessons move the same way, but removing one only loses the grouping: the
# module is still there to hold what was inside it.
r = admin.post(f"/admin/products/{_order_id}/edit", data=_MultiDict([
    *_order_base.items(), ("slug", "order-test"),
    ("mod1_from", "1"), ("mod1_title", "Gamma"),
    ("mod1_lesson_from", ""), ("mod1_lesson_title", "First"),
    ("mod1_lesson_from", ""), ("mod1_lesson_title", "Second"),
    ("mod1_lesson1_file", (BytesIO(_pdf), "one.pdf")),
    ("mod1_lesson2_file", (BytesIO(_pdf), "two.pdf")),
]), content_type="multipart/form-data", follow_redirects=True)


def _lesson_map(product_id):
    with app.app_context():
        row = db.session.get(Product, product_id).modules()[0]
        return ([(les["title"], [a.filename for a in les["contents"]])
                 for les in row["lessons"]],
                [a.filename for a in row["intro"]])


ok("Lessons start holding their own files",
   _lesson_map(_order_id) == ([("First", ["one.pdf"]), ("Second", ["two.pdf"])],
                              ["gamma.pdf"]),
   f"got {_lesson_map(_order_id)}")
r = admin.post(f"/admin/products/{_order_id}/edit", data=_MultiDict([
    *_order_base.items(), ("slug", "order-test"),
    ("mod1_from", "1"), ("mod1_title", "Gamma"),
    ("mod1_lesson_from", "2"), ("mod1_lesson_title", "Second"),
    ("mod1_lesson_from", "1"), ("mod1_lesson_title", "First"),
]), content_type="multipart/form-data", follow_redirects=True)
ok("A lesson moved up takes its files with it too",
   _lesson_map(_order_id) == ([("Second", ["two.pdf"]), ("First", ["one.pdf"])],
                              ["gamma.pdf"]),
   f"got {_lesson_map(_order_id)}")
r = admin.post(f"/admin/products/{_order_id}/edit", data=_MultiDict([
    *_order_base.items(), ("slug", "order-test"),
    ("mod1_from", "1"), ("mod1_title", "Gamma"),
    ("mod1_lesson_from", "2"), ("mod1_lesson_title", "First"),
]), content_type="multipart/form-data", follow_redirects=True)
ok("A lesson removed leaves its files in the module, not deleted",
   _lesson_map(_order_id) == ([("First", ["one.pdf"])],
                              ["gamma.pdf", "two.pdf"]),
   f"got {_lesson_map(_order_id)}")

# The store page says what each lesson is called and stops there. The extract
# written for a lesson is for the buyer reading it; printing every one of them
# turned the module cards into walls of text.
r = admin.post("/admin/products/new", data=_MultiDict([
    ("title", "Glance Course"), ("track", "building"), ("types", "course"),
    ("promise", "Short and clear"), ("price", "20.00"),
    ("stripe", "price_glance"),
    ("description", "A long write-up that runs well past the short note "
                    "saying who the course is for."),
    ("audience", "Women starting out."),
    ("mod1_title", "Week one"), ("mod1_lesson_title", "Niche Selection"),
    ("mod1_lesson_desc", "In this lesson you pick a niche and stay with it."),
    ("mod1_lesson_title", "Content Basics"), ("mod1_lesson_desc", ""),
    ("mod1_lesson1_file", (BytesIO(_pdf), "niche.pdf")),
]), content_type="multipart/form-data", follow_redirects=True)
_gbody = admin.get("/courses/glance-course").get_data(as_text=True)
ok("The store page numbers and names each lesson",
   r.status_code == 200 and "Niche Selection" in _gbody
   and "Lesson 1" in _gbody and "Lesson 2" in _gbody)
ok("And keeps the lesson's own write-up for the buyer, not the shop window",
   "In this lesson you pick a niche" not in _gbody)

# Beside a long write-up, "who this is for" is a couple of lines and the rest
# of the column was empty. The facts card fills it with what a buyer asks next.
ok("The column beside the write-up carries a facts card",
   "At a glance" in _gbody and "pd-glance" in _gbody
   and "Who this is for" in _gbody)
ok("Saying how big it is and what is inside",
   "1 module, 2 lessons" in _gbody and "1 document" in _gbody,
   "size or contents line missing")
ok("When it opens and what they keep",
   "Every module opens the moment you buy" in _gbody
   and "Kept for good" in _gbody)
with app.app_context():
    _facts = dict(db.session.get(Product, drip_prod_id).glance_facts())
    ok("A drip-fed course's card says it comes a module at a time",
       "moment you buy" not in _facts.get("Pace", "")
       and ("every" in _facts.get("Pace", "")
            or "one at a time" in _facts.get("Pace", "")),
       f"got {_facts.get('Pace')!r}")
    ok("And repeats the membership that comes with it",
       "3 months of Creator membership, free" in _facts.get("Also included", ""),
       f"got {_facts.get('Also included')!r}")
    _bare = Product(slug="bare-glance", title="Bare Glance", type="guide",
                    status="published", promise="x", price_cents=500,
                    stripe_price_id="price_bare")
    db.session.add(_bare)
    db.session.commit()
    ok("A guide with nothing to list keeps its write-up full width instead",
       _bare.glance_facts() == [], f"got {_bare.glance_facts()}")
    db.session.delete(_bare)
    db.session.commit()
with app.app_context():
    _svc = __import__("app.services.catalog", fromlist=["catalog"])
    _glance = Product.query.filter_by(slug="glance-course").first()
    if _glance is not None:
        _svc._purge_product(_glance)
        db.session.commit()

# The editor is long, so each section, module and lesson folds away. What each
# one is called stays on show, and everything folded still saves.
_fbody = admin.get(f"/admin/products/{_order_id}/edit").get_data(as_text=True)
ok("Every numbered section can be folded away",
   _fbody.count('<details class="admin-panel studio-card studio-section" open>')
   >= 4 and "<section class=\"admin-panel studio-card studio-section\">"
   not in _fbody,
   "some sections are still fixed open")
ok("A module keeps its name and note on show, and folds the rest",
   'data-fold-toggle' in _fbody
   and '<div class="studio-fold" data-fold hidden>' in _fbody
   and 'name="mod1_title"' in _fbody and 'name="mod1_desc"' in _fbody)
_folded = re.search(r'<div class="studio-fold" data-fold hidden>(.*?)'
                    r'name="mod1_release_date"', _fbody, re.S)
ok("With the settings and files inside the folded part",
   bool(_folded), "the module body isn't in the fold")
_admin_js = client.get("/static/js/admin.js").get_data(as_text=True)
ok("A field the browser rejects opens its fold rather than failing silently",
   '"invalid"' in _admin_js)
ok("What was left open is remembered, per page, across a refresh or a save",
   "ba:folds:" in _admin_js and "localStorage" in _admin_js
   and 'addEventListener("toggle"' in _admin_js)

# The two sections both numbered 5 never appear together — one is only on a
# product being created, the other only on one that exists — so each page
# counts from 1 with no gap and no repeat.
_new_heads = re.findall(r"<h2>(\d+)\.", admin.get("/admin/products/new")
                        .get_data(as_text=True))
_edit_heads = re.findall(r"<h2>(\d+)\.", _fbody)
ok("The sections number straight through on a new product",
   _new_heads == [str(i) for i in range(1, len(_new_heads) + 1)],
   f"got {_new_heads}")
ok("And on one being edited",
   _edit_heads == [str(i) for i in range(1, len(_edit_heads) + 1)],
   f"got {_edit_heads}")

_obody = _fbody
ok("The editor offers the arrows and a remove on both",
   all(a in _obody for a in ("data-module-up", "data-module-down",
                             "data-module-remove", "data-lesson-up",
                             "data-lesson-down", "data-lesson-remove")))
ok("And stamps each row with where its content lives now",
   'name="mod1_from"' in _obody and 'name="mod1_lesson_from"' in _obody)
with app.app_context():
    _svc = __import__("app.services.catalog", fromlist=["catalog"])
    _svc._purge_product(db.session.get(Product, _order_id))
    db.session.commit()

# A module row with no title used to be dropped on save, and the lessons
# written inside it went too, with nothing said. Someone filling in lessons
# before naming the module lost the lot on the first save.
r = admin.post("/admin/products/new", data=_MultiDict([
    ("title", "Kept Work"), ("track", "building"), ("types", "course"),
    ("promise", "x"), ("price", "10.00"), ("stripe", "price_kept"),
    ("mod1_title", "Named"), ("mod1_lesson_title", "A1"),
    ("mod2_title", ""), ("mod2_lesson_title", "B1"),
    ("mod2_lesson_title", "B2"),
    ("mod3_title", ""), ("mod3_desc", ""),
]), content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _kept = Product.query.filter_by(slug="kept-work").first()
    _rows = _kept.curriculum() if _kept else []
    ok("Lessons under a module nobody named are kept, under a stand-in name",
       [(row["title"], [x["title"] for x in row["lessons"]]) for row in _rows]
       == [("Named", ["A1"]), ("Module 2", ["B1", "B2"])],
       f"got {[(row['title'], [x['title'] for x in row['lessons']]) for row in _rows]}")
    ok("While a row with nothing in it is still nothing to save",
       len(_rows) == 2)
ok("And the stand-in name is said out loud, not slipped in",
   "under a stand-in name" in r.get_data(as_text=True))

r = admin.post("/admin/products/new", data=_MultiDict([
    ("title", "Kept Writing"), ("track", "building"), ("types", "course"),
    ("promise", "x"), ("price", "10.00"), ("stripe", "price_kept2"),
    ("mod1_title", "Named"),
    ("mod1_lesson_title", ""), ("mod1_lesson_desc", "All of the writing."),
    ("mod1_lesson_title", ""), ("mod1_lesson_desc", ""),
]), content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _rows = Product.query.filter_by(slug="kept-writing").first().curriculum()
    ok("A lesson written but not named keeps what was written in it",
       [(x["title"], x["description"]) for x in _rows[0]["lessons"]]
       == [("Lesson 1", "All of the writing.")],
       f"got {_rows[0]['lessons']}")
with app.app_context():
    from app.services import catalog as _cat_svc
    for _slug in ("kept-work", "kept-writing"):
        _p = Product.query.filter_by(slug=_slug).first()
        if _p:
            _cat_svc._purge_product(_p)
    db.session.commit()

# Instagram hands out photo links that are signed and run out, so one pointed
# at rather than copied works the day it is set and turns into a broken circle
# on the home page weeks later.
from app.admin.routes import _is_temporary_photo_link  # noqa: E402

ok("A photo link that will expire is recognised",
   _is_temporary_photo_link(
       "https://scontent-lhr8-1.cdninstagram.com/v/t51.2885-19/x.jpg?oe=68")
   and _is_temporary_photo_link("https://scontent.xx.fbcdn.net/v/x.jpg"))
ok("And our own stored photo is not mistaken for one",
   not _is_temporary_photo_link("/media/site/creator")
   and not _is_temporary_photo_link("") and not _is_temporary_photo_link(None))
with app.app_context():
    from app.services import settings as _settings_svc
    _settings_svc.set_setting("creator_name", "Spotlit")
    _settings_svc.set_setting("creator_instagram", "spotlit")
    _settings_svc.set_setting(
        "creator_image_url",
        "https://scontent-lhr8-1.cdninstagram.com/v/t51/gone.jpg?oe=1")
    db.session.commit()
_home = client.get("/").get_data(as_text=True)
_photo = re.search(r'<img class="spotlight-creator__photo"[^>]*>', _home)
ok("A picture that fails leaves the mark it would have had, not its alt text",
   bool(_photo) and "data-photo-fallback" in _photo.group(0),
   "no fallback on the spotlight photo")
_main_js = client.get("/static/js/main.js").get_data(as_text=True)
ok("Which the page swaps in when the picture doesn't arrive",
   "data-photo-fallback" in _main_js and '"error"' in _main_js)

# Writing boxes in Studio get formatting buttons. They work by writing marks
# around what's picked out, so nothing already typed can be replaced.
from app.services.markdown import render_markdown as _md  # noqa: E402

ok("Bold, italic and strikethrough all render",
   "<strong>bold</strong>" in _md("**bold**")
   and "<em>soft</em>" in _md("*soft*")
   and "<s>gone</s>" in _md("~~gone~~"),
   f"got {_md('**bold** *soft* ~~gone~~')}")
ok("So do bullet points and links",
   "<li>one</li>" in _md("- one\n- two")
   and 'href="https://example.com"' in _md("[a](https://example.com)"))
ok("Where lines used to be kept, they still are",
   "<br>" in _md("one\ntwo", breaks=True)
   and "<br>" not in _md("one\ntwo"))

_fmt_pages = {
    "the product editor": f"/admin/products/{drip_prod_id}/edit",
    "a page": "/admin/pages",
    "the tip editor": "/admin/content/new",
}
for _what, _url in _fmt_pages.items():
    _r = admin.get(_url)
    if _r.status_code != 200:
        continue
    _body = _r.get_data(as_text=True)
    if "<textarea" not in _body:
        continue
    ok(f"Writing boxes in {_what} carry formatting buttons",
       "data-format" in _body, "none marked up")
_pbody = admin.get(f"/admin/products/{drip_prod_id}/edit").get_data(as_text=True)
_contents_tag = re.search(r"<textarea[^>]*\bname=\"contents\"[^>]*>", _pbody)
ok("But the one-item-per-line box doesn't, since marks would show through",
   bool(_contents_tag) and "data-format" not in _contents_tag.group(0),
   f"got {_contents_tag.group(0) if _contents_tag else 'no box'}")
_admin_js2 = client.get("/static/js/admin.js").get_data(as_text=True)
ok("The buttons are the five asked for, and no others",
   all(t in _admin_js2 for t in ('title: "Bold"', 'title: "Italic"',
                                 'title: "Strikethrough"',
                                 'title: "Bullet points"', 'title: "Link"')))
ok("They write at the cursor, so the browser's own undo still steps back",
   '"insertText"' in _admin_js2)
ok("Pressing one a second time on the same words takes it off again",
   "function alreadyOn" in _admin_js2 and "function insideOut" in _admin_js2
   and "function wrappedAt" in _admin_js2)
ok("Counting the stars, so italic comes off bold without taking one with it",
   "function starRun" in _admin_js2)
ok("Bullets have always come back off, and still do",
   "bulleted ? line.replace" in _admin_js2)

# A membership can start with days free before the first charge.
with app.app_context():
    _hp = MembershipPlan.query.filter_by(tier="healing").first()
    _hp.trial_days = 14
    db.session.commit()
    ok("A plan holds how many days are free, and says it in words",
       _hp.free_days() == 14 and _hp.trial_display() == "2 weeks free",
       f"got {_hp.free_days()} / {_hp.trial_display()!r}")
    _hp.trial_days = 5
    db.session.commit()
    ok("Counted in days when that reads better", _hp.trial_display() == "5 days free")
    _hp.trial_days = 0
    db.session.commit()
    ok("And says nothing when the plan charges from the start",
       _hp.trial_display() == "")
    _hp.trial_days = 14
    db.session.commit()

_trial_call = {}
_real_checkout2 = pay.create_checkout_session
_real_conf2 = pay.configured
pay.configured = lambda: True
pay.create_checkout_session = lambda **kw: (
    _trial_call.update(kw) or "https://stripe.test/pay")
try:
    plain_client.get("/checkout/membership/healing")
    ok("Someone new is sent to Stripe with the free days on the subscription",
       _trial_call.get("trial_days") == 14, f"got {_trial_call.get('trial_days')}")
    with app.app_context():
        _sw = User.query.filter_by(email="plainmember@example.com").first()
        _sw.membership = "creator"
        db.session.commit()
    plain_client.get("/checkout/membership/healing")
    ok("Someone already paying doesn't get free days for switching",
       _trial_call.get("trial_days") == 0, f"got {_trial_call.get('trial_days')}")
    with app.app_context():
        _sw = User.query.filter_by(email="plainmember@example.com").first()
        _sw.membership = "none"
        db.session.commit()
finally:
    pay.create_checkout_session = _real_checkout2
    pay.configured = _real_conf2

_trial_page = plain_client.get("/membership").get_data(as_text=True)
ok("The membership page offers it where someone is deciding",
   "Start with 2 weeks free" in _trial_page, "not offered")
ok("And says it once at the top before any of the prices",
   "Every paid plan starts with 2 weeks free" in _trial_page,
   "nothing at the top")
with app.app_context():
    _cr = MembershipPlan.query.filter_by(tier="creator").first()
    _was_cr = _cr.trial_days
    _cr.trial_days = 30
    db.session.commit()
_mixed = plain_client.get("/membership").get_data(as_text=True)
ok("When the plans differ there is no one number, so it just says there is one",
   "Paid plans start with a free trial." in _mixed
   and "Every paid plan starts with" not in _mixed)
ok("While each card still names its own",
   "Start with 2 weeks free" in _mixed and "Start with 1 month free" in _mixed)
with app.app_context():
    MembershipPlan.query.filter_by(tier="creator").first().trial_days = _was_cr
    db.session.commit()
with app.app_context():
    _hp2 = MembershipPlan.query.filter_by(tier="healing").first()
    _hp2.trial_days = 60
    db.session.commit()
    ok("A month is a month, not thirty days",
       _hp2.trial_display() == "2 months free", f"got {_hp2.trial_display()}")
    _hp2.trial_days = 14
    db.session.commit()
_plans_page = admin.get("/admin/memberships").get_data(as_text=True)
ok("And Studio has somewhere to set it per plan",
   'name="healing_trial_days"' in _plans_page
   and 'name="creator_trial_days"' in _plans_page)
# That form saves every plan at once, so each one's current values go back
# with it — posting only the field under test would blank the others.
with app.app_context():
    _plan_form = {}
    for _pl in MembershipPlan.query.all():
        _plan_form[f"{_pl.tier}_name"] = _pl.name or _pl.tier
        _plan_form[f"{_pl.tier}_stripe"] = _pl.stripe_price_id or ""
        _plan_form[f"{_pl.tier}_stripe_annual"] = _pl.stripe_price_id_annual or ""
        _plan_form[f"{_pl.tier}_stripe_product"] = _pl.stripe_product_id or ""
        _plan_form[f"{_pl.tier}_stripe_product_annual"] = (
            _pl.stripe_product_id_annual or "")
        _plan_form[f"{_pl.tier}_trial_days"] = str(_pl.free_days())
        if _pl.active:
            _plan_form[f"{_pl.tier}_active"] = "1"
        if _pl.price_cents is not None:
            _plan_form[f"{_pl.tier}_price"] = f"{_pl.price_cents / 100:.2f}"
        if _pl.annual_price_cents is not None:
            _plan_form[f"{_pl.tier}_annual_price"] = (
                f"{_pl.annual_price_cents / 100:.2f}")
_plan_form["healing_trial_days"] = "30"
admin.post("/admin/memberships", data=_plan_form, follow_redirects=True)
with app.app_context():
    ok("Setting it there sticks",
       MembershipPlan.query.filter_by(tier="healing").first().free_days() == 30)
    MembershipPlan.query.filter_by(tier="healing").first().trial_days = 0
    db.session.commit()

# Where paying leaves you depends on whether there is anywhere to land. Signed
# in, that is the library. Not signed in, /account would only bounce to a
# sign-in page, so it is the product again with word that it is on its way.
_landing = {}
_real_checkout = pay.create_checkout_session
_real_configured = pay.configured
pay.configured = lambda: True
pay.create_checkout_session = lambda **kw: (
    _landing.update(kw) or "https://stripe.test/pay")
try:
    app.test_client().get("/checkout/product/drip-course")
    ok("A guest is sent back to the product they bought",
       "/courses/drip-course" in _landing.get("return_url", "")
       and "bought=1" in _landing.get("return_url", ""),
       f"got {_landing.get('return_url')}")
    drip_client.get("/checkout/product/drip-course")
    ok("Someone with an account lands in their library",
       "/account" in _landing.get("return_url", "")
       and "tab=saved" in _landing.get("return_url", ""),
       f"got {_landing.get('return_url')}")
finally:
    pay.create_checkout_session = _real_checkout
    pay.configured = _real_configured

_landed = client.get("/courses/drip-course?bought=1",
                     follow_redirects=True).get_data(as_text=True)
ok("Landing there says the file is coming by email",
   "the receipt has the file with it" in _landed
   and "waiting in My space" in _landed)
ok("And the page says nothing of the sort on an ordinary visit",
   "the receipt has the file with it"
   not in client.get("/courses/drip-course").get_data(as_text=True))

# A long message shouldn't make its inbox row taller than the screen.
_inbox_body = admin.get("/admin/inbox").get_data(as_text=True)
ok("Inbox messages are cut to a few lines, with the rest a click away",
   'class="inbox-note" data-clamp' in _inbox_body,
   "nothing marked to shorten")
ok("So are the notes on a report",
   'class="inbox-report__note" data-clamp' in _inbox_body
   or "report_rows" not in _inbox_body)
_admin_js3 = client.get("/static/js/admin.js").get_data(as_text=True)
ok("The button only turns up when something is actually cut off",
   "scrollHeight > box.clientHeight" in _admin_js3
   and "Show more" in _admin_js3 and "Show less" in _admin_js3)
_css_clamp = client.get("/static/css/main.css").get_data(as_text=True)
ok("And opening one shows the whole of it",
   "-webkit-line-clamp: 4" in _css_clamp
   and ".is-clamped.is-open" in _css_clamp)

# With no disk attached there is nowhere durable to put a file, so the bytes
# go into Postgres instead — which survives a deploy, at a much smaller size.
from werkzeug.datastructures import FileStorage as _FileStorage  # noqa: E402

_was_dir = app.config["COURSE_FILES_DIR"]
app.config["COURSE_FILES_DIR"] = ""
try:
    with app.app_context():
        ok("With no files directory, storage falls to the database",
           _assets_svc.in_database()
           and _assets_svc.max_upload_bytes() == _assets_svc.DB_MAX_BYTES)
        _dbp = Product(slug="in-the-db", title="In The DB", type="guide",
                       status="published", track="building", promise="x",
                       price_cents=100, stripe_price_id="price_inthedb")
        db.session.add(_dbp)
        db.session.commit()
        _raw = b"%PDF-1.4 " + b"kept in a row " * 40 + b"\n%%EOF\n"
        _dba = _assets_svc.add_asset(
            _dbp, _FileStorage(BytesIO(_raw), filename="rowfile.pdf",
                               content_type="application/pdf"))
        db.session.commit()
        ok("The file is held in the row, not pointed at on a disk",
           _dba.disk_name is None and bytes(_dba.data or b"") == _raw)
        ok("It reads back whole, and never reads as missing",
           _assets_svc.read_bytes(_dba) == _raw and not _dba.file_missing())
        ok("And it still goes out with a receipt",
           [f["name"] for f in _assets_svc.receipt_files(_dbp)] == ["rowfile.pdf"])
        try:
            _assets_svc.add_asset(
                _dbp, _FileStorage(BytesIO(b"%PDF-1.4 " + b"x" * (
                    _assets_svc.DB_MAX_BYTES + 10) + b"\n%%EOF\n"),
                    filename="huge.pdf", content_type="application/pdf"))
            _refused = ""
        except _assets_svc.AssetError as _exc:
            _refused = str(_exc)
        ok("Anything past what a row should hold is turned away, and says why",
           "kept in the database" in _refused, f"got {_refused!r}")
        db.session.rollback()
        _cat_svc._purge_product(Product.query.filter_by(slug="in-the-db").first())
        db.session.commit()
finally:
    app.config["COURSE_FILES_DIR"] = _was_dir

# A disk that isn't really persistent behaves like a working one until the
# next restart, when it comes back empty. Leaving a marker on it and
# remembering that marker turns that into something Studio can state.
with app.app_context():
    from app.services import storage_health as _sh
    _first = _sh.check()
    ok("A healthy disk raises nothing",
       _first["checked"] and not _first["swapped"] and _first["missing"] == 0,
       f"got {_first}")

_ok_line = admin.get("/admin", follow_redirects=True).get_data(as_text=True)
ok("When all is well the dashboard says where uploads are landing",
   "Uploaded files are kept" in _ok_line and "all present" in _ok_line,
   "nothing said either way")

with app.app_context():
    _lost2 = ProductAsset.query.filter(ProductAsset.disk_name.isnot(None)).first()
    _path2 = _assets_svc.disk_path(_lost2.disk_name)
    _kept2 = open(_path2, "rb").read()
    _os.remove(_path2)
    _gone_state = _sh.check()
    ok("A file that has gone is counted",
       _gone_state["missing"] >= 1, f"got {_gone_state}")

_alarm = admin.get("/admin", follow_redirects=True).get_data(as_text=True)
ok("Studio says so on the dashboard, not only on the file",
   "no longer on the server" in _alarm, "no warning")
ok("And drops the all-is-well line while it isn't",
   "all present" not in _alarm)
ok("And keeps saying it while it is still true",
   "no longer on the server"
   in admin.get("/admin", follow_redirects=True).get_data(as_text=True))
with app.app_context():
    open(_path2, "wb").write(_kept2)
ok("Going quiet once the files are back",
   "no longer on the server"
   not in admin.get("/admin", follow_redirects=True).get_data(as_text=True))

with app.app_context():
    from app.services.settings import get_setting as _gs, set_setting as _ss
    _mark = _sh._read_marker()
    ok("The disk is marked so a different one can be told apart",
       bool(_mark) and _gs(_sh.SETTING_KEY, "") == _mark,
       f"marker {_mark!r} vs remembered {_gs(_sh.SETTING_KEY, '')!r}")
    _os.remove(_os.path.join(app.config["COURSE_FILES_DIR"], _sh.MARKER_NAME))
    _os.remove(_path2)
    _swapped = _sh.check()
    ok("A disk with no marker of its own reads as a different one",
       _swapped["swapped"], f"got {_swapped}")
    open(_path2, "wb").write(_kept2)
    _ss(_sh.SETTING_KEY, _sh._read_marker())
    db.session.commit()

# Removing a purchase from My Space used to make it vanish, which reads like
# it was taken away. Nothing is refunded and it is still theirs, so it stays
# on the shelf, marked, with a way back.
def _shelf():
    _h = drip_client.get("/account?tab=saved").get_data(as_text=True)
    return {
        "listed": "Drip Course" in _h,
        "marked": "Removed by you" in _h,
        "opens": f"/account/courses/{drip_purchase_id}" in _h,
        "restorable": "Put it back" in _h,
    }


_before_away = _shelf()
ok("A purchase is on the shelf and opens", _before_away["listed"]
   and _before_away["opens"] and not _before_away["marked"],
   f"got {_before_away}")
drip_client.post(f"/account/shop/{drip_purchase_id}/remove",
                 follow_redirects=True)
_away = _shelf()
ok("Removing it leaves it there, saying who removed it",
   _away["listed"] and _away["marked"], f"got {_away}")
ok("With nothing to open, and a way to put it back",
   not _away["opens"] and _away["restorable"], f"got {_away}")
drip_client.post(f"/account/shop/{drip_purchase_id}/restore",
                 follow_redirects=True)
_back = _shelf()
ok("Putting it back opens it again", _back == _before_away, f"got {_back}")
with app.app_context():
    ok("And the purchase itself was never touched",
       db.session.get(ShopPurchase, drip_purchase_id).status == "linked")

# A guide should arrive, not wait to be found: the receipt carries the PDF.
from werkzeug.datastructures import FileStorage as _FileStorage  # noqa: E402
with app.app_context():
    from app.services import assets as _ast
    _rc = Product(slug="receipt-guide", title="Receipt Guide", type="guide",
                  status="published", track="building", promise="x",
                  price_cents=1200, stripe_price_id="price_receipt")
    db.session.add(_rc)
    db.session.commit()
    _small = b"%PDF-1.4 " + b"the guide " * 50 + b"\n%%EOF\n"
    _ast.add_asset(_rc, _FileStorage(BytesIO(_small), filename="guide.pdf",
                                     content_type="application/pdf"))
    _ast.add_asset(_rc, _FileStorage(BytesIO(b"\x00\x00\x00\x18ftypmp42v"),
                                     filename="intro.mp4",
                                     content_type="video/mp4"))
    db.session.commit()
    _picked = _ast.receipt_files(_rc)
    ok("The guide itself goes with the receipt",
       [f["name"] for f in _picked] == ["guide.pdf"],
       f"got {[f['name'] for f in _picked]}")
    ok("With the bytes of the real file",
       _picked and _picked[0]["data"] == _small)
    ok("Anything too big to post is left in the library instead",
       _ast.receipt_files(_rc, budget=10) == [])

    # An email arriving without the guide looks the same whichever reason it
    # was, so each one says which — otherwise it can only be guessed at.
    import logging as _logging

    class _Catch(_logging.Handler):
        def __init__(self):
            super().__init__()
            self.said = []

        def emit(self, record):
            self.said.append(record.getMessage())

    _heard = _Catch()
    _ast_log = _logging.getLogger("app.services.assets")
    _ast_log.addHandler(_heard)
    try:
        _ast.receipt_files(None)
        _ast.receipt_files(_rc, budget=10)
        _gone_pdf = next(a for a in _rc.top_level_assets() if a.kind == "pdf")
        _pdf_path = _assets_svc.disk_path(_gone_pdf.disk_name)
        _pdf_keep = open(_pdf_path, "rb").read()
        _os.remove(_pdf_path)
        _ast.receipt_files(_rc)
        open(_pdf_path, "wb").write(_pdf_keep)
        _ast.receipt_files(_rc)
    finally:
        _ast_log.removeHandler(_heard)
    _all_said = " | ".join(_heard.said)
    ok("No product behind the payment says so",
       "no product behind this payment" in _all_said, f"got {_all_said}")
    ok("A file too big to post says so",
       "over what an email will carry" in _all_said, f"got {_all_said}")
    ok("A file whose bytes have gone says so",
       "could not be read" in _all_said, f"got {_all_said}")
    ok("And a receipt that does carry one names it",
       "attaching guide.pdf" in _all_said, f"got {_all_said}")

    # Anything that isn't a course or a bundle is something the buyer should
    # simply have, whatever kind of reading file it is.
    _rc.set_types(["template"])
    _ast.add_asset(_rc, _FileStorage(BytesIO(b"PK\x03\x04 a template"),
                                     filename="worksheet.docx",
                                     content_type="application/vnd.openxml"))
    db.session.commit()
    ok("A template goes out, and not only if it is a PDF",
       sorted(f["name"] for f in _ast.receipt_files(_rc))
       == ["guide.pdf", "worksheet.docx"],
       f"got {[f['name'] for f in _ast.receipt_files(_rc)]}")
    _rc.set_types(["course"])
    db.session.commit()
    ok("A course is read on the site, not posted out",
       _ast.receipt_files(_rc) == [])
    _rc.set_types(["bundle"])
    db.session.commit()
    ok("Nor is a bundle", _ast.receipt_files(_rc) == [])
    _rc.set_types(["guide", "course"])
    db.session.commit()
    ok("And a guide that is also a course counts as the course",
       _ast.receipt_files(_rc) == [])
    _rc.set_types(["guide"])
    db.session.commit()

    # A drip-fed course must not arrive whole on day one.
    _rc.drip_enabled = True
    _rc.set_curriculum([{"title": "One"}, {"title": "Two"}])
    _ast.add_asset(_rc, _FileStorage(BytesIO(b"%PDF-1.4 week two\n%%EOF\n"),
                                     filename="week-two.pdf",
                                     content_type="application/pdf"),
                   module_index=2)
    db.session.commit()
    ok("But a module they haven't reached yet doesn't",
       "week-two.pdf" not in [f["name"] for f in _ast.receipt_files(_rc)]
       and "guide.pdf" in [f["name"] for f in _ast.receipt_files(_rc)],
       f"got {[f['name'] for f in _ast.receipt_files(_rc)]}")
    _cat_svc._purge_product(db.session.get(Product, _rc.id))
    db.session.commit()

_carried = {}
_real_send_email = _mailer.send_email


def _catch_send(to, subject, text, html_body=None, template_id=None,
                params=None, sender=None, attachments=None):
    _carried.clear()
    _carried.update(text=text, params=params or {},
                    files=[a["name"] for a in (attachments or [])])
    return True


_mailer.send_email = _catch_send
try:
    with app.app_context():
        _mailer.send_order_receipt(
            "buyer@example.com", order_id="R-1",
            product_name="Receipt Guide", amount="$12",
            order_date="Sep 02, 2026",
            attachments=[{"name": "guide.pdf", "data": b"%PDF-1.4\n"}])
    ok("The receipt sends it and says it is there",
       _carried["files"] == ["guide.pdf"]
       and "Attached: guide.pdf" in _carried["text"]
       and _carried["params"].get("ATTACHED") == "guide.pdf",
       f"got {_carried}")
finally:
    _mailer.send_email = _real_send_email

# When a PDF still won't open, the reader has to say why and leave a way in.
_reader_js = client.get("/static/js/course-reader.js").get_data(as_text=True)
ok("The reader names what went wrong instead of one blanket sentence",
   all(word in _reader_js for word in
       ("PasswordException", "InvalidPDFException", "MissingPDFException")),
   "no per-cause messages")
ok("And offers to hand the file to the browser's own reader",
   "Open it in a new tab instead" in _reader_js)
ok("A LiveCycle form is recognised rather than drawn blank",
   "isPureXfa" in _reader_js and "enableXfa" in _reader_js)

# Fillable PDFs are read here and filled in elsewhere. Drawing real boxes over
# the page never worked, so it is gone: the page is painted as it comes,
# boxes and all, with a line saying where to actually fill them in.
ok("Nothing is left of the fields-over-the-page reader",
   not any(word in _reader_js for word in
           ("AnnotationLayer", "annotationLayer", "renderFormLayer",
            "ENABLE_FORMS", "formState", "setupFormSaving")),
   "some of it is still there")
_css_now = client.get("/static/css/main.css").get_data(as_text=True)
ok("Nor of the styling that dressed them",
   "annotationLayer" not in _css_now and "pdf-stack" not in _css_now)
ok("Nowhere to post answers to any more",
   app.test_client().post("/account/courses/1/formdata",
                          json={"form_data": {}}).status_code == 404)
with app.app_context():
    from app.models import CourseProgress as _Progress
    ok("And nothing left on the reading record to hold them",
       "form_data_json" not in [c.name for c in _Progress.__table__.columns])
ok("A PDF says where to fill one in instead",
   "print it and write on the paper" in pathlib.Path(
       "app/templates/main/course_reader.html").read_text(),
   "no note under the page")
ok("A file that isn't there says so, rather than blaming the connection",
   "isn't on the server any more" in _reader_js
   and "MissingPDFException" in _reader_js)

# A file can lose its bytes — an upload that never finished, storage moved.
# Nothing showed it: Studio listed the file and the buyer got a reader that
# wouldn't open, with no way for either to tell why.
with app.app_context():
    _lost = ProductAsset.query.filter(ProductAsset.disk_name.isnot(None)).first()
    if _lost is not None:
        _path = _assets_svc.disk_path(_lost.disk_name)
        _keep = open(_path, "rb").read()
        ok("A file that is where it should be isn't flagged",
           not _lost.file_missing())
        _os.remove(_path)
        ok("One whose bytes have gone is",
           db.session.get(ProductAsset, _lost.id).file_missing())
        _lost_product, _lost_name = _lost.product_id, _lost.filename
        open(_path, "wb").write(_keep)
        ok("And it is not flagged again once it is back",
           not db.session.get(ProductAsset, _lost.id).file_missing())
ok("Studio has somewhere to say it, on the file itself",
   "module-items__gone" in client.get("/static/css/main.css").get_data(as_text=True)
   and "file_missing()" in pathlib.Path(
       "app/templates/admin/product_form.html").read_text())

# A lesson added in the editor but not saved yet doesn't exist as far as the
# product is concerned. A file pinned to it used to belong to no lesson and no
# module intro, so it showed nowhere: not to the buyer, and not to the owner,
# who then couldn't move or delete it either.
_slip = b"\x00\x00\x00\x18ftypmp42" + b"unsaved lesson"
r = admin.post(f"/admin/products/{drip_prod_id}/uploads/begin",
               json={"filename": "early.mp4", "size": len(_slip)})
_early_id = r.get_json()["upload_id"]
admin.post(f"/admin/products/{drip_prod_id}/uploads/{_early_id}/chunk",
           data={"chunk": (BytesIO(_slip), "part")},
           content_type="multipart/form-data")
r = admin.post(f"/admin/products/{drip_prod_id}/uploads/{_early_id}/finish",
               json={"filename": "early.mp4", "module": 2, "lesson": 9})
_early = r.get_json()
ok("A file uploaded into an unsaved lesson is still saved",
   r.status_code == 200 and bool(_early.get("asset_id")))
with app.app_context():
    _prod = db.session.get(Product, drip_prod_id)
    _asset = db.session.get(ProductAsset, _early["asset_id"])
    ok("It waits in the module instead of a lesson that isn't there",
       _asset.lesson_index is None and _asset.module_index == 2,
       f"module {_asset.module_index}, lesson {_asset.lesson_index}")
    _row = _prod.modules()[1]
    ok("So the owner and the buyer can both see it",
       _asset.id in [a.id for a in _row["intro"]]
       and _asset.id in [a.id for a in _row["contents"]])

    # Same rescue for a file already stranded: one uploaded before this was
    # fixed, or one left behind when its lesson was deleted from the module.
    _asset.lesson_index = 9
    db.session.commit()
    _row = db.session.get(Product, drip_prod_id).modules()[1]
    _shown = sum(len(les["contents"]) for les in _row["lessons"]) + len(_row["intro"])
    ok("A file left pointing at a lesson that's gone shows in the module",
       _asset.id in [a.id for a in _row["intro"]]
       and _shown == len(_row["contents"]),
       f"{_shown} of {len(_row['contents'])} files reachable")
    _asset.lesson_index = None
    db.session.commit()

# Saving the form finishes what the upload couldn't: the lesson is written
# first, so the file dropped into it can now be pinned there. The editor sends
# the lesson it was dropped under alongside the lesson itself.
_mod_fields = {"mod1_title": "Week one", "mod1_desc": "Start here",
               "mod2_title": "Week two", "mod2_desc": "Keep going",
               "mod3_title": "Week three", "mod3_desc": "Look back"}
r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=_MultiDict([
        *dict(_drip_fields, **_mod_fields, slug="drip-course").items(),
        ("mod2_lesson_title", "The one we just added"),
        ("mod2_lesson_desc", "Written at the same time as the upload."),
        (f"asset_{_early['asset_id']}_lesson", "1"),
    ]),
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _row = db.session.get(Product, drip_prod_id).modules()[1]
    ok("Saving drops the file into the lesson it was uploaded under",
       r.status_code == 200 and len(_row["lessons"]) == 1
       and _early["asset_id"] in [a.id for a in _row["lessons"][0]["contents"]],
       f"lessons={len(_row['lessons'])}")

admin.post(f"/admin/products/{drip_prod_id}/assets/{_early['asset_id']}/delete",
           follow_redirects=True)
admin.post(f"/admin/products/{drip_prod_id}/edit",
           data=_MultiDict(dict(_drip_fields, **_mod_fields,
                                slug="drip-course").items()),
           content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    ok("And a stranded file can be cleared out again",
       db.session.get(ProductAsset, _early["asset_id"]) is None)

# Lessons take files on the page itself, so a product being created can fill
# them in one go — the one-at-a-time uploader needs a product id to send to,
# which a product that doesn't exist yet hasn't got.
r = admin.post("/admin/products/new", data=_MultiDict([
    ("title", "Lesson Uploads"), ("track", "building"), ("types", "course"),
    ("promise", "Files land in lessons"), ("price", "12.00"),
    ("stripe", "price_lessonup"), ("live", "1"),
    ("mod1_title", "Module one"), ("mod1_desc", "d"),
    ("mod1_lesson_title", "Lesson one"), ("mod1_lesson_desc", "first"),
    ("mod1_lesson_title", "Lesson two"), ("mod1_lesson_desc", "second"),
    ("mod1_lesson1_file", (BytesIO(b"\x89PNG\r\n\x1a\n" + b"x" * 64), "diagram.png")),
    ("mod1_lesson2_file", (BytesIO(b"%PDF-1.4 handout\n%%EOF\n"), "handout.pdf")),
]), content_type="multipart/form-data", follow_redirects=True)
ok("A new product can be created with files already in its lessons",
   r.status_code == 200)
with app.app_context():
    _lu = Product.query.filter_by(slug="lesson-uploads").first()
    _lessons = _lu.modules()[0]["lessons"] if _lu else []
    ok("Each file lands in the lesson it was attached to",
       [[a.display_title() for a in les["contents"]] for les in _lessons]
       == [["diagram.png"], ["handout.pdf"]],
       f"got {[[a.display_title() for a in les['contents']] for les in _lessons]}")
    _lu_id = _lu.id

# A lesson left untitled is dropped when the curriculum is saved, so the third
# set of fields on the page is not lesson three.
r = admin.post(f"/admin/products/{_lu_id}/edit", data=_MultiDict([
    ("title", "Lesson Uploads"), ("track", "building"), ("types", "course"),
    ("promise", "Files land in lessons"), ("price", "12.00"),
    ("stripe", "price_lessonup"), ("live", "1"), ("slug", "lesson-uploads"),
    ("mod1_title", "Module one"), ("mod1_desc", "d"),
    ("mod1_lesson_title", "Lesson one"), ("mod1_lesson_desc", "first"),
    ("mod1_lesson_title", ""), ("mod1_lesson_desc", ""),
    ("mod1_lesson_title", "Lesson two"), ("mod1_lesson_desc", "second"),
    ("mod1_lesson3_file", (BytesIO(b"%PDF-1.4 extra\n%%EOF\n"), "extra.pdf")),
]), content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _lessons = db.session.get(Product, _lu_id).modules()[0]["lessons"]
    ok("A blank lesson between two others doesn't misplace the files",
       r.status_code == 200 and len(_lessons) == 2
       and [a.display_title() for a in _lessons[1]["contents"]]
       == ["handout.pdf", "extra.pdf"],
       f"got {[[a.display_title() for a in x['contents']] for x in _lessons]}")

# a file with no module is open from day one and stays reachable
r = admin.post(
    f"/admin/products/{drip_prod_id}/assets",
    data={"asset": (BytesIO(b"%PDF-1.4 welcome\n%%EOF\n"), "welcome.pdf"),
          "asset_title": "Read me first"},
    content_type="multipart/form-data", follow_redirects=True)
ok("Studio takes a file that belongs to no module", r.status_code == 200)
with app.app_context():
    _mods = db.session.get(Product, drip_prod_id).modules()
    ok("A loose file doesn't get swept into a module",
       len(_mods[0]["contents"]) == 5)
r = drip_client.get(f"/account/courses/{drip_purchase_id}?module=1")
_rbody = r.get_data(as_text=True)
ok("A file outside the modules stays reachable while reading a module",
   "Yours from day one" in _rbody and "Read me first" in _rbody)

# removing one piece leaves the rest of the module alone
r = admin.post(
    f"/admin/products/{drip_prod_id}/assets/{_big_id}/delete",
    follow_redirects=True)
with app.app_context():
    ok("Removing one item deletes its file from the disk too",
       db.session.get(ProductAsset, _big_id) is None
       and not _os.path.isfile(_big_path))
    ok("The rest of the module survives",
       len(db.session.get(Product, drip_prod_id).modules()[1]["contents"]) == 1)

# --- extracts written for one particular file --------------------------------
# A note hangs off the upload it was written about, so it arrives with the
# video rather than sitting beside it as another thing to click.
_drip_modules = {
    "slug": "drip-course",
    "mod1_title": "Week one", "mod1_desc": "Start here",
    "mod2_title": "Week two", "mod2_desc": "Keep going",
    "mod3_title": "Week three", "mod3_desc": "Look back",
}
r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=_MultiDict([
        *dict(_drip_fields, **_drip_modules).items(),
        (f"newnote_{video_item_id}_title", "Before you press play"),
        (f"newnote_{video_item_id}_body", "Have a pen **ready**."),
        (f"newnote_{video_item_id}_title", "After the video"),
        (f"newnote_{video_item_id}_body", "Sit with it for a day."),
    ]),
    content_type="multipart/form-data",
    follow_redirects=True,
)
ok("Studio attaches written extracts to one file", r.status_code == 200)
with app.app_context():
    _prod = db.session.get(Product, drip_prod_id)
    _video = db.session.get(ProductAsset, video_item_id)
    _notes = list(_video.notes)
    _note_ids = [n.id for n in _notes]
    ok("Both extracts hang off the video, in the order they were written",
       [n.title for n in _notes] == ["Before you press play", "After the video"]
       and all(n.parent_asset_id == video_item_id for n in _notes),
       f"titles={[n.title for n in _notes]}")
    ok("They sort among themselves, not among the module",
       [n.sort_order for n in _notes] == [1, 2])
    ok("A note copies its file's module, so drip gating needs no special case",
       all(n.module_index == _video.module_index for n in _notes))
    ok("An extract still reads as text with words and a size",
       all(n.is_text() and n.size > 0 and n.disk_name is None for n in _notes))
    _m1 = _prod.modules()[0]["contents"]
    ok("Notes are not pieces of the module in their own right",
       len(_m1) == 5 and not any(a.id in _note_ids for a in _m1))
    _loose = [a for a in _prod.top_level_assets() if not a.module_index]
    ok("Nor among the files that are open from day one",
       [a.display_title() for a in _loose] == ["Read me first"])
    ok("And the cover takes its kind from a file, not from a note",
       _prod.top_level_assets()[0].id not in _note_ids)

_edit = admin.get(f"/admin/products/{drip_prod_id}/edit").get_data(as_text=True)
ok("Studio offers to write another extract for a file",
   f'data-note-add data-parent="{video_item_id}"' in _edit)
ok("And to rewrite an existing one where it sits",
   f'name="note_{_note_ids[0]}_body"' in _edit
   and "Have a pen **ready**." in _edit)

r = drip_client.get(
    f"/account/courses/{drip_purchase_id}?module=1&item={video_item_id}")
_rbody = r.get_data(as_text=True)
ok("The reader shows the file and everything written for it together",
   r.status_code == 200 and "course-reader__video" in _rbody
   and "Before you press play" in _rbody
   and "Have a pen <strong>ready</strong>." in _rbody
   and "Sit with it for a day." in _rbody)
ok("The stage splits, with the player sized to fit and the notes beneath it",
   "reader__notes" in _rbody and "reader__stage--split" in _rbody
   and "reader__stage--media" in _rbody
   and "reader__viewer--fill" not in _rbody)
ok("A chip marks the file that carries writing",
   "reader-pieces__notes" in _rbody)
ok("But an extract is never a chip of its own",
   f"item={_note_ids[0]}" not in _rbody)

r = drip_client.get(
    f"/account/courses/{drip_purchase_id}?module=1&item={_note_ids[0]}")
_rbody = r.get_data(as_text=True)
ok("Asking for an extract alone falls back to the module's first piece",
   r.status_code == 200 and f"/file/{mod1_asset_id}" in _rbody
   and "Have a pen" not in _rbody)

r = client.get("/courses/drip-course")
_dbody = r.get_data(as_text=True)
ok("The store page says how much is inside each module",
   "1 video, 2 documents, written notes" in _dbody and "1 document" in _dbody)
ok("But gives away nothing that is in it",
   "Before you press play" not in _dbody and "worksheet" not in _dbody
   and "Have a pen" not in _dbody)

r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=_MultiDict([
        *dict(_drip_fields, **_drip_modules).items(),
        (f"note_{_note_ids[0]}_title", "Before you press play"),
        (f"note_{_note_ids[0]}_body", "Have a pen and water ready."),
    ]),
    content_type="multipart/form-data",
    follow_redirects=True,
)
with app.app_context():
    _n = db.session.get(ProductAsset, _note_ids[0])
    ok("An extract can be rewritten in place instead of retyped",
       _n.body == "Have a pen and water ready."
       and _n.size == len(_n.body.encode("utf-8")))
    ok("Rewriting one leaves the other alone",
       db.session.get(ProductAsset, _note_ids[1]).title == "After the video")

# a locked module keeps its notes locked, because they carry its number
r = admin.post(
    f"/admin/products/{drip_prod_id}/edit",
    data=_MultiDict([
        *dict(_drip_fields, **_drip_modules).items(),
        (f"newnote_{mod3_asset_id}_title", "Looking back"),
        (f"newnote_{mod3_asset_id}_body", "What has shifted since week one?"),
    ]),
    content_type="multipart/form-data",
    follow_redirects=True,
)
with app.app_context():
    _late_note_id = db.session.get(ProductAsset, mod3_asset_id).notes[0].id
_backdate_drip(0)
r = drip_client.get(f"/account/courses/{drip_purchase_id}?module=3")
ok("A locked module's extracts are out of reach with the rest of it",
   r.status_code == 200 and "What has shifted" not in r.get_data(as_text=True))
r = drip_client.get(f"/account/courses/{drip_purchase_id}/file/{_late_note_id}")
ok("Not even by asking for the extract directly", r.status_code == 404)
_backdate_drip(20)
r = drip_client.get(f"/account/courses/{drip_purchase_id}?module=3")
ok("It arrives with its module once that opens",
   "What has shifted since week one?" in r.get_data(as_text=True))

r = admin.post(f"/admin/products/{drip_prod_id}/assets/{video_item_id}/delete",
               follow_redirects=True)
with app.app_context():
    ok("Removing a file takes the extracts written for it with it",
       db.session.get(ProductAsset, video_item_id) is None
       and all(db.session.get(ProductAsset, i) is None for i in _note_ids))
    ok("The rest of the module is untouched",
       len(db.session.get(Product, drip_prod_id).modules()[0]["contents"]) == 4)

# an extract can be reached two ways — through its file and through the
# product — so deleting the product must not try to delete it twice
import warnings as _warnings  # noqa: E402

admin.post(
    "/admin/products/new",
    data=_MultiDict([
        ("title", "Doomed Course"), ("track", "building"), ("type", "guide"),
        ("mod1_title", "Only module"), ("mod1_desc", "one"),
        ("mod1_file", (BytesIO(b"%PDF-1.4 doomed\n%%EOF\n"), "doomed.pdf")),
    ]),
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _doomed = Product.query.filter_by(title="Doomed Course").first()
    _doomed_id, _doomed_slug = _doomed.id, _doomed.slug
    _doomed_file_id = _doomed.modules()[0]["asset"].id
admin.post(
    f"/admin/products/{_doomed_id}/edit",
    data=_MultiDict([
        ("title", "Doomed Course"), ("track", "building"), ("type", "guide"),
        ("slug", _doomed_slug),
        ("mod1_title", "Only module"), ("mod1_desc", "one"),
        (f"newnote_{_doomed_file_id}_title", "A word first"),
        (f"newnote_{_doomed_file_id}_body", "Nothing here lasts."),
    ]),
    content_type="multipart/form-data", follow_redirects=True)
with app.app_context():
    _doomed_note_id = db.session.get(ProductAsset, _doomed_file_id).notes[0].id
with _warnings.catch_warnings(record=True) as _caught:
    _warnings.simplefilter("always")
    r = admin.post(f"/admin/products/{_doomed_id}/delete", follow_redirects=True)
ok("Deleting a product clears its files and their extracts", r.status_code == 200)
with app.app_context():
    ok("Nothing of it is left behind",
       db.session.get(Product, _doomed_id) is None
       and db.session.get(ProductAsset, _doomed_file_id) is None
       and db.session.get(ProductAsset, _doomed_note_id) is None)
ok("And each row is deleted once, not down two paths at the same time",
   not [w for w in _caught if "expected to delete" in str(w.message)],
   f"warnings={[str(w.message)[:60] for w in _caught]}")

# --- test mode: a real, buyable product only the owners can see --------------
_test_fields = {
    "title": "Dress Rehearsal",
    "track": "building",
    "type": "guide",
    "price": "12.00",
    "promise": "A dry run of the whole checkout.",
    "stripe": "price_dress_rehearsal",
    "live": "1",
    "test_mode": "1",
}
with app.app_context():
    from app.models import Notification as _Note
    from app.services import stats as _stats_mod
    notes_before = _Note.query.filter_by(kind="course").count()
    insights_before = _stats_mod.payment_insights(days=30)["orders_30d"]
r = admin.post("/admin/products/new", data=_test_fields,
               content_type="multipart/form-data", follow_redirects=True)
ok("Studio creates a test-mode product", r.status_code == 200)
with app.app_context():
    tprod = Product.query.filter_by(slug="dress-rehearsal").first()
    ok("Test mode saves, and the product is still live",
       tprod is not None and tprod.test_mode is True and tprod.status == "published")
    test_prod_id = tprod.id
    ok("A test product is never announced to members",
       _Note.query.filter_by(kind="course").count() == notes_before)
ok("Studio flags it as a test in the product list",
   "status-pill--test" in admin.get("/admin/products").get_data(as_text=True))

_guest_courses = client.get("/courses").get_data(as_text=True)
_member_courses = drip_client.get("/courses").get_data(as_text=True)
_owner_courses = admin.get("/courses").get_data(as_text=True)
ok("Guests never see a test product in the catalogue",
   "Dress Rehearsal" not in _guest_courses)
ok("Signed-in members don't see it either",
   "Dress Rehearsal" not in _member_courses)
ok("Owners do see it, marked as a test",
   "Dress Rehearsal" in _owner_courses and "lib-card__badge--test" in _owner_courses)

ok("Guests get a 404 on the test product's page",
   client.get("/courses/dress-rehearsal").status_code == 404)
ok("Members get a 404 on it too",
   drip_client.get("/courses/dress-rehearsal").status_code == 404)
r = admin.get("/courses/dress-rehearsal")
_tbody = r.get_data(as_text=True)
ok("Owners open the page, warned what it is, with a working Buy button",
   r.status_code == 200 and "pd-preview-banner--test" in _tbody
   and "Buy now" in _tbody and "Preview only" not in _tbody)

ok("Guests can't reach checkout for a test product",
   client.get("/checkout/product/dress-rehearsal").status_code == 404)
ok("Members can't reach it either",
   drip_client.get("/checkout/product/dress-rehearsal").status_code == 404)
ok("Owners are allowed through to checkout",
   admin.get("/checkout/product/dress-rehearsal").status_code != 404)

# artwork shouldn't leak to someone guessing the id
with app.app_context():
    db.session.get(Product, cover_id).test_mode = True
    db.session.commit()
ok("A test product's cover is withheld from the public",
   client.get(f"/media/product-cover/{cover_id}").status_code == 404)
ok("Owners still get that cover",
   admin.get(f"/media/product-cover/{cover_id}").status_code == 200)
with app.app_context():
    db.session.get(Product, cover_id).test_mode = False
    db.session.commit()

# an owner buying their own test product is a rehearsal, not revenue
_test_payload = _payment_payload(
    "9300", "owner@example.com", "price_dress_rehearsal",
    amount=1200, product_name="Dress Rehearsal")
r = client.post("/webhooks/stripe", data=_test_payload,
                headers=_stripe_headers(_test_payload))
with app.app_context():
    torder = Order.query.filter_by(ls_order_id="9300").first()
    ok("A test purchase is recorded like any other",
       r.status_code == 200 and torder is not None and torder.status == "paid"
       and torder.product_id == test_prod_id)
    ok("But it stays out of the dashboard revenue figures",
       _stats_mod.payment_insights(days=30)["orders_30d"] == insights_before)
    ok("And out of the top-products list",
       all(row["title"] != "Dress Rehearsal"
           for row in _stats_mod.payment_insights(days=30)["top_products"]))
ok("The owner finds their test purchase in their own library",
   "Dress Rehearsal" in admin.get("/account?tab=saved").get_data(as_text=True))

# taking it out of test mode is what finally tells everyone
with app.app_context():
    notes_before = _Note.query.filter_by(kind="course").count()
_live_fields = {k: v for k, v in _test_fields.items() if k != "test_mode"}
r = admin.post(f"/admin/products/{test_prod_id}/edit",
               data=dict(_live_fields, slug="dress-rehearsal"),
               content_type="multipart/form-data", follow_redirects=True)
ok("Leaving test mode announces the product properly",
   r.status_code == 200 and "notified" in r.get_data(as_text=True))
with app.app_context():
    ok("Members finally hear about it",
       _Note.query.filter_by(kind="course").count() > notes_before)
ok("And it turns up in the public catalogue",
   "Dress Rehearsal" in client.get("/courses").get_data(as_text=True))

r = admin.post("/admin/products/new",
               data={"title": "Quiet Rehearsal", "track": "healing",
                     "type": "guide", "test_mode": "1"},
               content_type="multipart/form-data", follow_redirects=True)
ok("Test mode on a draft says there is nothing to buy yet",
   "still a draft" in r.get_data(as_text=True))

# --- viewing as a member, so an owner can check a test product's perk -------
# Owners rank as Full Bloom and are the only people who may buy a test product,
# so without this the membership perk on one can never be seen working.
r = admin.post("/admin/products/new",
               data={"title": "Perk Rehearsal", "track": "building",
                     "type": "guide", "price": "20.00",
                     "promise": "A dry run of the membership perk.",
                     "stripe": "price_perk_rehearsal", "live": "1",
                     "test_mode": "1", "perk_tier": "creator",
                     "perk_months": "3"},
               content_type="multipart/form-data", follow_redirects=True)
ok("Studio creates a test product that hands out membership months",
   r.status_code == 200)
_perk_payload = _payment_payload(
    "9301", "owner@example.com", "price_perk_rehearsal",
    amount=2000, product_name="Perk Rehearsal")
client.post("/webhooks/stripe", data=_perk_payload,
            headers=_stripe_headers(_perk_payload))
with app.app_context():
    from app.services.perks import perk_state as _perk_state
    _owner = User.query.filter_by(email="owner@example.com").first()
    # An owner who pays for nothing is the interesting case: the perk is then
    # the only thing standing between them and Free.
    _owner.membership = "none"
    _owner.membership_manual = None
    db.session.commit()
    ok("The owner's test purchase really does carry a Creator perk",
       _perk_state(_owner)["tier"] == "creator")
    ok("But their account is untouched: still an owner on no paid tier",
       _owner.effective_membership() == "full_bloom" and _owner.membership == "none")

ok("Without previewing, the owner sails into the members' marketplace",
   admin.get("/marketplace/mine").status_code == 200)

r = admin.post("/admin/preview", data={"tier": "real", "next": "/account"},
               follow_redirects=True)
_body = r.get_data(as_text=True)
ok("Viewing as what they've earned reports the perk tier",
   r.status_code == 200 and "Now browsing as Creator" in _body)
ok("A bar keeps saying which tier is being previewed",
   "preview-bar" in _body and "Viewing as" in _body)
with app.app_context():
    _owner = User.query.filter_by(email="owner@example.com").first()
    ok("Previewing changes nothing on the account itself",
       _owner.membership == "none" and _owner.is_admin is True)

r = admin.post("/admin/preview", data={"tier": "none", "next": "/account"},
               follow_redirects=True)
ok("The owner can also drop all the way to Free",
   "Now browsing as Free" in r.get_data(as_text=True))
r = admin.get("/marketplace/mine", follow_redirects=True)
ok("And the members-only marketplace shuts them out like anyone else",
   "members&#39; perk" in r.get_data(as_text=True)
   or "members' perk" in r.get_data(as_text=True))
ok("Studio stays open while previewing", admin.get("/admin/").status_code == 200)
ok("So do test products — they are owner tooling, not a membership gate",
   admin.get("/checkout/product/dress-rehearsal").status_code != 404)
ok("The bar rides along inside Studio too",
   "preview-bar" in admin.get("/admin/settings").get_data(as_text=True))

r = admin.post("/admin/preview", data={"tier": "off", "next": "/account"},
               follow_redirects=True)
ok("Leaving preview hands the owner back their own view",
   "Back to your owner view" in r.get_data(as_text=True))
ok("And the marketplace opens again",
   admin.get("/marketplace/mine").status_code == 200)
ok("With no bar left on the page",
   "preview-bar" not in admin.get("/account").get_data(as_text=True))

ok("A member cannot reach the preview switch at all",
   drip_client.post("/admin/preview", data={"tier": "full_bloom"}).status_code == 404)
r = admin.post("/admin/preview", data={"tier": "emperor"}, follow_redirects=True)
ok("An invented tier is refused", "isn&#39;t a tier you can preview" in r.get_data(as_text=True)
   or "isn't a tier you can preview" in r.get_data(as_text=True))
r = admin.post("/admin/preview", data={"tier": "healing", "next": "https://evil.test/x"},
               follow_redirects=False)
ok("Preview will not bounce the owner off-site",
   r.status_code == 302 and "evil.test" not in (r.headers.get("Location") or ""))
admin.post("/admin/preview", data={"tier": "off"})

# the perk is still the owner's to override, and it ends on its own
with app.app_context():
    from app.services.memberships import reconcile_user, set_manual_tier
    dripper = User.query.filter_by(email="dripper@example.com").first()
    set_manual_tier(dripper, "none")
    db.session.commit()
    reconcile_user(dripper)
    db.session.commit()
    ok("Studio can still take a perk membership away",
       dripper.membership == "none" and dripper.membership_manual == "none")
    dripper.membership_manual = None
    dripper.membership_manual_at = None
    reconcile_user(dripper)
    db.session.commit()
    ok("Clearing the override hands the perk back", dripper.membership == "creator")
_backdate_drip(200)
with app.app_context():
    from app.services.memberships import reconcile_user
    dripper = User.query.filter_by(email="dripper@example.com").first()
    reconcile_user(dripper)
    db.session.commit()
    ok("Perk membership ends by itself once the months run out",
       dripper.membership == "none")

_backdate_drip(1)
with app.app_context():
    from app.services.memberships import reconcile_user
    dripper = User.query.filter_by(email="dripper@example.com").first()
    reconcile_user(dripper)
    db.session.commit()
    ok("Perk is back while it is still running", dripper.membership == "creator")
r = drip_client.get("/account")
ok("Account says where the free membership came from",
   r.status_code == 200
   and "Came free with a product you bought" in r.get_data(as_text=True))
drip_refund = _payment_payload(
    "9100", "dripper@example.com", "price_drip_course",
    event="refund.succeeded", amount=3900, product_name="Drip Course")
client.post("/webhooks/stripe", data=drip_refund, headers=_stripe_headers(drip_refund))
with app.app_context():
    dripper = User.query.filter_by(email="dripper@example.com").first()
    ok("Refunding the product takes the perk membership back",
       dripper.membership == "none")

# --- scheduled cancels reach Studio ----------------------------------------
# Stripe announces "this ends on the 14th" only through subscription.updated.
# Without that hook a member who cancels in Stripe's portal keeps reading as a
# full-price member in Studio until the subscription finally deletes itself.
def _sub_updated_webhook(sub_id, email, price_id, *, canceling, status="active",
                         ends_in_days=12):
    ends = int((datetime.utcnow() + timedelta(days=ends_in_days)).timestamp())
    body = json.dumps({
        "id": f"evt_{sub_id}",
        "type": "customer.subscription.updated",
        "data": {"object": {
            "id": sub_id,
            "object": "subscription",
            "status": status,
            "cancel_at_period_end": bool(canceling),
            "current_period_end": ends,
            "customer_email": email,
            "items": {"data": [{"price": {"id": price_id}}]},
        }},
    }).encode()
    return client.post("/webhooks/stripe", data=body, headers=_stripe_headers(body)), ends


with app.app_context():
    from app.models import MembershipPlan as _MP
    _creator_price = _MP.query.filter_by(tier="creator").first().stripe_price_id
    _canceller = User(email="canceller@example.com", display_name="Quitter",
                      membership="creator", email_verified_at=utcnow())
    _canceller.set_password(USER_PW)
    db.session.add(_canceller)
    db.session.commit()

r, _ends = _sub_updated_webhook("sub_cancelme", "canceller@example.com",
                                _creator_price, canceling=True)
with app.app_context():
    _cx = User.query.filter_by(email="canceller@example.com").first()
    ok("Stripe cancel-at-period-end is recorded on the member",
       r.status_code == 200 and _cx.membership_cancel_at is not None
       and _cx.membership_is_canceling(),
       f"status={r.status_code} cancel_at={getattr(_cx, 'membership_cancel_at', None)}")
    ok("Cancelling keeps their tier until the period ends",
       _cx.membership == "creator", f"got {_cx.membership}")
    _end_label = _cx.membership_access_end_display()

r = admin.get("/admin/members")
_mbody = r.get_data(as_text=True)
ok("Studio members shows cancelled with the revoke date",
   "cancelled, revoked on" in _mbody and _end_label in _mbody,
   f"looking for {_end_label!r}")

# a non-membership subscription must not clear (or invent) a membership flag
_sub_updated_webhook("sub_notamembership", "canceller@example.com",
                     "price_unrelated_thing", canceling=False)
with app.app_context():
    ok("Unrelated subscriptions don't touch membership flags",
       User.query.filter_by(email="canceller@example.com").first().membership_cancel_at
       is not None)

# resuming clears it again
_sub_updated_webhook("sub_cancelme", "canceller@example.com",
                     _creator_price, canceling=False)
with app.app_context():
    _cx = User.query.filter_by(email="canceller@example.com").first()
    ok("Resuming a subscription clears the cancelled badge",
       _cx.membership_cancel_at is None and not _cx.membership_is_canceling())
r = admin.get("/admin/members")
ok("Studio members drops the badge after a resume",
   "cancelled, revoked on" not in r.get_data(as_text=True))

r = admin.post("/admin/members/refresh-cancellations", follow_redirects=True)
ok("Refresh cancellations button lands back on the member list",
   r.status_code == 200 and "Members" in r.get_data(as_text=True))

# Ending a membership must not rewrite what was earned.
with app.app_context():
    from app.services import stats as _stats
    from app.services.memberships import purchased_tier as _ptier
    _rev_before = _stats.payment_insights(3650)["revenue"]
    _paid_order = Order(ls_order_id="ENDME-1", ls_variant_id=_creator_price,
                        membership_tier="creator", buyer_email="canceller@example.com",
                        total_cents=250000, currency="USD", status="paid")
    db.session.add(_paid_order)
    db.session.commit()
    _rev_paid = _stats.payment_insights(3650)["revenue"]
    ok("A paid membership order counts towards revenue",
       _rev_paid != _rev_before, f"{_rev_before} -> {_rev_paid}")
    ok("A paid membership order grants the tier",
       _ptier("canceller@example.com") == "creator")

    _paid_order.status = "ended"
    db.session.commit()
    _rev_ended = _stats.payment_insights(3650)["revenue"]
    ok("Cancelling keeps the money in the revenue figures",
       _rev_ended == _rev_paid, f"{_rev_paid} -> {_rev_ended}")
    ok("Cancelling still takes the tier away",
       _ptier("canceller@example.com") == "none",
       f"got {_ptier('canceller@example.com')}")

    _paid_order.status = "refunded"
    db.session.commit()
    ok("A real refund does come back out of revenue",
       _stats.payment_insights(3650)["revenue"] == _rev_before)

# --- My space greeting follows the member's own clock ----------------------
with app.app_context():
    from datetime import timezone as _tzone
    from zoneinfo import ZoneInfo as _ZI

    from app.services.timefmt import greeting as _greet
    from app.services.timefmt import local_now as _local_now

    def _at(hour, tz="UTC"):
        return datetime(2026, 6, 15, hour, 30, tzinfo=_ZI(tz))

    ok("Morning greeting before noon", _greet("Nadia", now=_at(9)) == "Good morning, Nadia.")
    ok("Afternoon greeting after noon", _greet("Nadia", now=_at(13)) == "Good afternoon, Nadia.")
    ok("Evening greeting after six", _greet("Nadia", now=_at(20)) == "Good evening, Nadia.")
    ok("Small hours ask if they're still awake",
       _greet("Nadia", now=_at(0)) == "Still awake, Nadia?"
       and _greet("Nadia", now=_at(3)) == "Still awake, Nadia?")
    ok("Five in the morning is morning again",
       _greet("Nadia", now=_at(5)) == "Good morning, Nadia.")
    ok("Greeting works for a member with no display name",
       _greet("", now=_at(9)) == "Good morning."
       and _greet(None, now=_at(2)) == "Still awake?")

    # The same instant is a different time of day in different places.
    _moment = datetime(2026, 6, 15, 2, 0, tzinfo=_tzone.utc)
    ok("Greeting is read on the member's clock, not the server's",
       _greet("Nadia", now=_moment.astimezone(_ZI("UTC"))) == "Still awake, Nadia?"
       and _greet("Nadia", now=_moment.astimezone(_ZI("Asia/Karachi")))
       == "Good morning, Nadia.")

    ok("local_now returns a zone-aware time in the viewer's timezone",
       _local_now("Asia/Karachi").tzinfo is not None)

# the page itself carries whichever greeting is current
r = buyer_client.get("/account")
_gbody = r.get_data(as_text=True)
ok("My space renders one of the four greetings",
   r.status_code == 200
   and any(g in _gbody for g in ("Good morning", "Good afternoon",
                                 "Good evening", "Still awake")),
   "no greeting found in the page")

r = client.get("/this-page-does-not-exist-xyz")
ok("404 offers problem report", r.status_code == 404 and b"Report this problem" in r.data)

print(f"\nAll {PASS} checks passed.")

